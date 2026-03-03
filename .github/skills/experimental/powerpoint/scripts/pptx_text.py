"""Text frame, paragraph, and run property utilities for PowerPoint skill scripts.

Centralizes enhanced text properties (margins, auto-size, spacing, underline,
hyperlinks, bullets) used by build_deck.py and extract_content.py.
"""

from lxml import etree
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

AUTO_SIZE_MAP = {
    "none": MSO_AUTO_SIZE.NONE,
    "fit": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
    "shrink": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
}

AUTO_SIZE_REVERSE = {
    MSO_AUTO_SIZE.NONE: "none",
    MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT: "fit",
    MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE: "shrink",
}

VERTICAL_ANCHOR_MAP = {
    "top": MSO_VERTICAL_ANCHOR.TOP,
    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
}

VERTICAL_ANCHOR_REVERSE = {
    MSO_VERTICAL_ANCHOR.TOP: "top",
    MSO_VERTICAL_ANCHOR.MIDDLE: "middle",
    MSO_VERTICAL_ANCHOR.BOTTOM: "bottom",
}

# EMU per inch constant for margin conversions
_EMU_PER_INCH = 914400


def apply_text_properties(text_frame, elem: dict):
    """Apply text frame-level properties from element definition.

    Supports: margin_left/right/top/bottom (inches), auto_size, vertical_anchor, word_wrap.
    """
    if "margin_left" in elem:
        text_frame.margin_left = Inches(elem["margin_left"])
    if "margin_right" in elem:
        text_frame.margin_right = Inches(elem["margin_right"])
    if "margin_top" in elem:
        text_frame.margin_top = Inches(elem["margin_top"])
    if "margin_bottom" in elem:
        text_frame.margin_bottom = Inches(elem["margin_bottom"])
    if "auto_size" in elem:
        text_frame.auto_size = AUTO_SIZE_MAP.get(elem["auto_size"], MSO_AUTO_SIZE.NONE)
    if "vertical_anchor" in elem:
        text_frame.vertical_anchor = VERTICAL_ANCHOR_MAP.get(elem["vertical_anchor"])
    if "word_wrap" in elem:
        text_frame.word_wrap = elem["word_wrap"]


def apply_paragraph_properties(paragraph, elem: dict):
    """Apply paragraph-level properties.

    Supports: space_before, space_after (pts), line_spacing (pts or factor), level.
    """
    if "space_before" in elem:
        paragraph.space_before = Pt(elem["space_before"])
    if "space_after" in elem:
        paragraph.space_after = Pt(elem["space_after"])
    if "line_spacing" in elem:
        val = elem["line_spacing"]
        if isinstance(val, float) and val < 10:
            # Factor-based spacing (e.g. 1.5 = 150%)
            paragraph.line_spacing = val
        else:
            paragraph.line_spacing = Pt(val)
    if "level" in elem:
        paragraph.level = elem["level"]


def apply_run_properties(run, elem: dict, colors: dict):
    """Apply run-level font properties beyond basic font/size/color/bold/italic.

    Supports: underline, hyperlink, char_spacing, effect (outer shadow).
    When a hyperlink is set, the font color is re-applied afterward to prevent
    the automatic theme hyperlink color from overriding the intended color.
    """
    if elem.get("underline"):
        run.font.underline = True
    if "hyperlink" in elem:
        run.hyperlink.address = elem["hyperlink"]
        # Re-apply font color after hyperlink to override auto-coloring
        from pptx_colors import apply_color_to_font, resolve_color
        color_key = next((k for k in ("font_color", "text_color", "color") if k in elem), None)
        if color_key:
            apply_color_to_font(run.font.color, resolve_color(elem[color_key], colors))
    if "char_spacing" in elem:
        _apply_char_spacing(run.font, elem["char_spacing"])
    effect = elem.get("effect") or elem.get("text_effect")
    if effect:
        _apply_run_effect(run, effect)


def _apply_run_effect(run, effect: dict):
    """Apply outer shadow effect to a run's rPr element."""
    if not effect or effect.get("type") != "outer_shadow":
        return
    rpr = run.font._element
    # Remove existing effectLst
    existing = rpr.find(qn('a:effectLst'))
    if existing is not None:
        rpr.remove(existing)

    effect_lst = etree.SubElement(rpr, qn('a:effectLst'))
    shadow = etree.SubElement(effect_lst, qn('a:outerShdw'))
    for attr in ('blurRad', 'dist', 'dir', 'algn', 'rotWithShape'):
        if attr in effect:
            shadow.set(attr, str(effect[attr]))

    color_type = effect.get("color_type", "preset")
    color_val = effect.get("color", "black")
    if color_type == "preset":
        color_el = etree.SubElement(shadow, qn('a:prstClr'))
        color_el.set('val', color_val)
    elif color_type == "rgb":
        color_el = etree.SubElement(shadow, qn('a:srgbClr'))
        color_el.set('val', color_val.lstrip('#'))
    else:
        color_el = etree.SubElement(shadow, qn('a:prstClr'))
        color_el.set('val', 'black')

    if "alpha" in effect:
        alpha_sub = etree.SubElement(color_el, qn('a:alpha'))
        alpha_sub.set('val', str(int(effect["alpha"] * 1000)))


def _apply_char_spacing(font, spacing_pt: float):
    """Apply character spacing to a font via the spc attribute on a:rPr.

    Args:
        font: python-pptx font object.
        spacing_pt: Spacing in points (converted to hundredths of a point for XML).
    """
    rpr = font._element
    spc_val = str(int(spacing_pt * 100))
    rpr.set('spc', spc_val)


def extract_text_frame_properties(text_frame) -> dict:
    """Extract text frame-level properties (margins, auto_size, vertical_anchor)."""
    props = {}
    if text_frame.margin_left is not None:
        props["margin_left"] = round(text_frame.margin_left / _EMU_PER_INCH, 3)
    if text_frame.margin_right is not None:
        props["margin_right"] = round(text_frame.margin_right / _EMU_PER_INCH, 3)
    if text_frame.margin_top is not None:
        props["margin_top"] = round(text_frame.margin_top / _EMU_PER_INCH, 3)
    if text_frame.margin_bottom is not None:
        props["margin_bottom"] = round(text_frame.margin_bottom / _EMU_PER_INCH, 3)
    if text_frame.auto_size is not None:
        label = AUTO_SIZE_REVERSE.get(text_frame.auto_size)
        if label:
            props["auto_size"] = label
    if text_frame.vertical_anchor is not None:
        label = VERTICAL_ANCHOR_REVERSE.get(text_frame.vertical_anchor)
        if label:
            props["vertical_anchor"] = label
    return props


def extract_paragraph_properties(paragraph) -> dict:
    """Extract paragraph-level spacing properties."""
    props = {}
    if paragraph.space_before is not None:
        props["space_before"] = round(paragraph.space_before.pt, 1)
    if paragraph.space_after is not None:
        props["space_after"] = round(paragraph.space_after.pt, 1)
    if paragraph.line_spacing is not None:
        if isinstance(paragraph.line_spacing, float):
            props["line_spacing"] = paragraph.line_spacing
        else:
            props["line_spacing"] = round(paragraph.line_spacing.pt, 1)
    if paragraph.level and paragraph.level > 0:
        props["level"] = paragraph.level
    return props


def extract_run_properties(run) -> dict:
    """Extract run-level properties beyond basic font info (underline, hyperlink, char_spacing, effects)."""
    props = {}
    if run.font.underline:
        props["underline"] = True
    try:
        if run.hyperlink and run.hyperlink.address:
            props["hyperlink"] = run.hyperlink.address
    except (AttributeError, TypeError):
        pass
    # Character spacing
    from pptx_fonts import _extract_char_spacing
    spc = _extract_char_spacing(run.font)
    if spc is not None:
        props["char_spacing"] = spc
    # Outer shadow effect on text run
    effect = _extract_run_effect(run)
    if effect:
        props["effect"] = effect
    return props


def _extract_run_effect(run) -> dict | None:
    """Extract outer shadow effect from a run's rPr effectLst."""
    try:
        rpr = run.font._element
        effect_lst = rpr.find(qn('a:effectLst'))
        if effect_lst is None or len(effect_lst) == 0:
            return None
        shadow = effect_lst.find(qn('a:outerShdw'))
        if shadow is None:
            return None
        result = {"type": "outer_shadow"}
        for attr in ('blurRad', 'dist', 'dir', 'algn', 'rotWithShape'):
            val = shadow.get(attr)
            if val is not None:
                result[attr] = val
        color_el = shadow[0] if len(shadow) > 0 else None
        if color_el is not None:
            tag = color_el.tag.split('}')[-1]
            if tag == 'prstClr':
                result["color"] = color_el.get('val', 'black')
                result["color_type"] = "preset"
            elif tag == 'srgbClr':
                result["color"] = '#' + color_el.get('val', '000000')
                result["color_type"] = "rgb"
            alpha_el = color_el.find(qn('a:alpha'))
            if alpha_el is not None:
                result["alpha"] = round(int(alpha_el.get('val', '100000')) / 1000, 1)
        return result
    except (AttributeError, TypeError, IndexError):
        return None


_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def extract_bullet_properties(paragraph) -> dict:
    """Extract bullet properties from a paragraph's pPr element.

    Returns dict with bullet_char, bullet_font, bullet_size_pct, bullet_color
    when present. Returns {"bullet_none": True} when buNone is set.
    """
    props = {}
    pPr = paragraph._p.find(qn('a:pPr'))
    if pPr is None:
        return props

    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        props['bullet_none'] = True
        return props

    buChar = pPr.find(qn('a:buChar'))
    if buChar is not None:
        props['bullet_char'] = buChar.get('char', '•')

    buFont = pPr.find(qn('a:buFont'))
    if buFont is not None:
        typeface = buFont.get('typeface')
        if typeface:
            props['bullet_font'] = typeface

    buSzPct = pPr.find(qn('a:buSzPct'))
    if buSzPct is not None:
        val = buSzPct.get('val')
        if val:
            props['bullet_size_pct'] = int(val)

    buClr = pPr.find(qn('a:buClr'))
    if buClr is not None:
        srgb = buClr.find(qn('a:srgbClr'))
        if srgb is not None:
            props['bullet_color'] = f"#{srgb.get('val', '000000')}"

    return props


def apply_bullet_properties(paragraph, elem: dict):
    """Apply bullet properties to a paragraph via lxml.

    Reads bullet_char, bullet_font, bullet_size_pct, bullet_color from elem.
    """
    if 'bullet_char' not in elem and 'bullet_none' not in elem:
        return

    pPr = paragraph._p.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, qn('a:pPr'))
        paragraph._p.insert(0, pPr)

    if elem.get('bullet_none'):
        etree.SubElement(pPr, qn('a:buNone'))
        return

    if 'bullet_font' in elem:
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', elem['bullet_font'])

    if 'bullet_size_pct' in elem:
        buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
        buSzPct.set('val', str(elem['bullet_size_pct']))

    if 'bullet_color' in elem:
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', elem['bullet_color'].lstrip('#'))

    if 'bullet_char' in elem:
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', elem['bullet_char'])
