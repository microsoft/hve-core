"""Build a PowerPoint slide deck from YAML content and style definitions.

Usage:
    python build_deck.py --content-dir content/ --style content/global/style.yaml --output slide-deck/presentation.pptx
    python build_deck.py --content-dir content/ --style content/global/style.yaml --source existing.pptx --output slide-deck/presentation.pptx --slides 3,7,15
"""

import argparse
import importlib.util
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
}


def load_yaml(path: Path) -> dict:
    """Load a YAML file and return the parsed dictionary."""
    with open(path, encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def resolve_color(value: str, colors: dict) -> RGBColor:
    """Resolve a color reference ($name) or hex (#RRGGBB) to an RGBColor."""
    if value.startswith("$"):
        key = value[1:]
        value = colors.get(key, "#000000")
    hex_str = value.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def resolve_font(value: str, typography: dict) -> str:
    """Resolve a font reference ($name) or return the literal font name."""
    if value.startswith("$"):
        key = value[1:]
        return typography.get(key, "Segoe UI")
    return value


def merge_styles(global_style: dict, overrides: dict | None) -> dict:
    """Deep-merge per-slide style overrides into the global style."""
    if not overrides:
        return global_style.copy()
    merged = {}
    for key in global_style:
        if key in overrides and isinstance(global_style[key], dict) and isinstance(overrides[key], dict):
            merged[key] = {**global_style[key], **overrides[key]}
        elif key in overrides:
            merged[key] = overrides[key]
        else:
            merged[key] = global_style[key]
    for key in overrides:
        if key not in merged:
            merged[key] = overrides[key]
    return merged


def set_slide_bg(slide, color: RGBColor):
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name="Segoe UI",
                font_size=16, font_color=None, bold=False, italic=False,
                alignment=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    if font_color:
        run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_shape_element(slide, elem, colors, typography):
    """Add a shape element from a content.yaml definition."""
    shape_type = SHAPE_MAP.get(elem.get("shape", "rectangle"), MSO_SHAPE.RECTANGLE)
    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if "fill" in elem:
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolve_color(elem["fill"], colors)
    else:
        shape.fill.background()

    if "line_color" in elem:
        shape.line.color.rgb = resolve_color(elem["line_color"], colors)
        shape.line.width = Pt(elem.get("line_width", 1))
    else:
        shape.line.fill.background()

    if "corner_radius" in elem:
        shape.adjustments[0] = elem["corner_radius"]

    if "text" in elem:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = elem["text"]
        run = p.runs[0]
        run.font.name = resolve_font(elem.get("text_font", "$body_font"), typography)
        run.font.size = Pt(elem.get("text_size", 16))
        if "text_color" in elem:
            run.font.color.rgb = resolve_color(elem["text_color"], colors)
        run.font.bold = elem.get("text_bold", False)

    return shape


def add_image_element(slide, elem, content_dir: Path):
    """Add an image element from a content.yaml definition."""
    img_path = content_dir / elem["path"]
    if not img_path.exists():
        # Fallback: add a text box with the path as placeholder
        add_textbox(slide, elem["left"], elem["top"], elem["width"], elem["height"],
                    f"[Image: {elem['path']}]", font_size=12)
        return None

    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])
    pic = slide.shapes.add_picture(str(img_path), left, top, width, height)
    return pic


def add_rich_text_element(slide, elem, colors, typography):
    """Add a rich text element with mixed font/color segments."""
    txBox = slide.shapes.add_textbox(
        Inches(elem["left"]), Inches(elem["top"]),
        Inches(elem["width"]), Inches(elem["height"])
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    for i, seg in enumerate(elem.get("segments", [])):
        run = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
        run.text = seg["text"]
        run.font.name = resolve_font(seg.get("font", "$body_font"), typography)
        run.font.size = Pt(seg.get("size", 16))
        if "color" in seg:
            run.font.color.rgb = resolve_color(seg["color"], colors)
        run.font.bold = seg.get("bold", False)
        run.font.italic = seg.get("italic", False)

    return txBox


def add_card_element(slide, elem, colors, typography):
    """Add a card panel with optional title bar and bullet content."""
    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])

    # Card background
    card_fill = resolve_color(elem.get("fill", "$bg_card"), colors)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_fill
    if "border_color" in elem:
        shape.line.color.rgb = resolve_color(elem["border_color"], colors)
        shape.line.width = Pt(elem.get("border_width", 1))
    else:
        shape.line.fill.background()

    # Accent bar
    if elem.get("accent_bar"):
        bar_color = resolve_color(elem.get("accent_color", "$accent_blue"), colors)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(elem["left"] + 0.15), Inches(elem["top"] + 0.1),
            Inches(elem["width"] - 0.3), Inches(0.04)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

    # Title
    y_offset = 0.2
    if "title" in elem:
        add_textbox(
            slide, elem["left"] + 0.2, elem["top"] + y_offset,
            elem["width"] - 0.4, 0.4,
            elem["title"],
            font_name=resolve_font("$body_font", typography),
            font_size=elem.get("title_size", 16),
            font_color=resolve_color(elem.get("title_color", "$text_white"), colors),
            bold=elem.get("title_bold", True)
        )
        y_offset += 0.5

    # Content bullets
    for item in elem.get("content", []):
        bullet_text = f"\u2022 {item['bullet']}" if "bullet" in item else item.get("text", "")
        color = resolve_color(item.get("color", "$text_white"), colors)
        add_textbox(
            slide, elem["left"] + 0.2, elem["top"] + y_offset,
            elem["width"] - 0.4, 0.35,
            bullet_text,
            font_name=resolve_font("$body_font", typography),
            font_size=item.get("size", 14),
            font_color=color
        )
        y_offset += 0.35

    return shape


def add_arrow_flow_element(slide, elem, colors, typography):
    """Add a horizontal arrow flow diagram."""
    items = elem.get("items", [])
    if not items:
        return

    total_width = elem["width"]
    item_width = total_width / len(items) - 0.3
    x = elem["left"]

    for item in items:
        color = resolve_color(item.get("color", "$accent_blue"), colors)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(x), Inches(elem["top"]),
            Inches(item_width), Inches(elem["height"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["label"]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = resolve_font("$body_font", typography)
        run.font.size = Pt(14)
        run.font.color.rgb = resolve_color("$text_white", colors)
        run.font.bold = True

        x += item_width + 0.3


def add_numbered_step_element(slide, elem, colors, typography):
    """Add a numbered step with circle, label, and description."""
    number = elem.get("number", 1)
    accent = resolve_color(elem.get("accent_color", "$accent_blue"), colors)

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(elem["left"]), Inches(elem["top"]),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = resolve_font("$body_font", typography)
    run.font.size = Pt(16)
    run.font.color.rgb = resolve_color("$text_white", colors)
    run.font.bold = True

    # Label
    add_textbox(
        slide, elem["left"] + 0.6, elem["top"],
        elem["width"] - 0.6, 0.35,
        elem["label"],
        font_name=resolve_font("$body_font", typography),
        font_size=16,
        font_color=resolve_color("$text_white", colors),
        bold=True
    )

    # Description
    if "description" in elem:
        add_textbox(
            slide, elem["left"] + 0.6, elem["top"] + 0.35,
            elem["width"] - 0.6, 0.4,
            elem["description"],
            font_name=resolve_font("$body_font", typography),
            font_size=14,
            font_color=resolve_color("$text_gray", colors)
        )


def build_slide(prs, slide_content: dict, style: dict, content_dir: Path):
    """Build a single slide from content.yaml data and style context."""
    merged_style = merge_styles(style, slide_content.get("style_overrides"))
    colors = merged_style.get("colors", {})
    typography = merged_style.get("typography", {})

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background
    bg_color = resolve_color(colors.get("bg_dark", "#1B1B1F"), colors)
    set_slide_bg(slide, bg_color)

    # Process elements in order
    for elem in slide_content.get("elements", []):
        elem_type = elem.get("type", "textbox")

        if elem_type == "shape":
            add_shape_element(slide, elem, colors, typography)
        elif elem_type == "textbox":
            font_name = resolve_font(elem.get("font", "$body_font"), typography)
            font_color = resolve_color(elem.get("font_color", "$text_white"), colors) if "font_color" in elem else None
            add_textbox(
                slide, elem["left"], elem["top"], elem["width"], elem["height"],
                elem.get("text", ""),
                font_name=font_name,
                font_size=elem.get("font_size", typography.get("body_size", 16)),
                font_color=font_color,
                bold=elem.get("bold", False),
                italic=elem.get("italic", False),
                alignment=elem.get("alignment")
            )
        elif elem_type == "image":
            add_image_element(slide, elem, content_dir)
        elif elem_type == "rich_text":
            add_rich_text_element(slide, elem, colors, typography)
        elif elem_type == "card":
            add_card_element(slide, elem, colors, typography)
        elif elem_type == "arrow_flow":
            add_arrow_flow_element(slide, elem, colors, typography)
        elif elem_type == "numbered_step":
            add_numbered_step_element(slide, elem, colors, typography)

    # Execute content-extra.py if present
    extra_script = content_dir / "content-extra.py"
    if extra_script.exists():
        spec = importlib.util.spec_from_file_location("content_extra", str(extra_script))
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        if hasattr(mod, "render"):
            mod.render(slide, merged_style, content_dir)

    # Add speaker notes
    notes = slide_content.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def discover_slides(content_dir: Path) -> list[tuple[int, Path]]:
    """Discover slide content directories and return sorted (number, path) pairs."""
    slides = []
    for child in content_dir.iterdir():
        if child.is_dir() and child.name.startswith("slide-"):
            match = re.match(r"slide-(\d+)", child.name)
            if match:
                num = int(match.group(1))
                content_yaml = child / "content.yaml"
                if content_yaml.exists():
                    slides.append((num, child))
    return sorted(slides, key=lambda x: x[0])


def main():
    parser = argparse.ArgumentParser(description="Build a PowerPoint deck from YAML content")
    parser.add_argument("--content-dir", required=True, help="Path to the content/ directory")
    parser.add_argument("--style", required=True, help="Path to the global style.yaml")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    parser.add_argument("--source", help="Source PPTX to update (for partial rebuilds)")
    parser.add_argument("--slides", help="Comma-separated slide numbers to rebuild (requires --source)")
    args = parser.parse_args()

    content_dir = Path(args.content_dir)
    style = load_yaml(Path(args.style))
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    dims = style.get("dimensions", {})
    width = dims.get("width_inches", 13.333)
    height = dims.get("height_inches", 7.5)

    if args.source and args.slides:
        # Partial rebuild: open existing deck and replace specific slides
        prs = Presentation(args.source)
        slide_nums = [int(s.strip()) for s in args.slides.split(",")]
        slides_data = discover_slides(content_dir)
        slides_to_rebuild = {num: path for num, path in slides_data if num in slide_nums}

        for num in slide_nums:
            if num not in slides_to_rebuild:
                print(f"Warning: No content found for slide {num}, skipping")
                continue
            slide_dir = slides_to_rebuild[num]
            slide_content = load_yaml(slide_dir / "content.yaml")
            # Remove the existing slide at this position (0-indexed)
            idx = num - 1
            if idx < len(prs.slides):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
            # Build the new slide
            build_slide(prs, slide_content, style, slide_dir)
            print(f"Rebuilt slide {num}")
    else:
        # Full build
        prs = Presentation()
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)

        slides_data = discover_slides(content_dir)
        if not slides_data:
            print("No slide content found in", content_dir)
            sys.exit(1)

        for num, slide_dir in slides_data:
            slide_content = load_yaml(slide_dir / "content.yaml")
            build_slide(prs, slide_content, style, slide_dir)
            print(f"Built slide {num}: {slide_content.get('title', 'Untitled')}")

    prs.save(str(output_path))
    print(f"\nDeck saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
