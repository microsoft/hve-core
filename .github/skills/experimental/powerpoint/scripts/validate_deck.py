"""Validate a PowerPoint slide deck against content definitions and styling rules.

Usage:
    python validate_deck.py --input slide-deck/presentation.pptx --content-dir content/
"""

import argparse
from pathlib import Path

from pptx import Presentation

from pptx_fonts import font_family_matches
from pptx_utils import emu_to_inches, load_yaml


def check_text_overlay(slide, slide_num: int) -> list[str]:
    """Check for overlapping text elements on a slide."""
    issues = []
    text_elements = []

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            left = emu_to_inches(shape.left)
            top = emu_to_inches(shape.top)
            width = emu_to_inches(shape.width)
            height = emu_to_inches(shape.height)
            bottom = top + height
            right = left + width
            text_elements.append({
                "name": shape.name,
                "left": left,
                "right": right,
                "top": top,
                "bottom": bottom,
                "text": shape.text_frame.text[:50],
            })

    # Sort by vertical position
    text_elements.sort(key=lambda x: x["top"])

    for i in range(len(text_elements) - 1):
        current = text_elements[i]
        next_elem = text_elements[i + 1]
        # Only flag overlap when elements share horizontal space
        h_overlap = current["left"] < next_elem["right"] and next_elem["left"] < current["right"]
        if not h_overlap:
            continue
        gap = next_elem["top"] - current["bottom"]
        if gap < 0:
            issues.append(
                f"Slide {slide_num}: Text overlay between "
                f"'{current['name']}' (bottom={current['bottom']:.2f}) and "
                f"'{next_elem['name']}' (top={next_elem['top']:.2f}), "
                f"gap={gap:.2f}\""
            )

    return issues


def check_width_overflow(slide, slide_num: int, max_width: float = 13.333) -> list[str]:
    """Check for elements extending beyond the slide width."""
    issues = []
    for shape in slide.shapes:
        left = emu_to_inches(shape.left)
        width = emu_to_inches(shape.width)
        right = left + width
        if right > max_width + 0.01:
            issues.append(
                f"Slide {slide_num}: Width overflow on '{shape.name}' "
                f"(left={left:.2f} + width={width:.2f} = {right:.2f} > {max_width})"
            )
    return issues


def check_speaker_notes(slide, slide_num: int) -> list[str]:
    """Check for missing speaker notes."""
    issues = []
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes:
            issues.append(f"Slide {slide_num}: Empty speaker notes")
    except (AttributeError, TypeError):
        issues.append(f"Slide {slide_num}: Missing speaker notes")
    return issues


def check_font_consistency(slide, slide_num: int, expected_fonts: set[str] | None = None) -> list[str]:
    """Check for unexpected or inconsistent fonts."""
    issues = []
    if not expected_fonts:
        return issues

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and not font_family_matches(run.font.name, expected_fonts):
                        issues.append(
                            f"Slide {slide_num}: Unexpected font '{run.font.name}' "
                            f"in '{shape.name}' (expected: {expected_fonts})"
                        )
    return issues


def check_height_overflow(slide, slide_num: int, max_height: float = 7.5) -> list[str]:
    """Check for elements extending beyond the slide height."""
    issues = []
    for shape in slide.shapes:
        top = emu_to_inches(shape.top)
        height = emu_to_inches(shape.height)
        bottom = top + height
        if bottom > max_height + 0.01:
            issues.append(
                f"Slide {slide_num}: Height overflow on '{shape.name}' "
                f"(top={top:.2f} + height={height:.2f} = {bottom:.2f} > {max_height})"
            )
    return issues


def check_edge_margins(slide, slide_num: int, max_width: float = 13.333,
                       max_height: float = 7.5, min_margin: float = 0.5,
                       strict: bool = False) -> list[str]:
    """Check for elements too close to slide edges."""
    issues = []
    for shape in slide.shapes:
        left = emu_to_inches(shape.left)
        top = emu_to_inches(shape.top)
        width = emu_to_inches(shape.width)
        height = emu_to_inches(shape.height)

        # Skip full-bleed elements (backgrounds, banners) in non-strict mode
        if not strict:
            if width >= max_width * 0.95 or height >= max_height * 0.95:
                continue

        if left < min_margin:
            issues.append(f"Slide {slide_num}: '{shape.name}' too close to left edge (left={left:.2f}, min={min_margin})")
        if top < min_margin:
            issues.append(f"Slide {slide_num}: '{shape.name}' too close to top edge (top={top:.2f}, min={min_margin})")
        right = left + width
        if right > max_width - min_margin:
            issues.append(f"Slide {slide_num}: '{shape.name}' too close to right edge (right={right:.2f})")
        bottom = top + height
        if bottom > max_height - min_margin:
            issues.append(f"Slide {slide_num}: '{shape.name}' too close to bottom edge (bottom={bottom:.2f})")
    return issues


def check_element_spacing(slide, slide_num: int, min_spacing: float = 0.3,
                          max_width: float = 13.333, max_height: float = 7.5) -> list[str]:
    """Check for insufficient spacing between adjacent elements."""
    issues = []
    elements = []
    for shape in slide.shapes:
        elements.append({
            "name": shape.name,
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "right": emu_to_inches(shape.left) + emu_to_inches(shape.width),
            "bottom": emu_to_inches(shape.top) + emu_to_inches(shape.height),
        })

    # Filter out full-slide background elements that overlap everything
    elements = [e for e in elements
                if not (e["right"] - e["left"] >= max_width * 0.95
                        and e["bottom"] - e["top"] >= max_height * 0.95)]

    for i in range(len(elements)):
        for j in range(i + 1, len(elements)):
            a, b = elements[i], elements[j]
            h_overlap = a["left"] < b["right"] and b["left"] < a["right"]
            v_overlap = a["top"] < b["bottom"] and b["top"] < a["bottom"]

            if h_overlap and v_overlap:
                # Elements actually overlap — already caught by overlay check
                continue

            # Check vertical gap for horizontally overlapping elements
            if h_overlap:
                v_gap = max(b["top"] - a["bottom"], a["top"] - b["bottom"])
                if 0 < v_gap < min_spacing:
                    issues.append(
                        f"Slide {slide_num}: Insufficient vertical spacing ({v_gap:.2f}\") "
                        f"between '{a['name']}' and '{b['name']}' (min: {min_spacing}\")"
                    )

            # Check horizontal gap for vertically overlapping elements
            if v_overlap:
                h_gap = max(b["left"] - a["right"], a["left"] - b["right"])
                if 0 < h_gap < min_spacing:
                    issues.append(
                        f"Slide {slide_num}: Insufficient horizontal spacing ({h_gap:.2f}\") "
                        f"between '{a['name']}' and '{b['name']}' (min: {min_spacing}\")"
                    )
    return issues


def _luminance_contrast(color1, color2) -> float:
    """Calculate WCAG luminance contrast ratio between two RGBColor values."""
    def relative_luminance(rgb):
        r, g, b = rgb[0] / 255, rgb[1] / 255, rgb[2] / 255
        r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
        g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
        b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    l1 = relative_luminance(color1)
    l2 = relative_luminance(color2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def check_color_contrast(slide, slide_num: int) -> list[str]:
    """Check for low-contrast text on shapes."""
    issues = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    text_color = run.font.color.rgb
                    if text_color is None:
                        continue
                except (AttributeError, TypeError):
                    continue

                # Check against shape fill
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        bg_color = fill.fore_color.rgb
                        if bg_color and text_color:
                            contrast = _luminance_contrast(text_color, bg_color)
                            if contrast < 3.0:
                                issues.append(
                                    f"Slide {slide_num}: Low contrast ({contrast:.1f}:1) in "
                                    f"'{shape.name}' — text #{text_color} on #{bg_color}"
                                )
                except (AttributeError, TypeError):
                    pass
    return issues


def check_narrow_text_boxes(slide, slide_num: int) -> list[str]:
    """Check for text boxes too narrow for their content."""
    issues = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            width = emu_to_inches(shape.width)
            text = shape.text_frame.text
            if width < 1.5 and len(text) > 30:
                issues.append(
                    f"Slide {slide_num}: Narrow text box '{shape.name}' "
                    f"(width={width:.2f}\", text length={len(text)})"
                )
    return issues


def check_leftover_placeholders(slide, slide_num: int) -> list[str]:
    """Check for placeholder text that wasn't replaced."""
    issues = []
    placeholder_patterns = ["Click to add", "click to edit", "Add title", "Add subtitle"]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            for pattern in placeholder_patterns:
                if pattern.lower() in text.lower():
                    issues.append(
                        f"Slide {slide_num}: Leftover placeholder text in '{shape.name}': '{text[:50]}'"
                    )
                    break
    return issues


def validate_deck(pptx_path: Path, content_dir: Path | None = None,
                   slide_filter: set[int] | None = None,
                   strict: bool = False) -> list[str]:
    """Run all validation checks on a PPTX file."""
    prs = Presentation(str(pptx_path))
    all_issues = []

    # Load expected fonts from style if available
    expected_fonts = None
    if content_dir:
        style_path = content_dir / "global" / "style.yaml"
        if style_path.exists():
            style = load_yaml(style_path)
            typo = style.get("typography", {})
            expected_fonts = set()
            if "body_font" in typo:
                expected_fonts.add(typo["body_font"])
            if "code_font" in typo:
                expected_fonts.add(typo["code_font"])

    max_width = emu_to_inches(prs.slide_width)
    max_height = emu_to_inches(prs.slide_height)

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue
        all_issues.extend(check_text_overlay(slide, slide_num))
        all_issues.extend(check_width_overflow(slide, slide_num, max_width))
        all_issues.extend(check_speaker_notes(slide, slide_num))
        all_issues.extend(check_font_consistency(slide, slide_num, expected_fonts))
        all_issues.extend(check_height_overflow(slide, slide_num, max_height))
        all_issues.extend(check_edge_margins(slide, slide_num, max_width, max_height, strict=strict))
        all_issues.extend(check_element_spacing(slide, slide_num, max_width=max_width, max_height=max_height))
        all_issues.extend(check_color_contrast(slide, slide_num))
        all_issues.extend(check_narrow_text_boxes(slide, slide_num))
        all_issues.extend(check_leftover_placeholders(slide, slide_num))

    # Check slide count against content (only when all slides have content dirs)
    if content_dir and not slide_filter:
        slide_dirs = sorted(
            [d for d in content_dir.iterdir() if d.is_dir() and d.name.startswith("slide-")]
        )
        if len(slide_dirs) == len(prs.slides):
            pass  # Full content coverage, no mismatch
        elif len(slide_dirs) < len(prs.slides):
            # Partial content is expected during incremental updates
            all_issues.append(
                f"Info: Partial content detected — {len(prs.slides)} slides in PPTX, "
                f"{len(slide_dirs)} content directories (expected for incremental updates)"
            )
        else:
            all_issues.append(
                f"Slide count mismatch: {len(prs.slides)} slides in PPTX, "
                f"{len(slide_dirs)} content directories"
            )

    return all_issues


def main():
    parser = argparse.ArgumentParser(description="Validate a PowerPoint deck")
    parser.add_argument("--input", required=True, help="Input PPTX file path")
    parser.add_argument("--content-dir", help="Content directory for comparison")
    parser.add_argument("--slides", help="Comma-separated slide numbers to validate (default: all)")
    parser.add_argument("--strict", action="store_true",
                        help="Enable strict validation (flag full-bleed elements in edge checks)")
    args = parser.parse_args()

    pptx_path = Path(args.input)
    content_dir = Path(args.content_dir) if args.content_dir else None

    slide_filter = None
    if args.slides:
        slide_filter = {int(s.strip()) for s in args.slides.split(",")}

    print(f"Validating: {pptx_path}")
    issues = validate_deck(pptx_path, content_dir, slide_filter=slide_filter, strict=args.strict)

    if issues:
        print(f"\n{len(issues)} issue(s) found:\n")
        for issue in issues:
            print(f"  - {issue}")
        return 1
    else:
        print("\nAll validation checks passed.")
        return 0


if __name__ == "__main__":
    exit(main())
