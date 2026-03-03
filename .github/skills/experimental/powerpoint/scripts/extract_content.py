"""Extract content from an existing PPTX into YAML content and style definitions.

Usage:
    python extract_content.py --input existing-deck.pptx --output-dir content/
    python extract_content.py --input existing-deck.pptx --output-dir content/ --slides 3,7,15
"""

import argparse
from collections import Counter
from pathlib import Path

import yaml
from pptx import Presentation

from pptx_colors import extract_color, hex_brightness, rgb_to_hex
from pptx_fills import extract_fill, extract_line
from pptx_fonts import (
    extract_alignment,
    extract_font_info,
    extract_paragraph_font,
    normalize_font_family,
)
from pptx_shapes import AUTO_SHAPE_NAME_MAP, extract_rotation
from pptx_text import (
    extract_paragraph_properties,
    extract_run_properties,
    extract_text_frame_properties,
)
from pptx_utils import emu_to_inches


def extract_connector(shape) -> dict:
    """Extract a connector element definition."""
    elem = {
        "type": "connector",
        "begin_x": emu_to_inches(shape.begin_x),
        "begin_y": emu_to_inches(shape.begin_y),
        "end_x": emu_to_inches(shape.end_x),
        "end_y": emu_to_inches(shape.end_y),
        "name": shape.name,
    }
    line_props = extract_line(shape)
    if line_props:
        elem.update(line_props)
    return elem


def extract_group(shape, slide_num: int, output_dir, img_count: int) -> dict:
    """Extract a group shape and its nested child elements."""
    elem = {
        "type": "group",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
        "elements": [],
    }
    for child in shape.shapes:
        child_elem = extract_child_shape(child, slide_num, output_dir, img_count)
        if child_elem:
            elem["elements"].append(child_elem)
    return elem


def extract_child_shape(shape, slide_num: int, output_dir, img_count: int) -> dict | None:
    """Extract a single child shape within a group."""
    shape_type = shape.shape_type
    if shape_type == 13:  # PICTURE
        return extract_image(shape, output_dir, slide_num, img_count)
    elif shape_type == 17:  # TEXT_BOX
        return extract_textbox(shape)
    elif shape_type == 1:  # AUTO_SHAPE
        return extract_shape(shape)
    elif shape_type == 9:  # LINE / CONNECTOR
        return extract_connector(shape)
    elif shape_type == 6:  # Nested GROUP
        return extract_group(shape, slide_num, output_dir, img_count)
    elif hasattr(shape, "has_table") and shape.has_table:
        from pptx_tables import extract_table
        return extract_table(shape)
    elif hasattr(shape, "has_chart") and shape.has_chart:
        from pptx_charts import extract_chart
        return extract_chart(shape)
    return {
        "type": "shape",
        "shape": "rectangle",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
        "_unrecognized_shape_type": int(shape_type),
    }


def extract_shape(shape) -> dict:
    """Extract a shape element definition."""
    elem = {
        "type": "shape",
        "shape": "rectangle",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }

    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    # Detect shape type from auto_shape_type enum
    try:
        elem["shape"] = AUTO_SHAPE_NAME_MAP.get(shape.auto_shape_type, "rectangle")
    except (AttributeError, TypeError):
        elem["shape"] = "rectangle"

    # Extract fill
    try:
        fill_result = extract_fill(shape.fill)
        if fill_result is not None:
            elem["fill"] = fill_result
    except (AttributeError, TypeError):
        pass

    # Extract line properties
    line_props = extract_line(shape)
    if line_props:
        elem.update(line_props)

    # Extract text if present
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            elem["text"] = text

            # Extract text frame-level properties
            tf_props = extract_text_frame_properties(shape.text_frame)
            if tf_props:
                elem.update(tf_props)

            # Extract text styling: try run-level first, fall back to paragraph-level
            for para in shape.text_frame.paragraphs:
                run_info = {}
                for run in para.runs:
                    run_info = extract_font_info(run.font)
                    run_info.update(extract_run_properties(run))
                    break
                para_info = extract_paragraph_font(para)
                para_spacing = extract_paragraph_properties(para)
                # Merge: run-level wins, paragraph-level fills gaps
                merged = {**para_info, **run_info}
                if "font" in merged:
                    elem["text_font"] = merged["font"]
                if "size" in merged:
                    elem["text_size"] = merged["size"]
                if "color" in merged:
                    elem["text_color"] = merged["color"]
                if merged.get("bold"):
                    elem["text_bold"] = True
                if merged.get("underline"):
                    elem["underline"] = True
                if merged.get("hyperlink"):
                    elem["hyperlink"] = merged["hyperlink"]
                if para_spacing:
                    elem.update(para_spacing)
                break

    return elem


def extract_textbox(shape) -> dict:
    """Extract a text box element definition."""
    elem = {
        "type": "textbox",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "text": shape.text_frame.text.strip() if shape.has_text_frame else "",
        "name": shape.name,
    }

    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    # Check if this is a rich text element (multiple runs with different formatting)
    if shape.has_text_frame:
        # Extract text frame-level properties (margins, auto_size, vertical_anchor)
        tf_props = extract_text_frame_properties(shape.text_frame)
        if tf_props:
            elem.update(tf_props)

        runs = []
        para_info = {}
        alignment = None
        para_spacing = {}
        for para in shape.text_frame.paragraphs:
            if not para_info:
                para_info = extract_paragraph_font(para)
            if alignment is None:
                alignment = extract_alignment(para)
            if not para_spacing:
                para_spacing = extract_paragraph_properties(para)
            for run in para.runs:
                font_info = extract_font_info(run.font)
                run_extra = extract_run_properties(run)
                runs.append({"text": run.text, **font_info, **run_extra})

        # If multiple runs with different formatting, mark as rich_text
        if len(runs) > 1:
            fonts = {r.get("font") for r in runs if "font" in r}
            colors = {r.get("color") for r in runs if "color" in r}
            if len(fonts) > 1 or len(colors) > 1:
                elem["type"] = "rich_text"
                elem["segments"] = runs
                del elem["text"]
                if alignment:
                    elem["alignment"] = alignment
                if para_spacing:
                    elem.update(para_spacing)
                return elem

        # Single-style text box: merge run-level and paragraph-level font info
        merged = {**para_info}
        if runs:
            # Run-level properties override paragraph-level
            for key, val in runs[0].items():
                if key != "text" and val is not None:
                    merged[key] = val

        if "font" in merged:
            elem["font"] = merged["font"]
        if "size" in merged:
            elem["font_size"] = merged["size"]
        if "color" in merged:
            elem["font_color"] = merged["color"]
        if merged.get("bold"):
            elem["font_bold"] = True
        if merged.get("italic"):
            elem["italic"] = True
        if merged.get("underline"):
            elem["underline"] = True
        if merged.get("hyperlink"):
            elem["hyperlink"] = merged["hyperlink"]
        if alignment:
            elem["alignment"] = alignment
        if para_spacing:
            elem.update(para_spacing)

    return elem


def extract_image(shape, output_dir: Path, slide_num: int, img_count: int) -> dict:
    """Extract an image element and save the image file."""
    try:
        img = shape.image
    except ValueError:
        # Linked images have no embedded blob
        elem = {
            "type": "image",
            "path": "LINKED_IMAGE_NOT_EMBEDDED",
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
            "name": shape.name,
            "_note": "Image was linked, not embedded in the PPTX",
        }
        rot = extract_rotation(shape)
        if rot is not None:
            elem["rotation"] = rot
        return elem

    ext = img.content_type.split("/")[-1]
    if ext == "jpeg":
        ext = "jpg"

    img_name = f"image-{img_count:02d}.{ext}"
    img_path = output_dir / "images" / img_name
    img_path.parent.mkdir(parents=True, exist_ok=True)

    with open(img_path, "wb") as f:
        f.write(img.blob)

    elem = {
        "type": "image",
        "path": f"images/{img_name}",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }
    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot
    return elem


def detect_global_style(prs) -> dict:
    """Analyze the presentation to detect common styling patterns."""
    bg_colors = Counter()
    text_colors = Counter()
    accent_colors = Counter()
    fill_colors = Counter()
    font_names = Counter()
    font_sizes = Counter()

    for slide in prs.slides:
        # Detect background colors
        try:
            fill_result = extract_fill(slide.background.fill)
            if isinstance(fill_result, str):
                bg_colors[fill_result] += 1
        except (AttributeError, TypeError):
            pass

        for shape in slide.shapes:
            # Collect fill colors
            try:
                fill_result = extract_fill(shape.fill)
                if isinstance(fill_result, str):
                    # Thin horizontal bars are accent colors
                    h = emu_to_inches(shape.height)
                    if h < 0.1:
                        accent_colors[fill_result] += 1
                    else:
                        fill_colors[fill_result] += 1
            except (AttributeError, TypeError):
                pass

            # Collect font information
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            # Separate font family from weight suffix
                            base_name = normalize_font_family(run.font.name)
                            font_names[base_name] += 1
                        if run.font.size:
                            font_sizes[int(run.font.size.pt)] += 1
                        try:
                            color = extract_color(run.font.color)
                            if isinstance(color, str) and color.startswith("#"):
                                text_colors[color] += 1
                        except (AttributeError, TypeError):
                            pass

    # Build color map from frequency analysis
    colors = {}
    if bg_colors:
        colors["bg_dark"] = bg_colors.most_common(1)[0][0]
    if fill_colors:
        colors["bg_card"] = fill_colors.most_common(1)[0][0]

    # Assign text colors by brightness
    for color_hex, _count in text_colors.most_common(5):
        brightness = hex_brightness(color_hex)
        if brightness > 200 and "text_white" not in colors:
            colors["text_white"] = color_hex
        elif brightness < 80 and "text_dark" not in colors:
            colors["text_dark"] = color_hex
        elif 80 <= brightness <= 200 and "text_gray" not in colors:
            colors["text_gray"] = color_hex

    # Assign accent colors from thin bars
    accent_names = ["accent_blue", "accent_teal", "accent_green"]
    for i, (color_hex, _count) in enumerate(accent_colors.most_common(3)):
        if i < len(accent_names):
            colors[accent_names[i]] = color_hex

    # Determine primary fonts (base family, not weight-specific)
    body_font = "Segoe UI"
    code_font = "Cascadia Code"
    for f, _count in font_names.most_common():
        if any(kw in f.lower() for kw in ("cascadia", "consolas", "mono", "courier")):
            code_font = f
        else:
            body_font = f
            break

    # Determine font sizes using frequency-based median for body, 75th percentile for heading
    heading_size = 28
    body_size = 16
    if font_sizes:
        # Filter outliers: ignore sizes below 8pt and above 60pt
        filtered = {s: c for s, c in font_sizes.items() if 8 < s < 60}
        if filtered:
            sorted_sizes = sorted(filtered.keys())
            body_size = sorted_sizes[len(sorted_sizes) // 2]
            heading_size = sorted_sizes[int(len(sorted_sizes) * 0.85)]

    style = {
        "dimensions": {
            "width_inches": emu_to_inches(prs.slide_width),
            "height_inches": emu_to_inches(prs.slide_height),
            "format": "16:9",
        },
        "colors": colors,
        "typography": {
            "body_font": body_font,
            "code_font": code_font,
            "heading_size": heading_size,
            "body_size": body_size,
        },
        "defaults": {
            "speaker_notes_required": True,
        },
    }

    # Extract presentation metadata
    metadata = {}
    props = prs.core_properties
    for attr in ("title", "author", "subject", "keywords", "description", "category"):
        val = getattr(props, attr, None)
        if val:
            metadata[attr] = val
    if metadata:
        style["metadata"] = metadata

    return style


def extract_slide(slide, slide_num: int, output_dir: Path) -> dict:
    """Extract all elements from a slide into a content.yaml structure."""
    slide_dir = output_dir / f"slide-{slide_num:03d}"
    slide_dir.mkdir(parents=True, exist_ok=True)

    content = {
        "slide": slide_num,
        "title": "",
        "elements": [],
    }

    # Extract layout name
    try:
        layout_name = slide.slide_layout.name
        if layout_name:
            content["layout"] = layout_name
    except (AttributeError, TypeError):
        pass

    # Extract slide background
    try:
        if not slide.follow_master_background:
            fill_result = extract_fill(slide.background.fill)
            if fill_result is not None:
                content["background"] = {"fill": fill_result}
    except (AttributeError, TypeError):
        pass

    # Extract speaker notes (include empty string when notes slide exists)
    try:
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            content["speaker_notes"] = notes
    except (AttributeError, TypeError):
        pass

    img_count = 0
    for shape in slide.shapes:
        shape_type = shape.shape_type

        if shape_type == 13:  # PICTURE
            img_count += 1
            elem = extract_image(shape, slide_dir, slide_num, img_count)
            content["elements"].append(elem)
        elif shape_type == 17:  # TEXT_BOX
            elem = extract_textbox(shape)
            content["elements"].append(elem)
            # Detect title (first large text near top)
            if not content["title"] and emu_to_inches(shape.top) < 1.5:
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if text and len(text) < 100:
                    content["title"] = text
        elif shape_type == 1:  # AUTO_SHAPE
            elem = extract_shape(shape)
            content["elements"].append(elem)
        elif shape_type == 14:  # PLACEHOLDER
            # Placeholders may contain text; extract like textboxes
            if shape.has_text_frame:
                elem = extract_textbox(shape)
                elem["_placeholder"] = True
                content["elements"].append(elem)
        elif shape_type == 6:  # GROUP
            elem = extract_group(shape, slide_num, slide_dir, img_count)
            content["elements"].append(elem)
        elif shape_type == 9:  # LINE / CONNECTOR
            elem = extract_connector(shape)
            content["elements"].append(elem)
        elif hasattr(shape, "has_table") and shape.has_table:
            from pptx_tables import extract_table
            elem = extract_table(shape)
            content["elements"].append(elem)
        elif hasattr(shape, "has_chart") and shape.has_chart:
            from pptx_charts import extract_chart
            elem = extract_chart(shape)
            content["elements"].append(elem)
        else:
            # Log unrecognized shape types for manual review
            content["elements"].append({
                "type": "shape",
                "shape": "rectangle",
                "left": emu_to_inches(shape.left),
                "top": emu_to_inches(shape.top),
                "width": emu_to_inches(shape.width),
                "height": emu_to_inches(shape.height),
                "name": shape.name,
                "_unrecognized_shape_type": int(shape_type),
            })

    return content, slide_dir


def main():
    parser = argparse.ArgumentParser(description="Extract content from a PPTX into YAML")
    parser.add_argument("--input", required=True, help="Input PPTX file path")
    parser.add_argument("--output-dir", required=True, help="Output content directory")
    parser.add_argument("--slides", help="Comma-separated slide numbers to extract (default: all)")
    args = parser.parse_args()

    pptx_path = Path(args.input)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    slide_filter = None
    if args.slides:
        slide_filter = {int(s.strip()) for s in args.slides.split(",")}

    prs = Presentation(str(pptx_path))
    print(f"Extracting from: {pptx_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Dimensions: {emu_to_inches(prs.slide_width)}\" x {emu_to_inches(prs.slide_height)}\"")

    # Detect and save global style
    global_style = detect_global_style(prs)
    global_dir = output_dir / "global"
    global_dir.mkdir(parents=True, exist_ok=True)
    style_path = global_dir / "style.yaml"
    with open(style_path, "w", encoding="utf-8") as f:
        yaml.dump(global_style, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
    print(f"Global style saved to {style_path}")

    # Extract slides (filtered or all)
    extracted = 0
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue
        content, slide_dir = extract_slide(slide, slide_num, output_dir)

        content_path = slide_dir / "content.yaml"
        with open(content_path, "w", encoding="utf-8") as f:
            yaml.dump(content, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
        print(f"Slide {slide_num}: {content.get('title', 'Untitled')} -> {content_path}")
        extracted += 1

    print(f"\nExtraction complete. {extracted} slide(s) extracted to {output_dir}")


if __name__ == "__main__":
    main()
