"""Tests for pptx_fonts module."""

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx_fonts import (
    ALIGNMENT_MAP,
    ALIGNMENT_REVERSE_MAP,
    FONT_WEIGHT_SUFFIXES,
    _extract_char_spacing,
    extract_alignment,
    extract_font_info,
    extract_paragraph_font,
    font_family_matches,
    normalize_font_family,
    resolve_font,
)


class TestResolveFont:
    """Tests for resolve_font."""

    def test_returns_literal_name(self):
        assert resolve_font("Arial") == "Arial"

    def test_returns_theme_ref(self):
        assert resolve_font("+mj-lt") == "+mj-lt"

    def test_with_typography_dict(self):
        assert resolve_font("Calibri", {"heading": "Arial"}) == "Calibri"


class TestNormalizeFontFamily:
    """Tests for normalize_font_family."""

    def test_no_suffix(self):
        assert normalize_font_family("Segoe UI") == "Segoe UI"

    def test_semibold(self):
        assert normalize_font_family("Segoe UI Semibold") == "Segoe UI"

    def test_semibold_camelcase(self):
        assert normalize_font_family("Segoe UI SemiBold") == "Segoe UI"

    def test_bold(self):
        assert normalize_font_family("Arial Bold") == "Arial"

    def test_light(self):
        assert normalize_font_family("Segoe UI Light") == "Segoe UI"

    def test_thin(self):
        assert normalize_font_family("Roboto Thin") == "Roboto"

    def test_black(self):
        assert normalize_font_family("Montserrat Black") == "Montserrat"

    def test_medium(self):
        assert normalize_font_family("Inter Medium") == "Inter"

    def test_extrabold(self):
        assert normalize_font_family("Open Sans ExtraBold") == "Open Sans"

    def test_extralight(self):
        assert normalize_font_family("Noto Sans ExtraLight") == "Noto Sans"


class TestFontFamilyMatches:
    """Tests for font_family_matches."""

    def test_exact_match(self):
        assert font_family_matches("Arial", {"Arial", "Segoe UI"}) is True

    def test_no_match(self):
        assert font_family_matches("Comic Sans", {"Arial", "Segoe UI"}) is False

    def test_weight_variant_match(self):
        assert font_family_matches("Segoe UI Semibold", {"Segoe UI"}) is True

    def test_expected_has_weight_variant(self):
        assert font_family_matches("Segoe UI", {"Segoe UI Bold"}) is True

    def test_different_base_no_match(self):
        assert font_family_matches("Arial Bold", {"Segoe UI"}) is False

    def test_empty_expected(self):
        assert font_family_matches("Arial", set()) is False


class TestExtractFontInfo:
    """Tests for extract_font_info with real python-pptx font objects."""

    def test_name_and_size(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(24)
        info = extract_font_info(run.font)
        assert info["font"] == "Arial"
        assert info["size"] == 24

    def test_bold(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        info = extract_font_info(run.font)
        assert info["bold"] is True

    def test_italic(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.italic = True
        info = extract_font_info(run.font)
        assert info["italic"] is True

    def test_underline(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.underline = True
        info = extract_font_info(run.font)
        assert info["underline"] is True

    def test_color(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        info = extract_font_info(run.font)
        assert info["color"] == "#FF0000"

    def test_missing_properties_omitted(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        run = txBox.text_frame.paragraphs[0].runs[0]
        # Don't set any properties explicitly
        info = extract_font_info(run.font)
        # Only keys for set values should appear
        assert "bold" not in info or info.get("bold") is not True


class TestExtractCharSpacing:
    """Tests for _extract_char_spacing."""

    def test_with_spacing(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        rpr = run.font._element
        rpr.set("spc", "200")
        result = _extract_char_spacing(run.font)
        assert result == 2.0

    def test_without_spacing(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        run = txBox.text_frame.paragraphs[0].runs[0]
        result = _extract_char_spacing(run.font)
        assert result is None

    def test_negative_spacing(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        rpr = run.font._element
        rpr.set("spc", "-100")
        result = _extract_char_spacing(run.font)
        assert result == -1.0


class TestExtractParagraphFont:
    """Tests for extract_paragraph_font."""

    def test_paragraph_level_font(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.font.name = "Calibri"
        para.font.size = Pt(16)
        para.font.bold = True
        info = extract_paragraph_font(para)
        assert info["font"] == "Calibri"
        assert info["size"] == 16
        assert info["bold"] is True

    def test_empty_paragraph(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = ""
        para = txBox.text_frame.paragraphs[0]
        info = extract_paragraph_font(para)
        # Without explicit properties, dict should be empty or minimal
        assert isinstance(info, dict)


class TestExtractAlignment:
    """Tests for extract_alignment."""

    def test_left(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        assert extract_alignment(para) == "left"

    def test_center(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        assert extract_alignment(para) == "center"

    def test_right(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        assert extract_alignment(para) == "right"

    def test_justify(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.JUSTIFY
        assert extract_alignment(para) == "justify"

    def test_none(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        para = txBox.text_frame.paragraphs[0]
        assert extract_alignment(para) is None


class TestConstants:
    """Tests for module constants."""

    def test_alignment_map_keys(self):
        assert set(ALIGNMENT_MAP.keys()) == {"left", "center", "right", "justify"}

    def test_alignment_reverse_map_values(self):
        expected = {"left", "center", "right", "justify"}
        assert set(ALIGNMENT_REVERSE_MAP.values()) == expected

    def test_font_weight_suffixes_not_empty(self):
        assert len(FONT_WEIGHT_SUFFIXES) > 0
        for suffix in FONT_WEIGHT_SUFFIXES:
            assert suffix.startswith(" ")
