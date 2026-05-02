#!/usr/bin/env python3
# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Embed per-slide WAV voice-over files into a PowerPoint deck.

Reads slide-NNN.wav files from an audio directory and adds them as embedded
media objects in the corresponding slides of a PPTX file. Adds animation
timing XML so PowerPoint recognizes the audio as narrations, enabling
'Use Recorded Timings and Narrations' in File > Export > Create a Video.

Usage:
    python embed_audio.py --input deck.pptx --audio-dir voice-over
    python embed_audio.py --input deck.pptx --audio-dir voice-over \
        --output deck-narrated.pptx
"""

from __future__ import annotations

import argparse
import logging
import sys
import wave
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Inches

logger = logging.getLogger(__name__)

EXIT_SUCCESS = 0
EXIT_FAILURE = 1
EXIT_ERROR = 2

AUDIO_MIME_TYPE = "audio/wav"
ICON_SIZE = Inches(0.1)
TIMING_BUFFER_MS = 1500


def get_wav_duration_ms(wav_path: Path) -> int:
    """Return WAV file duration in milliseconds with buffer."""
    with wave.open(str(wav_path), "rb") as wf:
        frames = wf.getnframes()
        rate = wf.getframerate()
        return int((frames / float(rate)) * 1000) + TIMING_BUFFER_MS


def _add_narration_timing(slide: Slide, shape_id: int, duration_ms: int) -> None:
    """Add auto-play narration timing XML to a slide.

    Creates the p:timing element structure that PowerPoint generates
    when using Record Slide Show, enabling 'Use Recorded Timings and
    Narrations' in video export.
    """
    existing = slide._element.find(qn("p:timing"))
    if existing is not None:
        slide._element.remove(existing)

    timing_xml = (
        '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<p:tnLst><p:par>"
        '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        "<p:childTnLst>"
        '<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        "<p:childTnLst><p:par>"
        '<p:cTn id="3" fill="hold">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        "<p:childTnLst><p:par>"
        '<p:cTn id="4" fill="hold">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        "<p:childTnLst>"
        '<p:cmd type="call" cmd="playFrom(0)"><p:cBhvr>'
        f'<p:cTn id="5" dur="{duration_ms}" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        "</p:cBhvr></p:cmd>"
        "</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>"
        "</p:childTnLst></p:cTn>"
        "<p:prevCondLst>"
        '<p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        "</p:prevCondLst>"
        "<p:nextCondLst>"
        '<p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        "</p:nextCondLst>"
        "</p:seq></p:childTnLst></p:cTn>"
        "</p:par></p:tnLst></p:timing>"
    )
    slide._element.append(etree.fromstring(timing_xml))


def _set_slide_transition(slide: Slide, duration_ms: int) -> None:
    """Set slide auto-advance timing after audio duration."""
    existing = slide._element.find(qn("p:transition"))
    if existing is not None:
        slide._element.remove(existing)

    transition = slide._element.makeelement(
        qn("p:transition"),
        {"advClick": "1", "advTm": str(duration_ms)},
    )
    timing = slide._element.find(qn("p:timing"))
    if timing is not None:
        timing.addprevious(transition)
    else:
        slide._element.append(transition)


def embed_slide_audio(slide: Slide, wav_path: Path) -> bool:
    """Embed a WAV file into a slide as a media object.

    Adds narration timing XML and slide auto-advance so PowerPoint
    recognizes the audio for video export.

    Returns True on success, False on failure.
    """
    try:
        movie_shape = slide.shapes.add_movie(
            str(wav_path),
            left=0,
            top=0,
            width=ICON_SIZE,
            height=ICON_SIZE,
            mime_type=AUDIO_MIME_TYPE,
        )
        shape_id: int | None = movie_shape.shape_id
        if shape_id is not None:
            duration_ms = get_wav_duration_ms(wav_path)
            _add_narration_timing(slide, shape_id, duration_ms)
            _set_slide_transition(slide, duration_ms)
        else:
            # Remove the orphaned audio shape to avoid partial state
            try:
                movie_shape._element.getparent().remove(movie_shape._element)
            except Exception:
                logger.debug(
                    "Could not remove orphaned shape for %s",
                    wav_path.name,
                    exc_info=True,
                )
            logger.error(
                "Could not find audio shape for %s; removed orphaned embed",
                wav_path.name,
            )
            return False
        return True
    except Exception as exc:  # python-pptx raises varied internal exceptions
        logger.exception(
            "Failed to embed audio %s (%s)", wav_path.name, type(exc).__name__
        )
        return False


def create_parser() -> argparse.ArgumentParser:
    """Create and configure the argument parser."""
    parser = argparse.ArgumentParser(
        description="Embed per-slide WAV voice-over files into a PPTX deck"
    )
    parser.add_argument(
        "--input",
        type=Path,
        required=True,
        help="Source PPTX file path",
    )
    parser.add_argument(
        "--audio-dir",
        type=Path,
        default=Path("voice-over"),
        help="Directory containing slide-NNN.wav files (default: voice-over)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Output PPTX file path (default: input stem + '-narrated.pptx')",
    )
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable verbose (DEBUG) logging",
    )
    return parser


def _run(args: argparse.Namespace) -> int:
    """Execute audio embedding logic."""

    input_path: Path = args.input
    audio_dir: Path = args.audio_dir

    if not input_path.is_file():
        logger.error("Input PPTX not found: %s", input_path)
        return EXIT_FAILURE

    if not audio_dir.is_dir():
        logger.error("Audio directory not found: %s", audio_dir)
        return EXIT_FAILURE

    output_path: Path = args.output or input_path.with_name(
        f"{input_path.stem}-narrated.pptx"
    )

    if output_path.resolve() == input_path.resolve():
        logger.error(
            "Output path must differ from input path to avoid overwriting the source"
        )
        return EXIT_ERROR

    prs = Presentation(str(input_path))
    embedded_count = 0
    failed_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        wav_name = f"slide-{idx:03d}.wav"
        wav_path = audio_dir / wav_name
        if not wav_path.is_file():
            logger.info("SKIP slide %d: %s not found", idx, wav_name)
            continue

        if embed_slide_audio(slide, wav_path):
            embedded_count += 1
            logger.info("Embedded %s into slide %d", wav_name, idx)
        else:
            logger.error("FAILED to embed %s into slide %d", wav_name, idx)
            failed_count += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if embedded_count == 0:
        logger.error(
            "No audio files were embedded. Verify that slide-NNN.wav files exist in %s",
            audio_dir,
        )
        return EXIT_FAILURE

    try:
        prs.save(str(output_path))
    except OSError as exc:
        logger.error("Failed to save output PPTX %s: %s", output_path, exc)
        return EXIT_FAILURE

    logger.info("Saved %s with %d embedded audio files", output_path, embedded_count)

    if failed_count > 0:
        logger.warning(
            "Completed with %d failure(s); %d slide(s) embedded successfully.",
            failed_count,
            embedded_count,
        )
        return EXIT_FAILURE
    return EXIT_SUCCESS


def configure_logging(verbose: bool = False) -> None:
    """Configure logging based on verbosity level."""
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=level, format="%(levelname)s: %(message)s")


def main() -> int:
    """Entry point for audio embedding."""
    parser = create_parser()
    args = parser.parse_args()
    configure_logging(verbose=args.verbose)
    try:
        return _run(args)
    except KeyboardInterrupt:
        return 130
    except BrokenPipeError:
        sys.stderr.close()
        return 1


if __name__ == "__main__":
    sys.exit(main())
