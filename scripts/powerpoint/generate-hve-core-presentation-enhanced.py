#!/usr/bin/env python3
"""Generate enhanced HVE-Core and RPI Framework PowerPoint — 20 slides with D-RPI, results, and community."""

# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx",
#   "Pillow",
# ]
# ///

import os
import io
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

# ── Constants ──────────────────────────────────────────────

DARK_BG = RGBColor(0x1B, 0x1B, 0x1B)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_PURPLE = RGBColor(0x88, 0x64, 0xD8)
ACCENT_TEAL = RGBColor(0x03, 0x8E, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MED_GRAY = RGBColor(0x90, 0x90, 0x90)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
CARD_BG = RGBColor(0x2D, 0x2D, 0x2D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

EXIT_SUCCESS = 0
EXIT_FAILURE = 1

OUTPUT_DIR = os.path.join(os.path.dirname(
    os.path.abspath(__file__)), "..", "docs")
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(
    __file__)), "..", ".copilot-tracking", "presentation-assets")
os.makedirs(IMG_DIR, exist_ok=True)


# ── Helper Functions ───────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, bold, color = line, default_size, False, default_color
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


def add_shape_with_text(slide, left, top, width, height, text, fill_color,
                        font_size=11, font_color=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = "Segoe UI"
    p.font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def add_connector_line(slide, start_x, start_y, end_x, end_y, color=ACCENT_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def set_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_accent_bar(slide, top=Inches(1.35), color=ACCENT_BLUE):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.8), Pt(4))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


# ── Font Loader ────────────────────────────────────────────

def _get_fonts():
    try:
        return {
            "title": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22),
            "zone": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16),
            "comp": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 13),
            "small": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 11),
            "label": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 12),
            "big": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28),
            "num": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36),
        }
    except OSError:
        f = ImageFont.load_default()
        return {k: f for k in ("title", "zone", "comp", "small", "label", "big", "num")}


def _draw_rounded_box(draw, x, y, w, h, color, fill=(45, 45, 45)):
    draw.rounded_rectangle([(x, y), (x + w, y + h)],
                           radius=10, outline=color, width=2, fill=fill)


def _draw_dashed_rounded_box(draw, x, y, w, h, color, dash_len=10, gap_len=6):
    """Draw a rounded rectangle with a dashed border."""
    for dx in range(0, w, dash_len + gap_len):
        x1 = x + min(dx, w)
        x2 = x + min(dx + dash_len, w)
        draw.line([(x1, y), (x2, y)], fill=color, width=2)
    for dx in range(0, w, dash_len + gap_len):
        x1 = x + min(dx, w)
        x2 = x + min(dx + dash_len, w)
        draw.line([(x1, y + h), (x2, y + h)], fill=color, width=2)
    for dy in range(0, h, dash_len + gap_len):
        y1 = y + min(dy, h)
        y2 = y + min(dy + dash_len, h)
        draw.line([(x, y1), (x, y2)], fill=color, width=2)
    for dy in range(0, h, dash_len + gap_len):
        y1 = y + min(dy, h)
        y2 = y + min(dy + dash_len, h)
        draw.line([(x + w, y1), (x + w, y2)], fill=color, width=2)


def _draw_arrow(draw, x1, y1, x2, y2, color, label="", font=None):
    draw.line([(x1, y1), (x2, y2)], fill=color, width=2)
    dx, dy = x2 - x1, y2 - y1
    length = max(1, (dx * dx + dy * dy) ** 0.5)
    ux, uy = dx / length, dy / length
    px, py = -uy, ux
    draw.polygon([(x2, y2), (int(x2 - 8 * ux + 4 * px), int(y2 - 8 * uy + 4 * py)),
                  (int(x2 - 8 * ux - 4 * px), int(y2 - 8 * uy - 4 * py))], fill=color)
    if label and font:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        draw.text((mx, my - 15), label, fill=color, font=font)


def _draw_dashed_arrow(draw, x1, y1, x2, y2, color, dash_len=8, gap_len=6):
    """Draw a dashed line with an arrowhead at the end."""
    dx, dy = x2 - x1, y2 - y1
    length = max(1, (dx * dx + dy * dy) ** 0.5)
    ux, uy = dx / length, dy / length
    pos = 0
    while pos < length:
        sx = x1 + ux * pos
        sy = y1 + uy * pos
        end = min(pos + dash_len, length)
        ex = x1 + ux * end
        ey = y1 + uy * end
        draw.line([(int(sx), int(sy)), (int(ex), int(ey))], fill=color, width=2)
        pos += dash_len + gap_len
    px, py = -uy, ux
    draw.polygon([(x2, y2),
                  (int(x2 - 8 * ux + 4 * px), int(y2 - 8 * uy + 4 * py)),
                  (int(x2 - 8 * ux - 4 * px), int(y2 - 8 * uy - 4 * py))], fill=color)


# ── Diagram Generators ─────────────────────────────────────

def create_rpi_pipeline_diagram():
    """RPI pipeline flow: 4 phases connected by arrows with type labels and /clear markers."""
    w, h = 1600, 600
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "The RPI Pipeline",
              fill=(255, 255, 255), font=fonts["title"])

    phases = [
        ("Research", (0, 120, 212)),
        ("Plan", (16, 124, 16)),
        ("Implement", (255, 140, 0)),
        ("Review", (136, 100, 216)),
    ]

    type_labels = ["Uncertainty", "Knowledge",
                   "Strategy", "Working Code", "Validated Code"]

    box_w, box_h = 260, 100
    gap = 100
    start_x = 60
    y_top = 250

    for i, (name, color) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        _draw_rounded_box(draw, x, y_top, box_w, box_h, color)
        draw.text((x + box_w // 2 - len(name) * 5, y_top + 35),
                  name, fill=color, font=fonts["zone"])

        label_in = type_labels[i]
        draw.text((x + 10, y_top - 40), label_in,
                  fill=(200, 200, 200), font=fonts["comp"])

        if i < len(phases) - 1:
            ax1 = x + box_w + 4
            ax2 = x + box_w + gap - 4
            ay = y_top + box_h // 2
            _draw_arrow(draw, ax1, ay, ax2, ay, color)

            cx = x + box_w + gap // 2
            cy = ay
            r = 16
            draw.ellipse([(cx - r, cy - r), (cx + r, cy + r)],
                         fill=(209, 52, 56))
            draw.text((cx - 12, cy - 7), "/clear",
                      fill=(255, 255, 255), font=fonts["small"])

    last_x = start_x + 3 * (box_w + gap)
    draw.text((last_x + 100, y_top - 40),
              type_labels[4], fill=(200, 200, 200), font=fonts["comp"])

    # 88% stat callout
    draw.rounded_rectangle([(w // 2 - 150, y_top + box_h + 40),
                            (w // 2 + 150, y_top + box_h + 80)],
                           radius=8, fill=(16, 124, 16))
    draw.text((w // 2 - 130, y_top + box_h + 48),
              "Up to 88% faster task completion",
              fill=(255, 255, 255), font=fonts["comp"])

    draw.text((60, h - 60), "Red circles = /clear (context reset between phases)",
              fill=(209, 52, 56), font=fonts["comp"])
    draw.text((60, h - 35), "Each phase transforms one type of understanding into the next",
              fill=(200, 200, 200), font=fonts["small"])

    path = os.path.join(IMG_DIR, "hve-rpi-pipeline-enhanced.png")
    img.save(path, "PNG")
    return path


def create_drpi_pipeline_diagram():
    """D-RPI extended pipeline: optional Discovery + 4 RPI phases with iteration loops."""
    w, h = 1600, 650
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 250, 12), "D-RPI: Discovery + RPI Pipeline",
              fill=(255, 255, 255), font=fonts["title"])

    # Discovery box — dashed border
    disc_x, disc_y, disc_w, disc_h = 40, 260, 220, 100
    disc_color = (144, 144, 144)
    _draw_dashed_rounded_box(draw, disc_x, disc_y, disc_w, disc_h, disc_color)
    draw.text((disc_x + 40, disc_y + 35), "Discovery",
              fill=disc_color, font=fonts["zone"])
    draw.text((disc_x + 60, disc_y + disc_h + 8), "(optional)",
              fill=(144, 144, 144), font=fonts["small"])

    # 4 RPI phases
    phases = [
        ("Research", (0, 120, 212)),
        ("Plan", (16, 124, 16)),
        ("Implement", (255, 140, 0)),
        ("Review", (136, 100, 216)),
    ]
    type_labels = ["Empathize/\nDefine", "Uncertainty", "Knowledge",
                   "Strategy", "Working Code", "Validated Code"]

    box_w, box_h = 240, 100
    gap = 60
    rpi_start_x = 320
    y_top = 260

    for i, (name, color) in enumerate(phases):
        x = rpi_start_x + i * (box_w + gap)
        _draw_rounded_box(draw, x, y_top, box_w, box_h, color)
        draw.text((x + box_w // 2 - len(name) * 5, y_top + 35),
                  name, fill=color, font=fonts["zone"])

        label = type_labels[i + 1]
        draw.text((x + 10, y_top - 35), label,
                  fill=(200, 200, 200), font=fonts["comp"])

        if i < len(phases) - 1:
            ax1 = x + box_w + 4
            ax2 = x + box_w + gap - 4
            ay = y_top + box_h // 2
            _draw_arrow(draw, ax1, ay, ax2, ay, color)
            cx = x + box_w + gap // 2
            cy = ay
            r = 14
            draw.ellipse([(cx - r, cy - r), (cx + r, cy + r)],
                         fill=(209, 52, 56))
            draw.text((cx - 10, cy - 6), "/clear",
                      fill=(255, 255, 255), font=fonts["small"])

    # Type label for Discovery
    draw.text((disc_x + 10, disc_y - 35), type_labels[0],
              fill=(144, 144, 144), font=fonts["comp"])

    # Dashed arrow: Discovery → Research
    _draw_dashed_arrow(draw, disc_x + disc_w + 4, disc_y + disc_h // 2,
                       rpi_start_x - 4, y_top + box_h // 2, disc_color)

    # Output type label after last box
    last_x = rpi_start_x + 3 * (box_w + gap)
    draw.text((last_x + 80, y_top - 35),
              type_labels[5], fill=(200, 200, 200), font=fonts["comp"])

    # Iteration loop arrow: Review bottom → Research bottom
    loop_y = y_top + box_h + 40
    draw.line([(last_x + box_w // 2, y_top + box_h),
               (last_x + box_w // 2, loop_y)],
              fill=(136, 100, 216), width=2)
    draw.line([(last_x + box_w // 2, loop_y),
               (rpi_start_x + box_w // 2, loop_y)],
              fill=(136, 100, 216), width=2)
    draw.line([(rpi_start_x + box_w // 2, loop_y),
               (rpi_start_x + box_w // 2, y_top + box_h)],
              fill=(136, 100, 216), width=2)
    # Arrowhead pointing up to Research
    rx = rpi_start_x + box_w // 2
    draw.polygon([(rx, y_top + box_h),
                  (rx - 5, y_top + box_h + 10),
                  (rx + 5, y_top + box_h + 10)],
                 fill=(136, 100, 216))
    draw.text((w // 2 - 60, loop_y - 18), "iteration loops",
              fill=(136, 100, 216), font=fonts["small"])

    draw.text((40, h - 55), "Dashed border = optional Discovery phase",
              fill=(144, 144, 144), font=fonts["comp"])
    draw.text((40, h - 30), "Red circles = /clear (context reset)  |  Purple loop = Review can route back",
              fill=(200, 200, 200), font=fonts["small"])

    path = os.path.join(IMG_DIR, "hve-drpi-pipeline.png")
    img.save(path, "PNG")
    return path


def create_quality_comparison_diagram():
    """Side-by-side comparison: Traditional AI vs RPI approach across 5 dimensions."""
    w, h = 1600, 700
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "Quality Comparison",
              fill=(255, 255, 255), font=fonts["title"])

    col_left = 100
    col_right = 850
    col_w = 650

    draw.rounded_rectangle(
        [(col_left, 55), (col_left + col_w, 95)], radius=8, fill=(100, 30, 30))
    draw.text((col_left + col_w // 2 - 60, 65), "Traditional AI",
              fill=(255, 255, 255), font=fonts["zone"])

    draw.rounded_rectangle(
        [(col_right, 55), (col_right + col_w, 95)], radius=8, fill=(20, 80, 20))
    draw.text((col_right + col_w // 2 - 30, 65), "RPI",
              fill=(255, 255, 255), font=fonts["zone"])

    rows = [
        ("Pattern Matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code chain"),
        ("Knowledge Transfer", "Lost between sessions",
         "Preserved in research artifacts"),
        ("Rework", "30-50% of AI output discarded", "< 10% rework rate"),
        ("Validation", "Manual review only", "Structured multi-dimension review"),
    ]

    for i, (dimension, trad, rpi) in enumerate(rows):
        ry = 120 + i * 110

        draw.rounded_rectangle([(col_left, ry), (col_left + col_w + col_right - col_left + col_w - col_left, ry + 30)],
                               radius=4, fill=(45, 45, 45))
        draw.text((w // 2 - len(dimension) * 4, ry + 6), dimension,
                  fill=(255, 255, 255), font=fonts["label"])

        draw.rounded_rectangle([(col_left, ry + 35), (col_left + col_w, ry + 90)], radius=8,
                               fill=(50, 25, 25), outline=(209, 52, 56), width=1)
        draw.text((col_left + 20, ry + 50), trad,
                  fill=(255, 150, 150), font=fonts["comp"])

        draw.rounded_rectangle([(col_right, ry + 35), (col_right + col_w, ry + 90)], radius=8,
                               fill=(20, 50, 20), outline=(16, 124, 16), width=1)
        draw.text((col_right + 20, ry + 50), rpi,
                  fill=(150, 255, 150), font=fonts["comp"])

    path = os.path.join(IMG_DIR, "hve-quality-comparison-enhanced.png")
    img.save(path, "PNG")
    return path


def create_role_mapping_diagram():
    """8 roles mapped to collections and workflow patterns in a 2x4 grid."""
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "Eight Roles, One Framework",
              fill=(255, 255, 255), font=fonts["title"])

    roles = [
        ("Developer", "hve-core", "RPI Pipeline", (0, 120, 212)),
        ("TPM / Lead", "project-planning", "Reqs → PRD → WIT", (16, 124, 16)),
        ("Platform Engineer", "coding-standards",
         "Artifact authoring", (255, 140, 0)),
        ("OSS Contributor", "github", "Backlog management", (136, 100, 216)),
        ("Security Engineer", "security-planning",
         "Threat modeling", (209, 52, 56)),
        ("Data Scientist", "data-science",
         "Spec → Notebook → Dashboard", (3, 142, 142)),
        ("Project Planner", "project-planning", "PRD / BRD / ADR", (0, 120, 212)),
        ("UX / DT Practitioner", "design-thinking",
         "9-method coaching", (136, 100, 216)),
    ]

    card_w, card_h = 340, 170
    cols, rows_count = 4, 2
    x_start = 40
    y_start = 60
    x_gap = (w - 2 * x_start - cols * card_w) // (cols - 1)
    y_gap = 40

    for idx, (role, collection, workflow, color) in enumerate(roles):
        col = idx % cols
        row = idx // cols
        x = x_start + col * (card_w + x_gap)
        y = y_start + row * (card_h + y_gap)

        _draw_rounded_box(draw, x, y, card_w, card_h, color)

        draw.text((x + 15, y + 15), role, fill=color, font=fonts["zone"])
        draw.text((x + 15, y + 50),
                  f"Collection: {collection}", fill=(200, 200, 200), font=fonts["comp"])
        draw.rounded_rectangle([(x + 15, y + 85), (x + card_w - 15, y + 130)],
                               radius=6, fill=(27, 27, 27), outline=color, width=1)
        draw.text((x + 25, y + 95), workflow, fill=color, font=fonts["comp"])
        draw.line([(x + 15, y + 75), (x + card_w - 15, y + 75)],
                  fill=color, width=1)

    draw.text((40, h - 40), "Each role maps to a specific collection and workflow pattern",
              fill=(200, 200, 200), font=fonts["small"])

    path = os.path.join(IMG_DIR, "hve-role-mapping-enhanced.png")
    img.save(path, "PNG")
    return path


# ── Slide Builders: Part 1 — "Why RPI?" (Slides 1–10) ─────

def slide_01_title(prs):
    """Slide 1: Title card with industry credibility strip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
                 "HVE-Core and the RPI Framework", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(3.0), color=ACCENT_BLUE)
    add_multi_text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(2), [
        ("Turning AI from a Code Generator into a Research Partner", 20, False, LIGHT_GRAY),
        ("", 10),
        ("A constraint-based framework for AI-assisted engineering", 16, False, MED_GRAY),
    ])
    add_shape_with_text(slide, Inches(0.8), Inches(5.2), Inches(3.0), Inches(0.5),
                        "30-Minute Overview  •  Two Parts", ACCENT_BLUE, font_size=12)

    # Industry credibility strip
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
                 "Proven at AVEVA  •  BMW  •  Michelin  •  Hexagon  •  Kubota  •  Nvidia",
                 font_size=14, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "Welcome. Over the next 30 minutes I'll show you why the way most teams use AI "
        "coding assistants is fundamentally broken — and how a constraint-based framework "
        "called RPI fixes it. We'll cover the problem, the framework, a live demo, and "
        "how different roles use it today.\n\n"
        "HVE-Core has been industry-proven at companies like AVEVA, BMW, Michelin, Hexagon, "
        "Kubota, and Nvidia to improve coding accuracy, automate repetitive tasks, and "
        "personalize solutions."
    ))


def slide_02_the_problem(prs):
    """Slide 2: The universally relatable AI failure mode with pain point statistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Problem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_RED)

    add_multi_text(slide, Inches(1.0), Inches(1.5), Inches(7.5), Inches(4.0), [
        ("", 8),
        ('You: "Build me a Terraform module for Azure IoT"', 18, False, WHITE),
        ("", 12),
        ("AI: *immediately generates 2,000 lines of code*", 18, False, LIGHT_GRAY),
        ("", 12),
        ("Reality:", 20, True, ACCENT_RED),
        ("  • Missing provider dependencies", 14, False, ACCENT_RED),
        ("  • Wrong variable naming conventions", 14, False, ACCENT_RED),
        ("  • Patterns from 3 years ago", 14, False, ACCENT_RED),
        ("  • Broke two downstream services", 14, False, ACCENT_RED),
        ("  • Used modules that don't exist in your registry", 14, False, ACCENT_RED),
        ("", 12),
        ("It looked right. It compiled. Then you deployed.", 14, True, MED_GRAY),
    ], default_size=14, default_color=WHITE)

    # Pain point statistics panel
    add_shape_with_text(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(0.5),
                        "The Cost", ACCENT_RED, font_size=14)
    add_multi_text(slide, Inches(9.0), Inches(2.2), Inches(3.5), Inches(3.0), [
        ("15–20 hrs/week", 20, True, ACCENT_ORANGE),
        ("lost to repetitive tasks and", 12, False, LIGHT_GRAY),
        ("relearning context", 12, False, LIGHT_GRAY),
        ("", 10),
        ("40% productivity loss", 20, True, ACCENT_ORANGE),
        ("from unstructured AI interactions", 12, False, LIGHT_GRAY),
        ("and context-switching", 12, False, LIGHT_GRAY),
    ])

    add_text_box(slide, Inches(1.0), Inches(6.5), Inches(10), Inches(0.5),
                 "Raise your hand if this has happened to you.", font_size=14, color=MED_GRAY)
    set_speaker_notes(slide, (
        "Raise your hand if this has happened to you. The AI generates something that looks "
        "right. It compiles. Then you deploy and discover it used patterns from 3 years ago, "
        "missed your naming conventions, and broke two downstream services. "
        "This is the universal failure mode of AI-assisted development.\n\n"
        "The pain is quantifiable: developers lose 15 to 20 hours per week dealing with "
        "repetitive tasks and relearning context. Unstructured AI interactions and constant "
        "context-switching cost about 40% of productivity. Those aren't hypothetical numbers — "
        "they come from real engineering teams."
    ))


def slide_03_why_it_happens(prs):
    """Slide 3: Root cause — AI conflates investigation with implementation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Why It Happens", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_shape_with_text(slide, Inches(1.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Investigate?\n\n🔍\n\nSearch • Analyze • Understand",
                        CARD_BG, font_size=16, font_color=ACCENT_BLUE)

    add_shape_with_text(slide, Inches(7.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Implement!\n\n⌨️\n\nGenerate • Write • Ship",
                        CARD_BG, font_size=16, font_color=ACCENT_ORANGE)

    add_connector_line(slide, Inches(5.7), Inches(3.0), Inches(7.3), Inches(3.0),
                       color=ACCENT_RED, width=Pt(3))

    add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9), Inches(0.8),
                 "AI Writes First and Thinks Never", font_size=28, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.0), Inches(5.7), Inches(9), Inches(0.5),
                 '"Plausible" and "correct" aren\'t the same thing',
                 font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "The root cause: AI treats every request as an implementation request. "
        "It never separates investigation from implementation. It pattern-matches "
        "from training data instead of searching your actual codebase. "
        "The output looks plausible — but plausible and correct aren't the same thing."
    ))


def slide_04_counterintuitive_insight(prs):
    """Slide 4: The key insight — constraints, not intelligence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Counterintuitive Insight", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_BLUE)

    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(2.5), [
        ('"The solution isn\'t teaching AI to be smarter.', 22, False, WHITE),
        ("", 8),
        ('It\'s preventing AI from doing certain things', 22, False, WHITE),
        ('at certain times."', 22, False, WHITE),
    ])
    add_accent_bar(slide, top=Inches(3.8), color=ACCENT_BLUE)

    add_shape_with_text(slide, Inches(1.5), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for plausible code", ACCENT_RED, font_size=16)
    add_connector_line(slide, Inches(6.2), Inches(5.15), Inches(7.0), Inches(5.15),
                       color=WHITE, width=Pt(3))
    add_shape_with_text(slide, Inches(7.2), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for verified truth", ACCENT_GREEN, font_size=16)

    add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
                 "Constraints are features, not limitations",
                 font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "This is the philosophical foundation of HVE-Core. We don't try to make AI smarter. "
        "We constrain what it can do at each stage. A researcher that cannot write code will "
        "search instead of inventing. A planner that cannot implement will sequence instead "
        "of coding. Constraints force better behavior."
    ))


def slide_05_what_is_hve_core(prs):
    """Slide 5: Elevator pitch, component summary, and industry adoption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "What is HVE-Core?", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 "An enterprise-ready prompt engineering framework that transforms "
                 "GitHub Copilot from a code-completion tool into a structured engineering partner.",
                 font_size=15, color=LIGHT_GRAY)

    components = [
        ("Agents", "22", ACCENT_BLUE),
        ("Prompts", "27", ACCENT_GREEN),
        ("Instructions", "24", ACCENT_ORANGE),
        ("Skills", "1", ACCENT_TEAL),
        ("Collections", "10", ACCENT_PURPLE),
    ]
    card_w = Inches(2.1)
    for i, (name, count, color) in enumerate(components):
        left = Inches(0.8 + i * 2.4)
        add_shape_with_text(slide, left, Inches(2.4), card_w, Inches(1.2),
                            f"{name}\n{count}", color, font_size=18)

    # Delegation flow
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
                 "Delegation Flow", font_size=18, bold=True, color=WHITE)

    flow_items = ["User", "Prompt", "Agent", "Instructions"]
    flow_colors = [MED_GRAY, ACCENT_GREEN, ACCENT_BLUE, ACCENT_ORANGE]
    for i, (item, color) in enumerate(zip(flow_items, flow_colors)):
        left = Inches(1.0 + i * 2.8)
        add_shape_with_text(slide, left, Inches(4.6), Inches(2.0), Inches(0.6),
                            item, CARD_BG, font_size=14, font_color=color)
        if i < len(flow_items) - 1:
            add_connector_line(slide, left + Inches(2.1), Inches(4.9),
                               left + Inches(2.7), Inches(4.9), color=color, width=Pt(2))

    add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
        ("Agents define behavior and constraints. Instructions encode standards.",
         13, False, MED_GRAY),
        ("Prompts are user entry points. Skills package domain knowledge.",
         13, False, MED_GRAY),
        ("Collections bundle everything for specific roles.", 13, False, MED_GRAY),
        ("", 6),
        ("Industry-proven at AVEVA, BMW, Michelin, Hexagon, Kubota, and Nvidia",
         13, True, ACCENT_TEAL),
    ])

    set_speaker_notes(slide, (
        "HVE-Core is an enterprise-ready prompt engineering framework. 22 custom agents, "
        "27 prompts, 24 instruction files, 1 skill, and 10 collections. Agents define what "
        "AI can and cannot do. Instructions encode coding standards. Prompts are how users "
        "interact. Collections bundle artifacts for specific roles — a developer gets different "
        "tools than a TPM.\n\n"
        "This isn't just theory. HVE-Core has been used by industry leaders like AVEVA, "
        "BMW, Michelin, Hexagon, Kubota, and Nvidia to boost productivity and improve "
        "coding accuracy in real production environments."
    ))


def slide_06_rpi_pipeline(prs, rpi_img):
    """Slide 6: The RPI type transformation pipeline diagram with 88% stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "The RPI Pipeline", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(rpi_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "Each phase converts one form of understanding into the next. Uncertainty becomes "
        "knowledge. Knowledge becomes strategy. Strategy becomes working code. Working code "
        "becomes validated code. The /clear markers between phases prevent context "
        "contamination — the researcher's reasoning must not leak into the implementor.\n\n"
        "Internal trials show up to 88% faster task completion using RPI compared to "
        "unstructured AI prompting. That means what used to take a week might take one day."
    ))


def slide_07_phase_deep_dive(prs):
    """Slide 7: Attribute cards for all 4 RPI phases with best practice tip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Deep Dive: Each Phase", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    phases = [
        ("Research", ACCENT_BLUE,
         "Discover what exists",
         "CANNOT implement or modify code",
         "research.md with citations",
         "/task-research"),
        ("Plan", ACCENT_GREEN,
         "Sequence the work",
         "CANNOT implement or modify code",
         "plan.instructions.md + details.md",
         "/task-plan"),
        ("Implement", ACCENT_ORANGE,
         "Execute the plan",
         "MUST follow plan — no creative decisions",
         "Working code matching plan",
         "/task-implement"),
        ("Review", ACCENT_PURPLE,
         "Validate the output",
         "CANNOT modify code — only report findings",
         "8-dimension review report",
         "/task-review"),
    ]

    card_w = Inches(2.9)
    for i, (name, color, purpose, constraint, output, invocation) in enumerate(phases):
        left = Inches(0.5 + i * 3.1)
        top_start = Inches(1.5)

        add_shape_with_text(slide, left, top_start, card_w, Inches(0.5),
                            name, color, font_size=16)

        add_multi_text(slide, left + Inches(0.1), top_start + Inches(0.6), card_w - Inches(0.2), Inches(4.5), [
            ("Purpose", 11, True, color),
            (purpose, 10, False, LIGHT_GRAY),
            ("", 6),
            ("Core Constraint", 11, True, ACCENT_RED),
            (constraint, 10, True, WHITE),
            ("", 6),
            ("Key Output", 11, True, color),
            (output, 10, False, LIGHT_GRAY),
            ("", 6),
            ("Invocation", 11, True, color),
            (invocation, 10, False, MED_GRAY),
        ])

    # Best practice tip
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 'Best practice: Keep implementation to ~3 files per RPI cycle for manageable scope',
                 font_size=12, bold=True, color=ACCENT_TEAL)

    set_speaker_notes(slide, (
        "The constraints row is the key architectural element. The researcher cannot implement "
        "— so it searches instead of inventing. The planner cannot implement — so it sequences "
        "instead of coding. The implementor follows the plan — no creative decisions. The "
        "reviewer cannot modify code — only validate. Each constraint forces the AI into its "
        "correct mode of operation.\n\n"
        "Best practice: keep each RPI cycle scoped to about 3 files. If the plan is larger, "
        "split into multiple cycles. This keeps tasks manageable and prevents context overload."
    ))


def slide_08_quality_and_clear(prs, quality_img):
    """Slide 8: Quality comparison + /clear rule + paradigm shift quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Quality Comparison + The /clear Rule", font_size=28, bold=True)
    add_accent_bar(slide, top=Inches(0.85))

    dimensions = [
        ("Pattern matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code"),
        ("Knowledge transfer", "Lost between sessions", "Preserved in artifacts"),
        ("Rework", "30-50% discarded", "< 10% rework"),
        ("Validation", "Manual review only", "8-dimension structured review"),
    ]

    add_shape_with_text(slide, Inches(1.5), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Traditional AI", ACCENT_RED, font_size=11)
    add_shape_with_text(slide, Inches(4.7), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Dimension", DARK_GRAY, font_size=11)
    add_shape_with_text(slide, Inches(7.9), Inches(1.2), Inches(3.5), Inches(0.35),
                        "RPI Approach", ACCENT_GREEN, font_size=11)

    for i, (dim, trad, rpi) in enumerate(dimensions):
        row_top = Inches(1.6 + i * 0.42)
        add_text_box(slide, Inches(1.5), row_top, Inches(3.0), Inches(0.38),
                     trad, font_size=10, color=ACCENT_RED)
        add_text_box(slide, Inches(4.7), row_top, Inches(3.0), Inches(0.38),
                     dim, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.9), row_top, Inches(3.5), Inches(0.38),
                     rpi, font_size=10, color=ACCENT_GREEN)

    # /clear flow
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
                 "The /clear Rule: Context Reset Between Phases", font_size=16, bold=True)

    clear_flow = [
        ("Research", ACCENT_BLUE),
        ("/clear", ACCENT_RED),
        ("Plan", ACCENT_GREEN),
        ("/clear", ACCENT_RED),
        ("Implement", ACCENT_ORANGE),
        ("/clear", ACCENT_RED),
        ("Review", ACCENT_PURPLE),
    ]

    for i, (label, color) in enumerate(clear_flow):
        left = Inches(0.8 + i * 1.7)
        if label == "/clear":
            add_shape_with_text(slide, left, Inches(4.5), Inches(1.0), Inches(0.45),
                                label, ACCENT_RED, font_size=11,
                                shape=MSO_SHAPE.OVAL)
        else:
            add_shape_with_text(slide, left, Inches(4.5), Inches(1.4), Inches(0.45),
                                label, color, font_size=12)
        if i < len(clear_flow) - 1:
            add_connector_line(slide, left + Inches(1.4 if label != "/clear" else 1.0),
                               Inches(4.72),
                               left + Inches(1.6 if label != "/clear" else 1.2),
                               Inches(4.72), color=MED_GRAY, width=Pt(2))

    # Paradigm shift quote
    add_multi_text(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(1.2), [
        ("", 6),
        ('Stop asking AI: "Write this code."', 15, True, WHITE),
        ('Start asking: "Help me research, plan, then implement with evidence."',
         15, True, ACCENT_BLUE),
    ])

    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
                 "Context clearing prevents mode contamination — each phase gets a clean start.",
                 font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "For TPMs: the traceability row is your strongest argument for RPI adoption. "
        "Every line of code traces back through a plan to a research finding. "
        "For everyone: the /clear rule sounds counterintuitive — why throw away context? "
        "Because context contamination causes the researcher's reasoning patterns to leak "
        "into the implementor, creating mode confusion. Each phase needs a clean start "
        "with only its designated inputs.\n\n"
        "The paradigm shift is real: stop asking AI to write code. Start asking it to help "
        "you research, plan, and then implement with evidence. That single change transforms "
        "AI from a code generator into a research partner."
    ))


def slide_09_chat_modes(prs):
    """Slide 9: Overview of all 6 custom Copilot chat modes with rpi-agent note."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Six Custom Chat Modes", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    # Group 1: RPI Cycle
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                 "RPI Cycle", font_size=16, bold=True, color=ACCENT_BLUE)

    rpi_modes = [
        ("Task Researcher", "/task-research", "Cannot implement", ACCENT_BLUE),
        ("Task Planner", "/task-plan", "Cannot implement", ACCENT_GREEN),
        ("Task Implementor", "/task-implement",
         "Follows plan only", ACCENT_ORANGE),
        ("Task Reviewer", "/task-review", "Cannot modify code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(rpi_modes):
        top = Inches(2.0 + i * 0.95)
        add_shape_with_text(slide, Inches(0.8), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(3.5), top, Inches(1.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(3.5), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)
        if i < len(rpi_modes) - 1:
            add_connector_line(slide, Inches(2.0), top + Inches(0.75),
                               Inches(2.0), top + Inches(0.9),
                               color=MED_GRAY, width=Pt(1))

    # Group 2: Complementary
    add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
                 "Complementary", font_size=16, bold=True, color=ACCENT_TEAL)

    comp_modes = [
        ("Prompt Builder", "/prompt-build", "Orchestrates subagents", ACCENT_TEAL),
        ("PR Review", "Agent picker", "Never modifies code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(comp_modes):
        top = Inches(2.0 + i * 1.2)
        add_shape_with_text(slide, Inches(7.5), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(10.2), top, Inches(2.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(10.2), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)

    add_multi_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.0), [
        ("Every mode enforces specific constraints — AI cannot bypass its designated role",
         13, False, MED_GRAY),
        ("", 6),
        ("Advanced: rpi-agent orchestrates all 4 phases in a single session for familiar tasks",
         11, False, ACCENT_TEAL),
    ])

    set_speaker_notes(slide, (
        "Four modes form the RPI cycle. Two additional modes serve specialized functions. "
        "Prompt Builder creates the artifacts that configure the agents — it's meta-level. "
        "PR Review is an 8-dimension quality gate that complements Task Reviewer. "
        "Every mode has a constraint that prevents the AI from doing something specific.\n\n"
        "Advanced tip: for familiar tasks where you know the codebase well, the rpi-agent "
        "can orchestrate all four phases in a single session. But for anything complex or "
        "unfamiliar, running phases separately with /clear gives much better results."
    ))


def slide_10_demo_intro(prs):
    """Slide 10: Setup slide before switching to VS Code for live demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "Let's See It in Action", font_size=40, bold=True,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(2.3), color=ACCENT_GREEN)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
                 "Live Demo: RPI Workflow — Researching and Planning a Feature",
                 font_size=20, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.4),
                 "Watch for:", font_size=16, bold=True, color=WHITE)

    watch_items = [
        ("1.", "How the researcher finds existing patterns (not invents code)", ACCENT_BLUE),
        ("2.", "How it cites specific files and line numbers", ACCENT_GREEN),
        ("3.", "How the plan references the research artifacts", ACCENT_ORANGE),
        ("4.", "How /clear resets context between phases", ACCENT_RED),
        ("5.", "How the implementor follows the plan step by step", ACCENT_PURPLE),
    ]

    for i, (num, desc, color) in enumerate(watch_items):
        top = Inches(4.1 + i * 0.5)
        add_shape_with_text(slide, Inches(1.0), top, Inches(0.5), Inches(0.4),
                            num, color, font_size=14)
        add_text_box(slide, Inches(1.7), top, Inches(10), Inches(0.4),
                     desc, font_size=14, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
                 "Demo: ~5 minutes in VS Code  •  Research → /clear → Plan → /clear → Implement",
                 font_size=13, color=MED_GRAY)

    set_speaker_notes(slide, (
        "I'm going to run through the RPI workflow live. Watch for five things: "
        "how the researcher searches rather than invents, how it cites files and line numbers, "
        "how the plan references research artifacts, how /clear resets context, and how "
        "the implementor follows the plan step by step.\n\n"
        "Demo flow: open VS Code, run /task-research, show output, /clear, "
        "run /task-plan with research file, show plan and details output, /clear, "
        "run /task-implement with plan.\n\n"
        "Fallback: if live demo fails, switch to pre-captured screenshots or terminal output."
    ))


# ── Slide Builders: Part 2 — "HVE-Core in Practice" (Slides 11–20) ──

def slide_11_who_uses_hve_core(prs, role_img):
    """Slide 11: 8-role workflow mapping."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Eight Roles, One Framework", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(role_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "HVE-Core isn't just for developers. Eight distinct roles, each with a dedicated "
        "collection that filters the right tools for their job.\n\n"
        "For developers: the RPI pipeline. For TPMs and leads: requirements discovery → PRD "
        "authoring → ADO work item creation. For PMs and non-coding roles: use Researcher to "
        "gather data for status reports, Planner to structure narratives, Prompt Builder to "
        "formulate user stories. For OSS contributors: RPI as a self-service mentor — pick a "
        "'good first issue,' research the codebase, plan your approach, implement with guidance.\n\n"
        "Notice that roles don't share workflows — a TPM's tooling is different from a "
        "developer's. Collections make that separation automatic."
    ))


def slide_12_discovery_design_thinking(prs, drpi_img):
    """Slide 12: D-RPI extended pipeline and Design Thinking alignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Beyond RPI: Discovery + Design Thinking", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_TEAL)

    # D-RPI diagram
    slide.shapes.add_picture(drpi_img, Inches(
        0.3), Inches(1.3), Inches(12.7), Inches(3.5))

    # Design Thinking alignment table
    add_text_box(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.4),
                 "Design Thinking Alignment", font_size=16, bold=True, color=WHITE)

    dt_rows = [
        ("Discovery", "Empathize / Define", MED_GRAY),
        ("Research", "Ideate (evidence-constrained)", ACCENT_BLUE),
        ("Plan", "Prototype (on paper)", ACCENT_GREEN),
    ]
    dt_rows_right = [
        ("Implement", "Build / Test", ACCENT_ORANGE),
        ("Review", "Validate / Feedback", ACCENT_PURPLE),
    ]

    for i, (phase, dt_stage, color) in enumerate(dt_rows):
        row_top = Inches(5.4 + i * 0.35)
        add_shape_with_text(slide, Inches(1.0), row_top, Inches(1.8), Inches(0.3),
                            phase, color, font_size=10)
        add_text_box(slide, Inches(3.0), row_top, Inches(3.0), Inches(0.3),
                     f"→  {dt_stage}", font_size=10, color=LIGHT_GRAY)

    for i, (phase, dt_stage, color) in enumerate(dt_rows_right):
        row_top = Inches(5.4 + i * 0.35)
        add_shape_with_text(slide, Inches(7.0), row_top, Inches(1.8), Inches(0.3),
                            phase, color, font_size=10)
        add_text_box(slide, Inches(9.0), row_top, Inches(3.0), Inches(0.3),
                     f"→  {dt_stage}", font_size=10, color=LIGHT_GRAY)

    add_text_box(slide, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
                 "D-RPI and Design Thinking integration are forward-looking extensions. "
                 "Core four-phase RPI is the production workflow today.",
                 font_size=11, color=MED_GRAY)

    set_speaker_notes(slide, (
        "RPI adapts to different situations. For complex or ambiguous projects where "
        "requirements are unclear, you can add a Discovery phase before Research — we call "
        "this D-RPI. Discovery is a brainstorming session with the AI in Ask Mode: 'Who is "
        "the audience? What do they care about? What do they already know?' It mirrors the "
        "Empathize and Define stages of Design Thinking.\n\n"
        "This alignment isn't a coincidence. RPI maps naturally to Design Thinking — Research "
        "parallels Ideation, Plan parallels Prototyping, Implement parallels Building, Review "
        "parallels Testing with feedback loops. Context from Discovery pre-seeds the Research "
        "phase, so the AI already understands the high-level intent.\n\n"
        "Note that D-RPI and Design Thinking integration are forward-looking extensions. The "
        "core four-phase RPI is the production workflow today. And RPI is cyclical, not linear "
        "— the Review phase can route back to Research for knowledge gaps, back to Plan for "
        "scope changes, or back to Implement for fixes."
    ))


def slide_13_real_results(prs):
    """Slide 13: Concrete case study metrics from field deployments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Real Results with HVE-Core", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    # 5 metric cards
    metrics = [
        ("50%", "Faster\nDeployments", "Global telecom:\nAI-driven IaC", ACCENT_GREEN),
        ("90%", "Faster\nDocs", "Architecture docs\nand security plans:\ndays → hours", ACCENT_GREEN),
        ("2 Days", "vs 8 Weeks", "Internal hackathon:\nRPI + Copilot\nprototype", ACCENT_GREEN),
        ("Quality", "Improved", "AI PR reviews\ncatch subtle bugs\npre-merge", ACCENT_BLUE),
        ("Dev\nExperience", "Better", "Less rework,\nless context-switching,\nhigher satisfaction", ACCENT_TEAL),
    ]

    card_w = Inches(2.3)
    for i, (number, subtitle, detail, color) in enumerate(metrics):
        left = Inches(0.5 + i * 2.5)

        # Top strip with bold number
        add_shape_with_text(slide, left, Inches(1.6), card_w, Inches(1.0),
                            number, color, font_size=36)
        # Subtitle
        add_shape_with_text(slide, left, Inches(2.6), card_w, Inches(0.8),
                            subtitle, CARD_BG, font_size=14, font_color=LIGHT_GRAY)
        # Detail text
        add_multi_text(slide, left + Inches(0.1), Inches(3.6), card_w - Inches(0.2), Inches(1.5), [
            (detail, 10, False, MED_GRAY),
        ])

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                 "Real teams.  Real metrics.  Real improvements.",
                 font_size=16, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 88% stat banner
    add_shape_with_text(slide, Inches(3.0), Inches(6.2), Inches(7.0), Inches(0.5),
                        "Internal trials: up to 88% faster task completion with RPI",
                        ACCENT_GREEN, font_size=13)

    set_speaker_notes(slide, (
        "These aren't hypothetical — real teams have seen dramatic improvements with HVE-Core. "
        "A global telecom rolled out an HVE approach in their cloud deployment process and "
        "cut deployment times by 50%. Another team saw architecture documents and security "
        "plans go from taking several days to just a few hours — a 90% reduction.\n\n"
        "At an internal hackathon, a team delivered a working prototype in only 2 days using "
        "RPI and Copilot — something that traditionally would have taken 6 to 8 weeks. "
        "Quality improved too: the AI-driven PR Review agent helped developers catch subtle "
        "bugs and security issues before code was merged.\n\n"
        "And developers are happier. By eliminating tedious tasks and reducing context-switching, "
        "HVE-Core frees engineers to focus on creative problem-solving. Teams report significantly "
        "less frustration and higher morale."
    ))


def slide_14_dogfooding_and_validation(prs):
    """Slide 14: Self-referential credibility + CI/CD pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Built by HVE-Core, Validated by HVE-Core", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_TEAL)

    # Left column: Dogfooding
    add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        ("Self-Referential Credibility", 18, True, ACCENT_TEAL),
        ("", 8),
        ("• This presentation was planned using RPI", 14, False, LIGHT_GRAY),
        ("• Every agent was built by the Prompt Builder agent", 14, False, LIGHT_GRAY),
        ("• Instructions files follow their own rules", 14, False, LIGHT_GRAY),
        ("• Collection manifests are validated by their own CI", 14, False, LIGHT_GRAY),
        ("• Documentation is generated from the same artifacts", 14, False, LIGHT_GRAY),
        ("", 8),
        ("If HVE-Core doesn't work on itself,", 14, True, MED_GRAY),
        ("why would it work on your project?", 14, True, MED_GRAY),
    ])

    # Right column: Validation Pipeline
    add_multi_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5), [
        ("Validation Pipeline", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("Linting (7 jobs)", 14, True, ACCENT_BLUE),
        ("  Markdown, YAML, frontmatter, links,", 12, False, MED_GRAY),
        ("  tables, collections metadata, version consistency", 12, False, MED_GRAY),
        ("", 6),
        ("Analysis (2 jobs)", 14, True, ACCENT_GREEN),
        ("  PowerShell PSScriptAnalyzer, skill validation", 12, False, MED_GRAY),
        ("", 6),
        ("Security (3 jobs)", 14, True, ACCENT_RED),
        ("  Dependency pinning, SHA staleness, copyright", 12, False, MED_GRAY),
        ("", 6),
        ("Schema Validation", 14, True, ACCENT_PURPLE),
        ("  Collection YAML → plugin generation → extension packaging", 12, False, MED_GRAY),
    ])

    set_speaker_notes(slide, (
        "HVE-Core uses its own tools to build itself. This presentation was planned using "
        "the RPI workflow. Every agent was built by the Prompt Builder agent. The validation "
        "pipeline runs 12+ automated checks on every PR. If it doesn't work on itself, "
        "why would it work on your project? Self-referential credibility is the strongest "
        "argument for adoption."
    ))


def slide_15_extension_ecosystem(prs):
    """Slide 15: 8 VS Code extension packages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Extension Ecosystem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_PURPLE)

    headers = [("Extension", Inches(0.8), Inches(3.0)),
               ("Agents", Inches(4.0), Inches(1.2)),
               ("Prompts", Inches(5.3), Inches(1.2)),
               ("Instructions", Inches(6.6), Inches(1.5)),
               ("Focus", Inches(8.3), Inches(4.0))]

    for name, left, width in headers:
        add_shape_with_text(slide, left, Inches(1.5), width, Inches(0.4),
                            name, ACCENT_BLUE, font_size=11)

    extensions = [
        ("hve-core", "21", "23", "18", "Core RPI + all foundations", True),
        ("hve-ado", "9", "19", "10", "Azure DevOps integration", False),
        ("hve-github", "9", "19", "10", "GitHub backlog management", False),
        ("hve-project-planning", "13", "15", "5",
         "PRD / BRD / ADR workflows", False),
        ("hve-security-planning", "9", "16", "5", "Threat modeling", False),
        ("hve-rpi", "8", "14", "5", "Standalone RPI cycle", False),
        ("hve-prompt-engineering", "8", "14", "5", "Artifact authoring", False),
        ("hve-data-science", "1", "1", "0", "Data science workflows", False),
    ]

    for i, (ext, agents, prompts, instr, focus, highlight) in enumerate(extensions):
        row_top = Inches(2.0 + i * 0.5)
        bg_color = ACCENT_BLUE if highlight else CARD_BG

        if highlight:
            add_shape_with_text(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                                ext, bg_color, font_size=11)
        else:
            add_text_box(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                         ext, font_size=11, color=LIGHT_GRAY)

        add_text_box(slide, Inches(4.0), row_top, Inches(1.2), Inches(0.4),
                     agents, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(5.3), row_top, Inches(1.2), Inches(0.4),
                     prompts, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.6), row_top, Inches(1.5), Inches(0.4),
                     instr, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.3), row_top, Inches(4.0), Inches(0.4),
                     focus, font_size=11, color=MED_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Common base: 8 core agents + 5 core instructions shared across all extensions",
                 font_size=13, bold=True, color=ACCENT_TEAL)

    set_speaker_notes(slide, (
        "8 VS Code extensions organized by role. hve-core is the full package with 21 agents, "
        "23 prompts, and 18 instruction files. Specialized extensions like hve-ado and "
        "hve-github focus on specific workflows. All share a common base of 8 core agents "
        "and 5 core instructions. Pick the extension that matches your role."
    ))


def slide_16_learning_curve(prs):
    """Slide 16: Honest framing + payoff narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Learning Curve & Compounding Value", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_ORANGE)

    # Left: honest framing
    add_shape_with_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Honest Truth", ACCENT_ORANGE, font_size=14)
    add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("Your first RPI workflow will feel slower.", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("You'll wonder why you can't just ask AI to code it.", 14, False, LIGHT_GRAY),
        ("You'll feel like the constraints are slowing you down.", 14, False, LIGHT_GRAY),
        ("You'll want to skip the research phase.", 14, False, LIGHT_GRAY),
        ("", 8),
        ("That's normal. Stay with it.", 14, True, MED_GRAY),
    ])

    # Right: payoff
    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Payoff", ACCENT_GREEN, font_size=14)
    add_multi_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("By your third feature, the workflow", 18, True, ACCENT_GREEN),
        ("feels natural.", 18, True, ACCENT_GREEN),
        ("", 8),
        ("Research artifacts compound across features.", 14, False, LIGHT_GRAY),
        ("Plans reference previous research findings.", 14, False, LIGHT_GRAY),
        ("Rework drops below 10%.", 14, False, LIGHT_GRAY),
        ("", 8),
        ("The framework pays for itself.", 14, True, MED_GRAY),
    ])

    # Institutional memory callout
    add_multi_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1.0), [
        ("Research documents accumulate into institutional memory.", 15, True, WHITE),
        ("New team members can read how past decisions were made. You're not just solving ",
         12, False, LIGHT_GRAY),
        ("today's problem — you're building the knowledge base that accelerates tomorrow's.",
         12, False, LIGHT_GRAY),
    ])

    # Version timeline
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
                 "Evolution Timeline", font_size=14, bold=True, color=WHITE)

    versions = [
        ("v1.1.0", "Jan 19", "Initial release", ACCENT_BLUE),
        ("v2.0.0", "Jan 28", "Breaking: agent restructure", ACCENT_RED),
        ("v2.3.4", "Feb 13", "Stable + extensions", ACCENT_GREEN),
    ]

    for i, (ver, date, desc, color) in enumerate(versions):
        left = Inches(1.0 + i * 3.8)
        add_shape_with_text(slide, left, Inches(6.6), Inches(1.2), Inches(0.4),
                            ver, color, font_size=11)
        add_text_box(slide, left + Inches(1.4), Inches(6.6), Inches(2.2), Inches(0.4),
                     f"{date} — {desc}", font_size=10, color=MED_GRAY)
        if i < len(versions) - 1:
            add_connector_line(slide, left + Inches(3.5), Inches(6.8),
                               left + Inches(3.7), Inches(6.8),
                               color=MED_GRAY, width=Pt(2))

    set_speaker_notes(slide, (
        "Be honest with your team: the first RPI workflow feels slower. That's by design. "
        "You're building research artifacts that compound across features. By the third "
        "feature, the workflow feels natural and you'll wonder how you worked without it.\n\n"
        "The real value isn't the current task — it's what compounds. Research documents "
        "accumulate into institutional memory. New team members can read how past decisions "
        "were made. You're not just solving today's problem; you're building the knowledge "
        "base that accelerates tomorrow's."
    ))


def slide_17_learning_resources(prs):
    """Slide 17: Three installation paths + learning resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Three Paths to Start + Learn More", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    # 3 installation cards
    install_methods = [
        ("VS Code Extension ⭐", "10 seconds", "Individual users, TPMs,\nimmediate access", ACCENT_BLUE),
        ("Peer Clone", "2 minutes", "Developers needing\ncustomization", ACCENT_ORANGE),
        ("Codespaces", "1 click", "Contributors wanting\nzero-config", ACCENT_GREEN),
    ]

    for i, (method, time, desc, color) in enumerate(install_methods):
        left = Inches(0.8 + i * 4.0)
        add_shape_with_text(slide, left, Inches(1.5), Inches(3.5), Inches(0.5),
                            f"{method}  •  {time}", color, font_size=13)
        add_multi_text(slide, left + Inches(0.2), Inches(2.2), Inches(3.2), Inches(1.0), [
            (desc, 12, False, LIGHT_GRAY),
        ])

    # Resources header
    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(11), Inches(0.4),
                 "Learning Resources", font_size=18, bold=True, color=WHITE)

    # Left resource items
    resources_left = [
        ("HVE-Learning Repository", "Self-paced modules on prompt engineering, RPI, backlog management", ACCENT_BLUE),
        ("Customer Zero Katas", "Hands-on practice exercises for real-world scenarios", ACCENT_BLUE),
    ]
    for i, (name, desc, color) in enumerate(resources_left):
        row_top = Inches(3.9 + i * 0.6)
        add_shape_with_text(slide, Inches(0.8), row_top, Inches(0.3), Inches(0.3),
                            "•", color, font_size=14, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(1.3), row_top, Inches(5.0), Inches(0.5), [
            (name, 13, True, WHITE),
            (desc, 11, False, MED_GRAY),
        ])

    # Right resource items
    resources_right = [
        ("First Workflow Tutorial", "15-minute guided RPI exercise", ACCENT_GREEN),
        ("Documentation", "aka.ms/hve-core — install, usage, FAQs", ACCENT_GREEN),
    ]
    for i, (name, desc, color) in enumerate(resources_right):
        row_top = Inches(3.9 + i * 0.6)
        add_shape_with_text(slide, Inches(7.0), row_top, Inches(0.3), Inches(0.3),
                            "•", color, font_size=14, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(7.5), row_top, Inches(5.0), Inches(0.5), [
            (name, 13, True, WHITE),
            (desc, 11, False, MED_GRAY),
        ])

    # Links bar
    add_shape_with_text(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5),
                        "microsoft/hve-learning  •  aka.ms/cz-repo-katas  •  aka.ms/hve-core",
                        ACCENT_BLUE, font_size=13)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Start with one RPI workflow. The research artifacts compound from there.",
                 font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "Three paths to get started, based on your role. The VS Code Extension is the fastest "
        "— literally 10 seconds from the Marketplace. Search 'HVE Core,' install, and you "
        "have all six chat modes immediately. If you need to customize artifacts, clone the "
        "repo. For zero-config, use GitHub Codespaces — one click.\n\n"
        "For learning, the first-workflow tutorial takes 15 minutes — it's a guided RPI "
        "exercise. The HVE-Learning repository has full self-paced modules on prompt "
        "engineering, RPI, and backlog management. Customer Zero Katas provide hands-on "
        "practice exercises for real-world scenarios.\n\n"
        "All documentation lives at aka.ms/hve-core — installation steps, usage examples, "
        "FAQs, and deep-dive guides."
    ))


def slide_18_community(prs):
    """Slide 18: Community engagement and contribution pathways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Join the Community", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_PURPLE)

    # 4 engagement cards in 2x2 grid
    cards = [
        ("Bi-Weekly Syncs", "HVE Community Syncs — live demos,\nQ&A, roadmap updates,\nand best practice sharing",
         ACCENT_BLUE, Inches(0.8), Inches(1.5)),
        ("Teams Channels", "#Hypervelocity and #SeasonOfHVE\n— real-time discussions,\ntips, troubleshooting",
         ACCENT_TEAL, Inches(7.0), Inches(1.5)),
        ("Good First Issues", 'Look for "good first issue" tags\non HVE-Core and accelerator\nrepos (Edge AI, Robotics)',
         ACCENT_GREEN, Inches(0.8), Inches(3.3)),
        ("Contribution Types", "Code • Prompts • Instructions\n• Skills • Docs • Bug reports\n• Feature requests",
         ACCENT_ORANGE, Inches(7.0), Inches(3.3)),
    ]

    for title, desc, color, left, top in cards:
        add_shape_with_text(slide, left, top, Inches(5.5), Inches(0.5),
                            title, color, font_size=14)
        add_multi_text(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.2), Inches(1.0), [
            (desc, 12, False, LIGHT_GRAY),
        ])

    # Contribution flow
    flow_items = [
        ("Find Issue", ACCENT_BLUE),
        ("Use RPI to Research", ACCENT_GREEN),
        ("Submit PR", ACCENT_ORANGE),
        ("Get Review", ACCENT_PURPLE),
    ]
    for i, (label, color) in enumerate(flow_items):
        left = Inches(1.0 + i * 2.8)
        add_shape_with_text(slide, left, Inches(5.3), Inches(2.2), Inches(0.5),
                            label, color, font_size=11)
        if i < len(flow_items) - 1:
            add_connector_line(slide, left + Inches(2.3), Inches(5.55),
                               left + Inches(2.7), Inches(5.55),
                               color=MED_GRAY, width=Pt(2))

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Use RPI to learn the repo and make your first PR — the agents keep you on track",
                 font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "We'd love to have you be part of the HVE community. Bi-weekly community syncs are "
        "the best way to see what's new — live demos of new features, Q&A with the core team, "
        "and roadmap discussions. These are open to everyone.\n\n"
        "For day-to-day collaboration, join the #Hypervelocity and #SeasonOfHVE Teams channels. "
        "You'll find tips, troubleshooting help, and discussions about best practices.\n\n"
        "Contributing is straightforward. Issues labeled 'good first issue' are curated for "
        "newcomers. And here's the meta play: use RPI itself to make your first contribution. "
        "Run Task Researcher to understand the codebase, Task Planner to outline your approach, "
        "then implement. Contributions aren't limited to code — prompts, instructions, "
        "documentation, and bug reports are all valuable."
    ))


def slide_19_getting_started(prs):
    """Slide 19: Streamlined getting started with GitHub URL."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Getting Started", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    # Quick start steps
    steps = [
        ("1", "Install HVE-Core extension from VS Code Marketplace (10 seconds)", ACCENT_BLUE),
        ("2", "Open Copilot Chat and select a chat mode (e.g., Task Researcher)", ACCENT_GREEN),
        ("3", "Run /task-research on your next feature or bug", ACCENT_ORANGE),
        ("4", "Experience the difference — then try the full RPI cycle", ACCENT_PURPLE),
    ]

    for i, (num, desc, color) in enumerate(steps):
        top = Inches(1.6 + i * 0.8)
        add_shape_with_text(slide, Inches(1.0), top, Inches(0.6), Inches(0.6),
                            num, color, font_size=20, shape=MSO_SHAPE.OVAL)
        add_text_box(slide, Inches(2.0), top + Inches(0.1), Inches(10), Inches(0.5),
                     desc, font_size=16, color=LIGHT_GRAY)

    # GitHub URL
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
                 "github.com/microsoft/hve-core", font_size=22, bold=True,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # CTA
    add_shape_with_text(slide, Inches(3.0), Inches(5.8), Inches(7), Inches(0.6),
                        "Try /task-research on your next feature", ACCENT_GREEN, font_size=16)

    set_speaker_notes(slide, (
        "Four steps to get started. Install the extension — 10 seconds. Open Copilot Chat. "
        "Select Task Researcher. Ask it about your next feature. That's it.\n\n"
        "Don't change your whole workflow — just use the researcher once and see what happens. "
        "When you see how it searches the codebase instead of inventing code, you'll understand "
        "why constraints are features."
    ))


def slide_20_key_takeaways(prs):
    """Slide 20: Four key takeaways, closing quote, CTA, and Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Key Takeaways", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_BLUE)

    # 4 numbered takeaways
    takeaways = [
        ("1", "Accelerate delivery and quality", "HVE-Core combines AI tools with structured "
         "practices — tasks up to 88% faster without sacrificing quality or security", ACCENT_BLUE),
        ("2", "RPI is the game-changer", "Research → Plan → Implement → Review turns Copilot "
         "from a nifty helper into a reliable partner. Structured phases kill the AI rework loop", ACCENT_GREEN),
        ("3", "Empower every role", "Not just for coders — PMs, TPMs, security engineers, "
         "data scientists, and OSS contributors each have dedicated workflows", ACCENT_ORANGE),
        ("4", "Start your hypervelocity journey", "It only takes 10 seconds to install. "
         "Pick a small task, run it through RPI, and see the difference", ACCENT_PURPLE),
    ]

    for i, (num, title, desc, color) in enumerate(takeaways):
        top = Inches(1.5 + i * 0.8)
        add_shape_with_text(slide, Inches(0.8), top, Inches(0.5), Inches(0.5),
                            num, color, font_size=16, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(1.5), top, Inches(10.5), Inches(0.7), [
            (title, 15, True, WHITE),
            (desc, 12, False, LIGHT_GRAY),
        ])

    # Closing quote with accent border
    quote_shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.8), Pt(4), Inches(0.7))
    quote_shp.fill.solid()
    quote_shp.fill.fore_color.rgb = ACCENT_BLUE
    quote_shp.line.fill.background()

    add_text_box(slide, Inches(1.8), Inches(4.8), Inches(10), Inches(0.7),
                 '"The code comes last, after the hard work of understanding is complete."',
                 font_size=18, bold=True, color=WHITE)

    # CTA flow: Install → Try RPI → Join Community
    cta_items = [
        ("Install", ACCENT_GREEN),
        ("Try RPI", ACCENT_BLUE),
        ("Join Community", ACCENT_PURPLE),
    ]
    for i, (label, color) in enumerate(cta_items):
        left = Inches(2.5 + i * 3.0)
        add_shape_with_text(slide, left, Inches(5.8), Inches(2.5), Inches(0.5),
                            label, color, font_size=14)
        if i < len(cta_items) - 1:
            add_connector_line(slide, left + Inches(2.6), Inches(6.05),
                               left + Inches(2.9), Inches(6.05),
                               color=MED_GRAY, width=Pt(2))

    # Q&A
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.8),
                 "Questions?", font_size=36, bold=True,
                 alignment=PP_ALIGN.CENTER, color=WHITE)

    set_speaker_notes(slide, (
        "To wrap up: Hypervelocity Engineering is about making your whole team dramatically "
        "faster and better. First, HVE-Core delivers real results — some teams have seen nearly "
        "a tenfold speed-up on certain tasks, while also catching quality issues that could "
        "have slipped through.\n\n"
        "Second, RPI is the heart of this approach. It might feel counterintuitive at first to "
        "slow down and do separate research and planning steps, but that structure is exactly "
        "what unlocks the speed later. Structured phases kill the AI rework loop.\n\n"
        "Third, HVE-Core empowers everyone, not just developers. PMs create PRDs. TPMs build "
        "backlog hierarchies. Security engineers generate threat models. Our call to action: "
        "try HVE-Core on one task this week. The setup is trivial — 10 seconds. 'The code "
        "comes last, after the hard work of understanding is complete.' Thank you — let's go "
        "forth and build at hypervelocity. Questions?"
    ))


# ── Main Assembly ──────────────────────────────────────────

def main() -> int:
    """Generate the enhanced HVE-Core and RPI presentation deck (20 slides)."""
    print("Generating enhanced HVE-Core and RPI presentation...")

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(IMG_DIR, exist_ok=True)

    print("  Creating diagrams...")
    rpi_img = create_rpi_pipeline_diagram()
    drpi_img = create_drpi_pipeline_diagram()
    quality_img = create_quality_comparison_diagram()
    role_img = create_role_mapping_diagram()

    print("  Building slides...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Part 1: Why RPI? (Slides 1-10)
    slide_01_title(prs)
    slide_02_the_problem(prs)
    slide_03_why_it_happens(prs)
    slide_04_counterintuitive_insight(prs)
    slide_05_what_is_hve_core(prs)
    slide_06_rpi_pipeline(prs, rpi_img)
    slide_07_phase_deep_dive(prs)
    slide_08_quality_and_clear(prs, quality_img)
    slide_09_chat_modes(prs)
    slide_10_demo_intro(prs)

    # Part 2: HVE-Core in Practice (Slides 11-20)
    slide_11_who_uses_hve_core(prs, role_img)
    slide_12_discovery_design_thinking(prs, drpi_img)
    slide_13_real_results(prs)
    slide_14_dogfooding_and_validation(prs)
    slide_15_extension_ecosystem(prs)
    slide_16_learning_curve(prs)
    slide_17_learning_resources(prs)
    slide_18_community(prs)
    slide_19_getting_started(prs)
    slide_20_key_takeaways(prs)

    output_path = os.path.join(OUTPUT_DIR, "hve-core-rpi-presentation-enhanced.pptx")
    prs.save(output_path)
    print(f"\n✅ Saved presentation: {output_path}")
    print(f"   {len(prs.slides)} slides generated")
    return EXIT_SUCCESS


if __name__ == "__main__":
    sys.exit(main())
