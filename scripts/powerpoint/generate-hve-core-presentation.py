#!/usr/bin/env python3
"""Generate HVE-Core and RPI Framework PowerPoint presentation with diagrams and speaker notes."""

# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx",
#   "Pillow",
# ]
# ///

import os
import io
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

# ── Constants ──────────────────────────────────────────────

DARK_BG = RGBColor(0x1B, 0x1B, 0x1B)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_PURPLE = RGBColor(0x88, 0x64, 0xD8)
ACCENT_TEAL = RGBColor(0x03, 0x8E, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MED_GRAY = RGBColor(0x90, 0x90, 0x90)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
CARD_BG = RGBColor(0x2D, 0x2D, 0x2D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

EXIT_SUCCESS = 0
EXIT_FAILURE = 1

OUTPUT_DIR = os.path.join(os.path.dirname(
    os.path.abspath(__file__)), "..", "docs")
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(
    __file__)), "..", ".copilot-tracking", "presentation-assets")
os.makedirs(IMG_DIR, exist_ok=True)


# ── Helper Functions ───────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, bold, color = line, default_size, False, default_color
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


def add_shape_with_text(slide, left, top, width, height, text, fill_color,
                        font_size=11, font_color=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = "Segoe UI"
    p.font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def add_connector_line(slide, start_x, start_y, end_x, end_y, color=ACCENT_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def set_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_accent_bar(slide, top=Inches(1.35), color=ACCENT_BLUE):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.8), Pt(4))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


# ── Font Loader ────────────────────────────────────────────

def _get_fonts():
    try:
        return {
            "title": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22),
            "zone": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16),
            "comp": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 13),
            "small": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 11),
            "label": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 12),
            "big": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28),
            "num": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36),
        }
    except OSError:
        f = ImageFont.load_default()
        return {k: f for k in ("title", "zone", "comp", "small", "label", "big", "num")}


def _draw_rounded_box(draw, x, y, w, h, color, fill=(45, 45, 45)):
    draw.rounded_rectangle([(x, y), (x + w, y + h)],
                           radius=10, outline=color, width=2, fill=fill)


def _draw_arrow(draw, x1, y1, x2, y2, color, label="", font=None):
    draw.line([(x1, y1), (x2, y2)], fill=color, width=2)
    dx, dy = x2 - x1, y2 - y1
    length = max(1, (dx * dx + dy * dy) ** 0.5)
    ux, uy = dx / length, dy / length
    px, py = -uy, ux
    draw.polygon([(x2, y2), (int(x2 - 8 * ux + 4 * px), int(y2 - 8 * uy + 4 * py)),
                  (int(x2 - 8 * ux - 4 * px), int(y2 - 8 * uy - 4 * py))], fill=color)
    if label and font:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        draw.text((mx, my - 15), label, fill=color, font=font)


# ── Diagram Generators ─────────────────────────────────────

def create_rpi_pipeline_diagram():
    """RPI pipeline flow: 4 phases connected by arrows with type labels and /clear markers."""
    w, h = 1600, 600
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "The RPI Pipeline",
              fill=(255, 255, 255), font=fonts["title"])

    phases = [
        ("Research", (0, 120, 212)),
        ("Plan", (16, 124, 16)),
        ("Implement", (255, 140, 0)),
        ("Review", (136, 100, 216)),
    ]

    type_labels = ["Uncertainty", "Knowledge",
                   "Strategy", "Working Code", "Validated Code"]

    box_w, box_h = 260, 100
    gap = 100
    start_x = 60
    y_top = 250

    for i, (name, color) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        _draw_rounded_box(draw, x, y_top, box_w, box_h, color)
        draw.text((x + box_w // 2 - len(name) * 5, y_top + 35),
                  name, fill=color, font=fonts["zone"])

        # Type label above the box (input type)
        label_in = type_labels[i]
        draw.text((x + 10, y_top - 40), label_in,
                  fill=(200, 200, 200), font=fonts["comp"])

        # Arrow and /clear marker between phases
        if i < len(phases) - 1:
            ax1 = x + box_w + 4
            ax2 = x + box_w + gap - 4
            ay = y_top + box_h // 2
            _draw_arrow(draw, ax1, ay, ax2, ay, color)

            # /clear marker as red circle
            cx = x + box_w + gap // 2
            cy = ay
            r = 16
            draw.ellipse([(cx - r, cy - r), (cx + r, cy + r)],
                         fill=(209, 52, 56))
            draw.text((cx - 12, cy - 7), "/clear",
                      fill=(255, 255, 255), font=fonts["small"])

    # Output type label after last box
    last_x = start_x + 3 * (box_w + gap)
    draw.text((last_x + 100, y_top - 40),
              type_labels[4], fill=(200, 200, 200), font=fonts["comp"])

    # Legend
    draw.text((60, h - 60), "Red circles = /clear (context reset between phases)",
              fill=(209, 52, 56), font=fonts["comp"])
    draw.text((60, h - 35), "Each phase transforms one type of understanding into the next",
              fill=(200, 200, 200), font=fonts["small"])

    path = os.path.join(IMG_DIR, "hve-rpi-pipeline.png")
    img.save(path, "PNG")
    return path


def create_quality_comparison_diagram():
    """Side-by-side comparison: Traditional AI vs RPI approach across 5 dimensions."""
    w, h = 1600, 700
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "Quality Comparison",
              fill=(255, 255, 255), font=fonts["title"])

    # Column headers
    col_left = 100
    col_right = 850
    col_w = 650

    draw.rounded_rectangle(
        [(col_left, 55), (col_left + col_w, 95)], radius=8, fill=(100, 30, 30))
    draw.text((col_left + col_w // 2 - 60, 65), "Traditional AI",
              fill=(255, 255, 255), font=fonts["zone"])

    draw.rounded_rectangle(
        [(col_right, 55), (col_right + col_w, 95)], radius=8, fill=(20, 80, 20))
    draw.text((col_right + col_w // 2 - 30, 65), "RPI",
              fill=(255, 255, 255), font=fonts["zone"])

    rows = [
        ("Pattern Matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code chain"),
        ("Knowledge Transfer", "Lost between sessions",
         "Preserved in research artifacts"),
        ("Rework", "30-50% of AI output discarded", "< 10% rework rate"),
        ("Validation", "Manual review only", "Structured multi-dimension review"),
    ]

    for i, (dimension, trad, rpi) in enumerate(rows):
        ry = 120 + i * 110

        # Dimension label
        draw.rounded_rectangle([(col_left, ry), (col_left + col_w + col_right - col_left + col_w - col_left, ry + 30)],
                               radius=4, fill=(45, 45, 45))
        draw.text((w // 2 - len(dimension) * 4, ry + 6), dimension,
                  fill=(255, 255, 255), font=fonts["label"])

        # Traditional value
        draw.rounded_rectangle([(col_left, ry + 35), (col_left + col_w, ry + 90)], radius=8,
                               fill=(50, 25, 25), outline=(209, 52, 56), width=1)
        draw.text((col_left + 20, ry + 50), trad,
                  fill=(255, 150, 150), font=fonts["comp"])

        # RPI value
        draw.rounded_rectangle([(col_right, ry + 35), (col_right + col_w, ry + 90)], radius=8,
                               fill=(20, 50, 20), outline=(16, 124, 16), width=1)
        draw.text((col_right + 20, ry + 50), rpi,
                  fill=(150, 255, 150), font=fonts["comp"])

    path = os.path.join(IMG_DIR, "hve-quality-comparison.png")
    img.save(path, "PNG")
    return path


def create_role_mapping_diagram():
    """8 roles mapped to collections and workflow patterns in a 2x4 grid."""
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (27, 27, 27))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 150, 12), "Who Uses HVE-Core?",
              fill=(255, 255, 255), font=fonts["title"])

    roles = [
        ("Developer", "hve-core", "RPI Pipeline", (0, 120, 212)),
        ("TPM / Lead", "project-planning", "Reqs → PRD → WIT", (16, 124, 16)),
        ("Platform Engineer", "coding-standards",
         "Artifact authoring", (255, 140, 0)),
        ("OSS Contributor", "github", "Backlog management", (136, 100, 216)),
        ("Security Engineer", "security-planning",
         "Threat modeling", (209, 52, 56)),
        ("Data Scientist", "data-science",
         "Spec → Notebook → Dashboard", (3, 142, 142)),
        ("Project Planner", "project-planning", "PRD / BRD / ADR", (0, 120, 212)),
        ("UX / DT Practitioner", "design-thinking",
         "9-method coaching", (136, 100, 216)),
    ]

    card_w, card_h = 340, 170
    cols, rows_count = 4, 2
    x_start = 40
    y_start = 60
    x_gap = (w - 2 * x_start - cols * card_w) // (cols - 1)
    y_gap = 40

    for idx, (role, collection, workflow, color) in enumerate(roles):
        col = idx % cols
        row = idx // cols
        x = x_start + col * (card_w + x_gap)
        y = y_start + row * (card_h + y_gap)

        _draw_rounded_box(draw, x, y, card_w, card_h, color)

        # Role name
        draw.text((x + 15, y + 15), role, fill=color, font=fonts["zone"])
        # Collection
        draw.text((x + 15, y + 50),
                  f"Collection: {collection}", fill=(200, 200, 200), font=fonts["comp"])
        # Workflow
        draw.rounded_rectangle([(x + 15, y + 85), (x + card_w - 15, y + 130)],
                               radius=6, fill=(27, 27, 27), outline=color, width=1)
        draw.text((x + 25, y + 95), workflow, fill=color, font=fonts["comp"])

        # Separator line
        draw.line([(x + 15, y + 75), (x + card_w - 15, y + 75)],
                  fill=color, width=1)

    # Legend
    draw.text((40, h - 40), "Each role maps to a specific collection and workflow pattern — no overlap, no confusion",
              fill=(200, 200, 200), font=fonts["small"])

    path = os.path.join(IMG_DIR, "hve-role-mapping.png")
    img.save(path, "PNG")
    return path


# ── Slide Builders: Part 1 — "Why RPI?" (Slides 1–8) ──────

def slide_01_title(prs):
    """Slide 1: Title card."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
                 "HVE-Core and the RPI Framework", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(3.2), color=ACCENT_BLUE)
    add_multi_text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(2), [
        ("Turning AI from a Code Generator into a Research Partner", 20, False, LIGHT_GRAY),
        ("", 10),
        ("A constraint-based framework for AI-assisted engineering", 16, False, MED_GRAY),
    ])
    add_shape_with_text(slide, Inches(0.8), Inches(6.0), Inches(3.0), Inches(0.5),
                        "30-Minute Overview  •  Two Parts", ACCENT_BLUE, font_size=12)
    set_speaker_notes(slide, (
        "Welcome. Over the next 30 minutes I'll show you why the way most teams use AI "
        "coding assistants is fundamentally broken — and how a constraint-based framework "
        "called RPI fixes it. We'll cover the problem, the framework, a live demo, and "
        "how different roles use it today."
    ))


def slide_02_the_problem(prs):
    """Slide 2: The universally relatable AI failure mode."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Problem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_RED)

    # Terminal-style monospace prompt
    add_multi_text(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5), [
        ("", 8),
        ('You: "Build me a Terraform module for Azure IoT"', 18, False, WHITE, ),
        ("", 12),
        ("AI: *immediately generates 2,000 lines of code*", 18, False, LIGHT_GRAY),
        ("", 12),
        ("Reality:", 20, True, ACCENT_RED),
        ("  • Missing provider dependencies", 16, False, ACCENT_RED),
        ("  • Wrong variable naming conventions", 16, False, ACCENT_RED),
        ("  • Patterns from 3 years ago", 16, False, ACCENT_RED),
        ("  • Broke two downstream services", 16, False, ACCENT_RED),
        ("  • Used modules that don't exist in your registry", 16, False, ACCENT_RED),
        ("", 12),
        ("It looked right. It compiled. Then you deployed.", 16, True, MED_GRAY),
    ], default_size=14, default_color=WHITE)

    add_text_box(slide, Inches(1.0), Inches(6.5), Inches(10), Inches(0.5),
                 "Raise your hand if this has happened to you.", font_size=14, color=MED_GRAY,
                 font_name="Segoe UI")
    set_speaker_notes(slide, (
        "Raise your hand if this has happened to you. The AI generates something that looks "
        "right. It compiles. Then you deploy and discover it used patterns from 3 years ago, "
        "missed your naming conventions, and broke two downstream services. "
        "This is the universal failure mode of AI-assisted development."
    ))


def slide_03_why_it_happens(prs):
    """Slide 3: Root cause — AI conflates investigation with implementation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Why It Happens", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    # Left column: Investigate?
    add_shape_with_text(slide, Inches(1.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Investigate?\n\n🔍\n\nSearch • Analyze • Understand",
                        CARD_BG, font_size=16, font_color=ACCENT_BLUE)

    # Right column: Implement!
    add_shape_with_text(slide, Inches(7.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Implement!\n\n⌨️\n\nGenerate • Write • Ship",
                        CARD_BG, font_size=16, font_color=ACCENT_ORANGE)

    # Center: crossed arrow
    add_connector_line(slide, Inches(5.7), Inches(3.0), Inches(7.3), Inches(3.0),
                       color=ACCENT_RED, width=Pt(3))

    # Central message
    add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9), Inches(0.8),
                 "AI Writes First and Thinks Never", font_size=28, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(2.0), Inches(5.7), Inches(9), Inches(0.5),
                 '"Plausible" and "correct" aren\'t the same thing',
                 font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "The root cause: AI treats every request as an implementation request. "
        "It never separates investigation from implementation. It pattern-matches "
        "from training data instead of searching your actual codebase. "
        "The output looks plausible — but plausible and correct aren't the same thing."
    ))


def slide_04_counterintuitive_insight(prs):
    """Slide 4: The key insight — constraints, not intelligence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Counterintuitive Insight", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_BLUE)

    # Large centered quote
    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(2.5), [
        ('"The solution isn\'t teaching AI to be smarter.', 22, False, WHITE),
        ("", 8),
        ('It\'s preventing AI from doing certain things', 22, False, WHITE),
        ('at certain times."', 22, False, WHITE),
    ])
    add_accent_bar(slide, top=Inches(3.8), color=ACCENT_BLUE)

    # Transformation arrow
    add_shape_with_text(slide, Inches(1.5), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for plausible code", ACCENT_RED, font_size=16)
    add_connector_line(slide, Inches(6.2), Inches(5.15), Inches(7.0), Inches(5.15),
                       color=WHITE, width=Pt(3))
    add_shape_with_text(slide, Inches(7.2), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for verified truth", ACCENT_GREEN, font_size=16)

    add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
                 "Constraints are features, not limitations",
                 font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "This is the philosophical foundation of HVE-Core. We don't try to make AI smarter. "
        "We constrain what it can do at each stage. A researcher that cannot write code will "
        "search instead of inventing. A planner that cannot implement will sequence instead "
        "of coding. Constraints force better behavior."
    ))


def slide_05_what_is_hve_core(prs):
    """Slide 5: Elevator pitch and component summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "What is HVE-Core?", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 "An enterprise-ready prompt engineering framework that transforms "
                 "GitHub Copilot from a code-completion tool into a structured engineering partner.",
                 font_size=15, color=LIGHT_GRAY)

    # 5 component cards
    components = [
        ("Agents", "22", ACCENT_BLUE),
        ("Prompts", "27", ACCENT_GREEN),
        ("Instructions", "24", ACCENT_ORANGE),
        ("Skills", "1", ACCENT_TEAL),
        ("Collections", "10", ACCENT_PURPLE),
    ]
    card_w = Inches(2.1)
    for i, (name, count, color) in enumerate(components):
        left = Inches(0.8 + i * 2.4)
        add_shape_with_text(slide, left, Inches(2.4), card_w, Inches(1.2),
                            f"{name}\n{count}", color, font_size=18)

    # Delegation flow
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
                 "Delegation Flow", font_size=18, bold=True, color=WHITE)

    flow_items = ["User", "Prompt", "Agent", "Instructions"]
    flow_colors = [MED_GRAY, ACCENT_GREEN, ACCENT_BLUE, ACCENT_ORANGE]
    for i, (item, color) in enumerate(zip(flow_items, flow_colors)):
        left = Inches(1.0 + i * 2.8)
        add_shape_with_text(slide, left, Inches(4.6), Inches(2.0), Inches(0.6),
                            item, CARD_BG, font_size=14, font_color=color)
        if i < len(flow_items) - 1:
            add_connector_line(slide, left + Inches(2.1), Inches(4.9),
                               left + Inches(2.7), Inches(4.9), color=color, width=Pt(2))

    add_multi_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.5), [
        ("Agents define behavior and constraints. Instructions encode standards.",
         13, False, MED_GRAY),
        ("Prompts are user entry points. Skills package domain knowledge.",
         13, False, MED_GRAY),
        ("Collections bundle everything for specific roles.", 13, False, MED_GRAY),
    ])

    set_speaker_notes(slide, (
        "HVE-Core is an enterprise-ready prompt engineering framework. 22 custom agents, "
        "27 prompts, 24 instruction files, 1 skill, and 10 collections. Agents define what "
        "AI can and cannot do. Instructions encode coding standards. Prompts are how users "
        "interact. Collections bundle artifacts for specific roles — a developer gets different "
        "tools than a TPM."
    ))


def slide_06_rpi_pipeline(prs, rpi_img):
    """Slide 6: The RPI type transformation pipeline diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "The RPI Pipeline", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(rpi_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "Each phase converts one form of understanding into the next. Uncertainty becomes "
        "knowledge. Knowledge becomes strategy. Strategy becomes working code. Working code "
        "becomes validated code. The /clear markers between phases prevent context "
        "contamination — the researcher's reasoning must not leak into the implementor."
    ))


def slide_07_phase_deep_dive(prs):
    """Slide 7: Attribute cards for all 4 RPI phases."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Deep Dive: Each Phase", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    phases = [
        ("Research", ACCENT_BLUE,
         "Discover what exists",
         "CANNOT implement or modify code",
         "research.md with citations",
         "/task-research"),
        ("Plan", ACCENT_GREEN,
         "Sequence the work",
         "CANNOT implement or modify code",
         "plan.instructions.md + details.md",
         "/task-plan"),
        ("Implement", ACCENT_ORANGE,
         "Execute the plan",
         "MUST follow plan — no creative decisions",
         "Working code matching plan",
         "/task-implement"),
        ("Review", ACCENT_PURPLE,
         "Validate the output",
         "CANNOT modify code — only report findings",
         "8-dimension review report",
         "/task-review"),
    ]

    card_w = Inches(2.9)
    for i, (name, color, purpose, constraint, output, invocation) in enumerate(phases):
        left = Inches(0.5 + i * 3.1)
        top_start = Inches(1.5)

        # Phase header
        add_shape_with_text(slide, left, top_start, card_w, Inches(0.5),
                            name, color, font_size=16)

        # Card body
        add_multi_text(slide, left + Inches(0.1), top_start + Inches(0.6), card_w - Inches(0.2), Inches(4.5), [
            ("Purpose", 11, True, color),
            (purpose, 10, False, LIGHT_GRAY),
            ("", 6),
            ("Core Constraint", 11, True, ACCENT_RED),
            (constraint, 10, True, WHITE),
            ("", 6),
            ("Key Output", 11, True, color),
            (output, 10, False, LIGHT_GRAY),
            ("", 6),
            ("Invocation", 11, True, color),
            (invocation, 10, False, MED_GRAY),
        ])

    set_speaker_notes(slide, (
        "The constraints row is the key architectural element. The researcher cannot implement "
        "— so it searches instead of inventing. The planner cannot implement — so it sequences "
        "instead of coding. The implementor follows the plan — no creative decisions. The "
        "reviewer cannot modify code — only validate. Each constraint forces the AI into its "
        "correct mode of operation."
    ))


def slide_08_quality_and_clear(prs, quality_img):
    """Slide 8: Quality comparison + the /clear rule explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Quality Comparison + The /clear Rule", font_size=28, bold=True)
    add_accent_bar(slide, top=Inches(0.85))

    # Top: quality comparison table (built with shapes since we want it compact)
    dimensions = [
        ("Pattern matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code"),
        ("Knowledge transfer", "Lost between sessions", "Preserved in artifacts"),
        ("Rework", "30-50% discarded", "< 10% rework"),
        ("Validation", "Manual review only", "8-dimension structured review"),
    ]

    # Table headers
    add_shape_with_text(slide, Inches(1.5), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Traditional AI", ACCENT_RED, font_size=11)
    add_shape_with_text(slide, Inches(4.7), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Dimension", DARK_GRAY, font_size=11)
    add_shape_with_text(slide, Inches(7.9), Inches(1.2), Inches(3.5), Inches(0.35),
                        "RPI Approach", ACCENT_GREEN, font_size=11)

    for i, (dim, trad, rpi) in enumerate(dimensions):
        row_top = Inches(1.6 + i * 0.42)
        add_text_box(slide, Inches(1.5), row_top, Inches(3.0), Inches(0.38),
                     trad, font_size=10, color=ACCENT_RED)
        add_text_box(slide, Inches(4.7), row_top, Inches(3.0), Inches(0.38),
                     dim, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.9), row_top, Inches(3.5), Inches(0.38),
                     rpi, font_size=10, color=ACCENT_GREEN)

    # Bottom: /clear flow
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
                 "The /clear Rule: Context Reset Between Phases", font_size=16, bold=True)

    clear_flow = [
        ("Research", ACCENT_BLUE),
        ("/clear", ACCENT_RED),
        ("Plan", ACCENT_GREEN),
        ("/clear", ACCENT_RED),
        ("Implement", ACCENT_ORANGE),
        ("/clear", ACCENT_RED),
        ("Review", ACCENT_PURPLE),
    ]

    for i, (label, color) in enumerate(clear_flow):
        left = Inches(0.8 + i * 1.7)
        if label == "/clear":
            add_shape_with_text(slide, left, Inches(4.9), Inches(1.0), Inches(0.5),
                                label, ACCENT_RED, font_size=11,
                                shape=MSO_SHAPE.OVAL)
        else:
            add_shape_with_text(slide, left, Inches(4.9), Inches(1.4), Inches(0.5),
                                label, color, font_size=12)
        if i < len(clear_flow) - 1:
            add_connector_line(slide, left + Inches(1.4 if label != "/clear" else 1.0),
                               Inches(5.15),
                               left + Inches(1.6 if label !=
                                             "/clear" else 1.2),
                               Inches(5.15), color=MED_GRAY, width=Pt(2))

    add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.8),
                 "Context clearing prevents mode contamination — the researcher's reasoning "
                 "must not leak into the implementor's execution context.",
                 font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "For TPMs: the traceability row is your strongest argument for RPI adoption. "
        "Every line of code traces back through a plan to a research finding. "
        "For everyone: the /clear rule sounds counterintuitive — why throw away context? "
        "Because context contamination causes the researcher's reasoning patterns to leak "
        "into the implementor, creating mode confusion. Each phase needs a clean start "
        "with only its designated inputs."
    ))


# ── Slide Builders: Part 2 — "HVE-Core in Practice" (Slides 9–16) ──

def slide_09_chat_modes(prs):
    """Slide 9: Overview of all 6 custom Copilot chat modes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Six Custom Chat Modes", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    # Group 1: RPI Cycle
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                 "RPI Cycle", font_size=16, bold=True, color=ACCENT_BLUE)

    rpi_modes = [
        ("Task Researcher", "/task-research", "Cannot implement", ACCENT_BLUE),
        ("Task Planner", "/task-plan", "Cannot implement", ACCENT_GREEN),
        ("Task Implementor", "/task-implement",
         "Follows plan only", ACCENT_ORANGE),
        ("Task Reviewer", "/task-review", "Cannot modify code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(rpi_modes):
        top = Inches(2.0 + i * 0.95)
        add_shape_with_text(slide, Inches(0.8), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(3.5), top, Inches(1.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(3.5), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)
        # Arrow to next
        if i < len(rpi_modes) - 1:
            add_connector_line(slide, Inches(2.0), top + Inches(0.75),
                               Inches(2.0), top + Inches(0.9),
                               color=MED_GRAY, width=Pt(1))

    # Group 2: Complementary
    add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
                 "Complementary", font_size=16, bold=True, color=ACCENT_TEAL)

    comp_modes = [
        ("Prompt Builder", "/prompt-build", "Orchestrates subagents", ACCENT_TEAL),
        ("PR Review", "Agent picker", "Never modifies code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(comp_modes):
        top = Inches(2.0 + i * 1.2)
        add_shape_with_text(slide, Inches(7.5), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(10.2), top, Inches(2.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(10.2), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
                 "Every mode enforces specific constraints — AI cannot bypass its designated role",
                 font_size=13, color=MED_GRAY)

    set_speaker_notes(slide, (
        "Four modes form the RPI cycle. Two additional modes serve specialized functions. "
        "Prompt Builder creates the artifacts that configure the agents — it's meta-level. "
        "PR Review is an 8-dimension quality gate that complements Task Reviewer. "
        "Every mode has a constraint that prevents the AI from doing something specific. "
        "That's what makes the framework reliable."
    ))


def slide_10_demo_intro(prs):
    """Slide 10: Setup slide before switching to VS Code for live demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "Let's See It in Action", font_size=40, bold=True,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(2.3), color=ACCENT_GREEN)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
                 "Live Demo: RPI Workflow — Researching and Planning a Feature",
                 font_size=20, color=LIGHT_GRAY)

    # Watch-for items
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.4),
                 "Watch for:", font_size=16, bold=True, color=WHITE)

    watch_items = [
        ("1.", "How the researcher finds existing patterns (not invents code)", ACCENT_BLUE),
        ("2.", "How it cites specific files and line numbers", ACCENT_GREEN),
        ("3.", "How the plan references the research artifacts", ACCENT_ORANGE),
    ]

    for i, (num, desc, color) in enumerate(watch_items):
        top = Inches(4.2 + i * 0.7)
        add_shape_with_text(slide, Inches(1.0), top, Inches(0.5), Inches(0.5),
                            num, color, font_size=16)
        add_text_box(slide, Inches(1.7), top, Inches(10), Inches(0.5),
                     desc, font_size=16, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
                 "Demo: ~5 minutes in VS Code  •  Research → /clear → Plan",
                 font_size=13, color=MED_GRAY)

    set_speaker_notes(slide, (
        "I'm going to run through Research and Plan live. Watch for three things: "
        "First, how the researcher searches the codebase instead of inventing code. "
        "Second, how it cites specific files and line numbers. "
        "Third, how the plan references the research artifacts. "
        "Demo flow: open VS Code, run /task-research, show output, /clear, "
        "run /task-plan with research file, show plan and details output."
    ))


def slide_11_who_uses_hve_core(prs, role_img):
    """Slide 12 in deck: 8-role workflow mapping."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Who Uses HVE-Core?", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(role_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "HVE-Core isn't just for developers. TPMs use project-planning for Requirements → "
        "PRD → ADO work items. Platform engineers use coding-standards for artifact authoring. "
        "Security engineers use security-planning for threat models. Data scientists have "
        "their own Spec → Notebook → Dashboard workflow. Roles don't share workflows — "
        "collections make that separation automatic."
    ))


def slide_12_dogfooding_and_validation(prs):
    """Slide 13 in deck: Self-referential credibility + CI/CD pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Dogfooding + Enterprise Validation", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_TEAL)

    # Left column: Dogfooding
    add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        ("Self-Referential Credibility", 18, True, ACCENT_TEAL),
        ("", 8),
        ("• This presentation was planned using RPI", 14, False, LIGHT_GRAY),
        ("• Every agent was built by the Prompt Builder agent", 14, False, LIGHT_GRAY),
        ("• Instructions files follow their own rules", 14, False, LIGHT_GRAY),
        ("• Collection manifests are validated by their own CI", 14, False, LIGHT_GRAY),
        ("• Documentation is generated from the same artifacts", 14, False, LIGHT_GRAY),
        ("", 8),
        ("If HVE-Core doesn't work on itself,", 14, True, MED_GRAY),
        ("why would it work on your project?", 14, True, MED_GRAY),
    ])

    # Right column: Validation Pipeline
    add_multi_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5), [
        ("Validation Pipeline", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("Linting (7 jobs)", 14, True, ACCENT_BLUE),
        ("  Markdown, YAML, frontmatter, links,", 12, False, MED_GRAY),
        ("  tables, collections metadata, version consistency", 12, False, MED_GRAY),
        ("", 6),
        ("Analysis (2 jobs)", 14, True, ACCENT_GREEN),
        ("  PowerShell PSScriptAnalyzer, skill validation", 12, False, MED_GRAY),
        ("", 6),
        ("Security (3 jobs)", 14, True, ACCENT_RED),
        ("  Dependency pinning, SHA staleness, copyright", 12, False, MED_GRAY),
        ("", 6),
        ("Schema Validation", 14, True, ACCENT_PURPLE),
        ("  Collection YAML → plugin generation → extension packaging", 12, False, MED_GRAY),
    ])

    set_speaker_notes(slide, (
        "HVE-Core uses its own tools to build itself. This presentation was planned using "
        "the RPI workflow. Every agent was built by the Prompt Builder agent. The validation "
        "pipeline runs 12+ automated checks on every PR. If it doesn't work on itself, "
        "why would it work on your project? Self-referential credibility is the strongest "
        "argument for adoption."
    ))


def slide_13_extension_ecosystem(prs):
    """Slide 14 in deck: 8 VS Code extension packages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Extension Ecosystem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_PURPLE)

    # Table header
    headers = [("Extension", Inches(0.8), Inches(3.0)),
               ("Agents", Inches(4.0), Inches(1.2)),
               ("Prompts", Inches(5.3), Inches(1.2)),
               ("Instructions", Inches(6.6), Inches(1.5)),
               ("Focus", Inches(8.3), Inches(4.0))]

    for name, left, width in headers:
        add_shape_with_text(slide, left, Inches(1.5), width, Inches(0.4),
                            name, ACCENT_BLUE, font_size=11)

    extensions = [
        ("hve-core", "21", "23", "18", "Core RPI + all foundations", True),
        ("hve-ado", "9", "19", "10", "Azure DevOps integration", False),
        ("hve-github", "9", "19", "10", "GitHub backlog management", False),
        ("hve-project-planning", "13", "15", "5",
         "PRD / BRD / ADR workflows", False),
        ("hve-security-planning", "9", "16", "5", "Threat modeling", False),
        ("hve-rpi", "8", "14", "5", "Standalone RPI cycle", False),
        ("hve-prompt-engineering", "8", "14", "5", "Artifact authoring", False),
        ("hve-data-science", "1", "1", "0", "Data science workflows", False),
    ]

    for i, (ext, agents, prompts, instr, focus, highlight) in enumerate(extensions):
        row_top = Inches(2.0 + i * 0.5)
        bg_color = ACCENT_BLUE if highlight else CARD_BG

        # Extension name
        if highlight:
            add_shape_with_text(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                                ext, bg_color, font_size=11)
        else:
            add_text_box(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                         ext, font_size=11, color=LIGHT_GRAY)

        # Counts
        add_text_box(slide, Inches(4.0), row_top, Inches(1.2), Inches(0.4),
                     agents, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(5.3), row_top, Inches(1.2), Inches(0.4),
                     prompts, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.6), row_top, Inches(1.5), Inches(0.4),
                     instr, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Focus
        add_text_box(slide, Inches(8.3), row_top, Inches(4.0), Inches(0.4),
                     focus, font_size=11, color=MED_GRAY)

    # Bottom note
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Common base: 8 core agents + 5 core instructions shared across all extensions",
                 font_size=13, bold=True, color=ACCENT_TEAL)

    set_speaker_notes(slide, (
        "8 VS Code extensions organized by role. hve-core is the full package with 21 agents, "
        "23 prompts, and 18 instruction files. Specialized extensions like hve-ado and "
        "hve-github focus on specific workflows. All share a common base of 8 core agents "
        "and 5 core instructions. Pick the extension that matches your role."
    ))


def slide_14_learning_curve(prs):
    """Slide 15 in deck: Honest framing + payoff narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Learning Curve & Compounding Value", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_ORANGE)

    # Left: honest framing
    add_shape_with_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Honest Truth", ACCENT_ORANGE, font_size=14)
    add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("Your first RPI workflow will feel slower.", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("You'll wonder why you can't just ask AI to code it.", 14, False, LIGHT_GRAY),
        ("You'll feel like the constraints are slowing you down.", 14, False, LIGHT_GRAY),
        ("You'll want to skip the research phase.", 14, False, LIGHT_GRAY),
        ("", 8),
        ("That's normal. Stay with it.", 14, True, MED_GRAY),
    ])

    # Right: payoff
    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Payoff", ACCENT_GREEN, font_size=14)
    add_multi_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("By your third feature, the workflow", 18, True, ACCENT_GREEN),
        ("feels natural.", 18, True, ACCENT_GREEN),
        ("", 8),
        ("Research artifacts compound across features.", 14, False, LIGHT_GRAY),
        ("Plans reference previous research findings.", 14, False, LIGHT_GRAY),
        ("Rework drops below 10%.", 14, False, LIGHT_GRAY),
        ("", 8),
        ("The framework pays for itself.", 14, True, MED_GRAY),
    ])

    # Bottom: Version timeline
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
                 "Evolution Timeline", font_size=16, bold=True, color=WHITE)

    versions = [
        ("v1.1.0", "Jan 19", "Initial release", ACCENT_BLUE),
        ("v2.0.0", "Jan 28", "Breaking: agent restructure", ACCENT_RED),
        ("v2.3.4", "Feb 13", "Stable + extensions", ACCENT_GREEN),
    ]

    for i, (ver, date, desc, color) in enumerate(versions):
        left = Inches(1.0 + i * 3.8)
        add_shape_with_text(slide, left, Inches(5.7), Inches(1.2), Inches(0.45),
                            ver, color, font_size=12)
        add_text_box(slide, left + Inches(1.4), Inches(5.7), Inches(2.2), Inches(0.45),
                     f"{date} — {desc}", font_size=11, color=MED_GRAY)
        if i < len(versions) - 1:
            add_connector_line(slide, left + Inches(3.5), Inches(5.92),
                               left + Inches(3.7), Inches(5.92),
                               color=MED_GRAY, width=Pt(2))

    set_speaker_notes(slide, (
        "Be honest with your team: the first RPI workflow feels slower. That's by design. "
        "You're building research artifacts that compound across features. By the third "
        "feature, the workflow feels natural and you'll wonder how you worked without it. "
        "The version timeline shows rapid iteration — v1.1.0 in January, breaking restructure "
        "at v2.0.0 nine days later, stable at v2.3.4 two weeks after that."
    ))


def slide_15_getting_started(prs):
    """Slide 16 in deck: Three installation paths, CTA, and Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Getting Started", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    # 3 installation cards
    install_methods = [
        ("VS Code Extension", "10 seconds", "Zero-config\nInstall from marketplace",
         ACCENT_BLUE, "Recommended for most users"),
        ("Peer Clone", "2 minutes", "Customizable\nClone .github/ into your repo",
         ACCENT_ORANGE, "For teams with custom needs"),
        ("Codespaces", "1 click", "Zero-config\nPre-built dev environment",
         ACCENT_GREEN, "For quick exploration"),
    ]

    for i, (method, time, desc, color, note) in enumerate(install_methods):
        left = Inches(0.8 + i * 4.0)
        add_shape_with_text(slide, left, Inches(1.5), Inches(3.5), Inches(0.6),
                            f"{method}  •  {time}", color, font_size=14)
        add_multi_text(slide, left + Inches(0.2), Inches(2.3), Inches(3.2), Inches(1.5), [
            (desc, 13, False, LIGHT_GRAY),
            ("", 6),
            (note, 11, False, MED_GRAY),
        ])

    # GitHub URL
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
                 "github.com/microsoft/hve-core", font_size=20, bold=True,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # CTA
    add_shape_with_text(slide, Inches(3.0), Inches(4.9), Inches(7), Inches(0.6),
                        'Try /task-research on your next feature', ACCENT_GREEN, font_size=16)

    # Q&A
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1.0),
                 "Questions?", font_size=36, bold=True,
                 alignment=PP_ALIGN.CENTER, color=WHITE)

    set_speaker_notes(slide, (
        "Three ways to get started. VS Code Extension is fastest — 10 seconds, zero config. "
        "Peer Clone gives you customization — clone our .github/ folder into your repo. "
        "Codespaces gives you a full pre-built environment in one click. "
        "The call to action: try /task-research on your next feature. "
        "Don't change your whole workflow — just use the researcher once and see what happens. "
        "Resources: getting-started guide, first-workflow tutorial, docs/rpi/ documentation."
    ))


# ── Main Assembly ──────────────────────────────────────────

def main() -> int:
    """Generate the HVE-Core and RPI presentation deck."""
    print("Generating HVE-Core and RPI presentation...")

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(IMG_DIR, exist_ok=True)

    print("  Creating diagrams...")
    rpi_img = create_rpi_pipeline_diagram()
    quality_img = create_quality_comparison_diagram()
    role_img = create_role_mapping_diagram()

    print("  Building slides...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Part 1: Why RPI?
    slide_01_title(prs)
    slide_02_the_problem(prs)
    slide_03_why_it_happens(prs)
    slide_04_counterintuitive_insight(prs)
    slide_05_what_is_hve_core(prs)
    slide_06_rpi_pipeline(prs, rpi_img)
    slide_07_phase_deep_dive(prs)
    slide_08_quality_and_clear(prs, quality_img)

    # Part 2: HVE-Core in Practice
    slide_09_chat_modes(prs)
    slide_10_demo_intro(prs)
    # Slide 11 is live demo — no generated slide
    slide_11_who_uses_hve_core(prs, role_img)
    slide_12_dogfooding_and_validation(prs)
    slide_13_extension_ecosystem(prs)
    slide_14_learning_curve(prs)
    slide_15_getting_started(prs)

    output_path = os.path.join(OUTPUT_DIR, "hve-core-rpi-presentation.pptx")
    prs.save(output_path)
    print(f"\n✅ Saved presentation: {output_path}")
    print(f"   {len(prs.slides)} slides generated")
    return EXIT_SUCCESS


if __name__ == "__main__":
    sys.exit(main())
