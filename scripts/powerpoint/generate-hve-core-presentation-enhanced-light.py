#!/usr/bin/env python3
"""Generate enhanced HVE-Core and RPI Framework PowerPoint (light theme) — 20 slides with D-RPI, results, and community."""

# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx",
#   "Pillow",
# ]
# ///

import os
import io
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

# ── Constants ──────────────────────────────────────────────

LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_PURPLE = RGBColor(0x88, 0x64, 0xD8)
ACCENT_TEAL = RGBColor(0x03, 0x8E, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x24, 0x24, 0x24)
SECONDARY_TEXT = RGBColor(0x44, 0x44, 0x44)
TERTIARY_TEXT = RGBColor(0x66, 0x66, 0x66)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

EXIT_SUCCESS = 0
EXIT_FAILURE = 1

OUTPUT_DIR = os.path.join(os.path.dirname(
    os.path.abspath(__file__)), "..", "docs")
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(
    __file__)), "..", ".copilot-tracking", "presentation-assets")
os.makedirs(IMG_DIR, exist_ok=True)


# ── Helper Functions ───────────────────────────────────────

def set_slide_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, bold, color = line, default_size, False, default_color
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


def add_shape_with_text(slide, left, top, width, height, text, fill_color,
                        font_size=11, font_color=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = "Segoe UI"
    p.font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def add_connector_line(slide, start_x, start_y, end_x, end_y, color=ACCENT_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def set_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_accent_bar(slide, top=Inches(1.35), color=ACCENT_BLUE):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.8), Pt(4))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


# ── Font Loader ────────────────────────────────────────────

def _get_fonts():
    try:
        return {
            "title": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22),
            "zone": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16),
            "comp": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 13),
            "small": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 11),
            "label": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 12),
            "big": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28),
            "num": ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36),
        }
    except OSError:
        f = ImageFont.load_default()
        return {k: f for k in ("title", "zone", "comp", "small", "label", "big", "num")}


def _draw_rounded_box(draw, x, y, w, h, color, fill=(235, 235, 235)):
    draw.rounded_rectangle([(x, y), (x + w, y + h)],
                           radius=10, outline=color, width=2, fill=fill)


def _draw_dashed_rounded_box(draw, x, y, w, h, color, dash_len=10, gap_len=6):
    """Draw a rounded rectangle with a dashed border."""
    for dx in range(0, w, dash_len + gap_len):
        x1 = x + min(dx, w)
        x2 = x + min(dx + dash_len, w)
        draw.line([(x1, y), (x2, y)], fill=color, width=2)
    for dx in range(0, w, dash_len + gap_len):
        x1 = x + min(dx, w)
        x2 = x + min(dx + dash_len, w)
        draw.line([(x1, y + h), (x2, y + h)], fill=color, width=2)
    for dy in range(0, h, dash_len + gap_len):
        y1 = y + min(dy, h)
        y2 = y + min(dy + dash_len, h)
        draw.line([(x, y1), (x, y2)], fill=color, width=2)
    for dy in range(0, h, dash_len + gap_len):
        y1 = y + min(dy, h)
        y2 = y + min(dy + dash_len, h)
        draw.line([(x + w, y1), (x + w, y2)], fill=color, width=2)


def _draw_arrow(draw, x1, y1, x2, y2, color, label="", font=None):
    draw.line([(x1, y1), (x2, y2)], fill=color, width=2)
    dx, dy = x2 - x1, y2 - y1
    length = max(1, (dx * dx + dy * dy) ** 0.5)
    ux, uy = dx / length, dy / length
    px, py = -uy, ux
    draw.polygon([(x2, y2), (int(x2 - 8 * ux + 4 * px), int(y2 - 8 * uy + 4 * py)),
                  (int(x2 - 8 * ux - 4 * px), int(y2 - 8 * uy - 4 * py))], fill=color)
    if label and font:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        draw.text((mx, my - 15), label, fill=color, font=font)


def _draw_dashed_arrow(draw, x1, y1, x2, y2, color, dash_len=8, gap_len=6):
    """Draw a dashed line with an arrowhead at the end."""
    dx, dy = x2 - x1, y2 - y1
    length = max(1, (dx * dx + dy * dy) ** 0.5)
    ux, uy = dx / length, dy / length
    pos = 0
    while pos < length:
        sx = x1 + ux * pos
        sy = y1 + uy * pos
        end = min(pos + dash_len, length)
        ex = x1 + ux * end
        ey = y1 + uy * end
        draw.line([(int(sx), int(sy)), (int(ex), int(ey))], fill=color, width=2)
        pos += dash_len + gap_len
    px, py = -uy, ux
    draw.polygon([(x2, y2),
                  (int(x2 - 8 * ux + 4 * px), int(y2 - 8 * uy + 4 * py)),
                  (int(x2 - 8 * ux - 4 * px), int(y2 - 8 * uy - 4 * py))], fill=color)


# ── Diagram Generators ─────────────────────────────────────

def create_rpi_pipeline_diagram():
    """RPI pipeline flow: 4 phases connected by arrows with type labels and /clear markers."""
    w, h = 1600, 600
    img = Image.new("RGB", (w, h), (248, 248, 248))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "The RPI Pipeline",
              fill=(36, 36, 36), font=fonts["title"])

    phases = [
        ("Research", (0, 120, 212)),
        ("Plan", (16, 124, 16)),
        ("Implement", (255, 140, 0)),
        ("Review", (136, 100, 216)),
    ]

    type_labels = ["Uncertainty", "Knowledge",
                   "Strategy", "Working Code", "Validated Code"]

    box_w, box_h = 260, 100
    gap = 100
    start_x = 60
    y_top = 250

    for i, (name, color) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        _draw_rounded_box(draw, x, y_top, box_w, box_h, color)
        draw.text((x + box_w // 2 - len(name) * 5, y_top + 35),
                  name, fill=color, font=fonts["zone"])

        label_in = type_labels[i]
        draw.text((x + 10, y_top - 40), label_in,
                  fill=(100, 100, 100), font=fonts["comp"])

        if i < len(phases) - 1:
            ax1 = x + box_w + 4
            ax2 = x + box_w + gap - 4
            ay = y_top + box_h // 2
            _draw_arrow(draw, ax1, ay, ax2, ay, color)

            cx = x + box_w + gap // 2
            cy = ay
            r = 16
            draw.ellipse([(cx - r, cy - r), (cx + r, cy + r)],
                         fill=(209, 52, 56))
            draw.text((cx - 12, cy - 7), "/clear",
                      fill=(255, 255, 255), font=fonts["small"])

    last_x = start_x + 3 * (box_w + gap)
    draw.text((last_x + 100, y_top - 40),
              type_labels[4], fill=(100, 100, 100), font=fonts["comp"])

    # 88% stat callout
    draw.rounded_rectangle([(w // 2 - 150, y_top + box_h + 40),
                            (w // 2 + 150, y_top + box_h + 80)],
                           radius=8, fill=(16, 124, 16))
    draw.text((w // 2 - 130, y_top + box_h + 48),
              "Up to 88% faster task completion",
              fill=(255, 255, 255), font=fonts["comp"])

    draw.text((60, h - 60), "Red circles = /clear (context reset between phases)",
              fill=(209, 52, 56), font=fonts["comp"])
    draw.text((60, h - 35), "Each phase transforms one type of understanding into the next",
              fill=(100, 100, 100), font=fonts["small"])

    path = os.path.join(IMG_DIR, "light-hve-rpi-pipeline-enhanced.png")
    img.save(path, "PNG")
    return path


def create_drpi_pipeline_diagram():
    """D-RPI extended pipeline: optional Discovery + 4 RPI phases with iteration loops."""
    w, h = 1600, 650
    img = Image.new("RGB", (w, h), (248, 248, 248))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 250, 12), "D-RPI: Discovery + RPI Pipeline",
              fill=(36, 36, 36), font=fonts["title"])

    # Discovery box — dashed border
    disc_x, disc_y, disc_w, disc_h = 40, 260, 220, 100
    disc_color = (100, 100, 100)
    _draw_dashed_rounded_box(draw, disc_x, disc_y, disc_w, disc_h, disc_color)
    draw.text((disc_x + 40, disc_y + 35), "Discovery",
              fill=disc_color, font=fonts["zone"])
    draw.text((disc_x + 60, disc_y + disc_h + 8), "(optional)",
              fill=(100, 100, 100), font=fonts["small"])

    # 4 RPI phases
    phases = [
        ("Research", (0, 120, 212)),
        ("Plan", (16, 124, 16)),
        ("Implement", (255, 140, 0)),
        ("Review", (136, 100, 216)),
    ]
    type_labels = ["Empathize/\nDefine", "Uncertainty", "Knowledge",
                   "Strategy", "Working Code", "Validated Code"]

    box_w, box_h = 240, 100
    gap = 60
    rpi_start_x = 320
    y_top = 260

    for i, (name, color) in enumerate(phases):
        x = rpi_start_x + i * (box_w + gap)
        _draw_rounded_box(draw, x, y_top, box_w, box_h, color)
        draw.text((x + box_w // 2 - len(name) * 5, y_top + 35),
                  name, fill=color, font=fonts["zone"])

        label = type_labels[i + 1]
        draw.text((x + 10, y_top - 35), label,
                  fill=(100, 100, 100), font=fonts["comp"])

        if i < len(phases) - 1:
            ax1 = x + box_w + 4
            ax2 = x + box_w + gap - 4
            ay = y_top + box_h // 2
            _draw_arrow(draw, ax1, ay, ax2, ay, color)
            cx = x + box_w + gap // 2
            cy = ay
            r = 14
            draw.ellipse([(cx - r, cy - r), (cx + r, cy + r)],
                         fill=(209, 52, 56))
            draw.text((cx - 10, cy - 6), "/clear",
                      fill=(255, 255, 255), font=fonts["small"])

    # Type label for Discovery
    draw.text((disc_x + 10, disc_y - 35), type_labels[0],
              fill=(100, 100, 100), font=fonts["comp"])

    # Dashed arrow: Discovery → Research
    _draw_dashed_arrow(draw, disc_x + disc_w + 4, disc_y + disc_h // 2,
                       rpi_start_x - 4, y_top + box_h // 2, disc_color)

    # Output type label after last box
    last_x = rpi_start_x + 3 * (box_w + gap)
    draw.text((last_x + 80, y_top - 35),
              type_labels[5], fill=(100, 100, 100), font=fonts["comp"])

    # Iteration loop arrow: Review bottom → Research bottom
    loop_y = y_top + box_h + 40
    draw.line([(last_x + box_w // 2, y_top + box_h),
               (last_x + box_w // 2, loop_y)],
              fill=(136, 100, 216), width=2)
    draw.line([(last_x + box_w // 2, loop_y),
               (rpi_start_x + box_w // 2, loop_y)],
              fill=(136, 100, 216), width=2)
    draw.line([(rpi_start_x + box_w // 2, loop_y),
               (rpi_start_x + box_w // 2, y_top + box_h)],
              fill=(136, 100, 216), width=2)
    rx = rpi_start_x + box_w // 2
    draw.polygon([(rx, y_top + box_h),
                  (rx - 5, y_top + box_h + 10),
                  (rx + 5, y_top + box_h + 10)],
                 fill=(136, 100, 216))
    draw.text((w // 2 - 60, loop_y - 18), "iteration loops",
              fill=(136, 100, 216), font=fonts["small"])

    draw.text((40, h - 55), "Dashed border = optional Discovery phase",
              fill=(100, 100, 100), font=fonts["comp"])
    draw.text((40, h - 30), "Red circles = /clear (context reset)  |  Purple loop = Review can route back",
              fill=(100, 100, 100), font=fonts["small"])

    path = os.path.join(IMG_DIR, "light-hve-drpi-pipeline.png")
    img.save(path, "PNG")
    return path


def create_quality_comparison_diagram():
    """Side-by-side comparison: Traditional AI vs RPI approach across 5 dimensions."""
    w, h = 1600, 700
    img = Image.new("RGB", (w, h), (248, 248, 248))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "Quality Comparison",
              fill=(36, 36, 36), font=fonts["title"])

    col_left = 100
    col_right = 850
    col_w = 650

    draw.rounded_rectangle(
        [(col_left, 55), (col_left + col_w, 95)], radius=8, fill=(255, 230, 230))
    draw.text((col_left + col_w // 2 - 60, 65), "Traditional AI",
              fill=(36, 36, 36), font=fonts["zone"])

    draw.rounded_rectangle(
        [(col_right, 55), (col_right + col_w, 95)], radius=8, fill=(230, 255, 230))
    draw.text((col_right + col_w // 2 - 30, 65), "RPI",
              fill=(36, 36, 36), font=fonts["zone"])

    rows = [
        ("Pattern Matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code chain"),
        ("Knowledge Transfer", "Lost between sessions",
         "Preserved in research artifacts"),
        ("Rework", "30-50% of AI output discarded", "< 10% rework rate"),
        ("Validation", "Manual review only", "Structured multi-dimension review"),
    ]

    for i, (dimension, trad, rpi) in enumerate(rows):
        ry = 120 + i * 110

        draw.rounded_rectangle([(col_left, ry), (col_left + col_w + col_right - col_left + col_w - col_left, ry + 30)],
                               radius=4, fill=(235, 235, 235))
        draw.text((w // 2 - len(dimension) * 4, ry + 6), dimension,
                  fill=(36, 36, 36), font=fonts["label"])

        draw.rounded_rectangle([(col_left, ry + 35), (col_left + col_w, ry + 90)], radius=8,
                               fill=(255, 240, 240), outline=(209, 52, 56), width=1)
        draw.text((col_left + 20, ry + 50), trad,
                  fill=(180, 40, 40), font=fonts["comp"])

        draw.rounded_rectangle([(col_right, ry + 35), (col_right + col_w, ry + 90)], radius=8,
                               fill=(240, 255, 240), outline=(16, 124, 16), width=1)
        draw.text((col_right + 20, ry + 50), rpi,
                  fill=(15, 100, 15), font=fonts["comp"])

    path = os.path.join(IMG_DIR, "light-hve-quality-comparison-enhanced.png")
    img.save(path, "PNG")
    return path


def create_role_mapping_diagram():
    """8 roles mapped to collections and workflow patterns in a 2x4 grid."""
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 248, 248))
    draw = ImageDraw.Draw(img)
    fonts = _get_fonts()

    draw.text((w // 2 - 200, 12), "Eight Roles, One Framework",
              fill=(36, 36, 36), font=fonts["title"])

    roles = [
        ("Developer", "hve-core", "RPI Pipeline", (0, 120, 212)),
        ("TPM / Lead", "project-planning", "Reqs → PRD → WIT", (16, 124, 16)),
        ("Platform Engineer", "coding-standards",
         "Artifact authoring", (255, 140, 0)),
        ("OSS Contributor", "github", "Backlog management", (136, 100, 216)),
        ("Security Engineer", "security-planning",
         "Threat modeling", (209, 52, 56)),
        ("Data Scientist", "data-science",
         "Spec → Notebook → Dashboard", (3, 142, 142)),
        ("Project Planner", "project-planning", "PRD / BRD / ADR", (0, 120, 212)),
        ("UX / DT Practitioner", "design-thinking",
         "9-method coaching", (136, 100, 216)),
    ]

    card_w, card_h = 340, 170
    cols, rows_count = 4, 2
    x_start = 40
    y_start = 60
    x_gap = (w - 2 * x_start - cols * card_w) // (cols - 1)
    y_gap = 40

    for idx, (role, collection, workflow, color) in enumerate(roles):
        col = idx % cols
        row = idx // cols
        x = x_start + col * (card_w + x_gap)
        y = y_start + row * (card_h + y_gap)

        _draw_rounded_box(draw, x, y, card_w, card_h, color)

        draw.text((x + 15, y + 15), role, fill=color, font=fonts["zone"])
        draw.text((x + 15, y + 50),
                  f"Collection: {collection}", fill=(100, 100, 100), font=fonts["comp"])
        draw.rounded_rectangle([(x + 15, y + 85), (x + card_w - 15, y + 130)],
                               radius=6, fill=(248, 248, 248), outline=color, width=1)
        draw.text((x + 25, y + 95), workflow, fill=color, font=fonts["comp"])
        draw.line([(x + 15, y + 75), (x + card_w - 15, y + 75)],
                  fill=color, width=1)

    draw.text((40, h - 40), "Each role maps to a specific collection and workflow pattern",
              fill=(100, 100, 100), font=fonts["small"])

    path = os.path.join(IMG_DIR, "light-hve-role-mapping-enhanced.png")
    img.save(path, "PNG")
    return path


# ── Slide Builders: Part 1 — "Why RPI?" (Slides 1–10) ─────

def slide_01_title(prs):
    """Slide 1: Title card with industry credibility strip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
                 "HVE-Core and the RPI Framework", font_size=44, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(3.0), color=ACCENT_BLUE)
    add_multi_text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(2), [
        ("Turning AI from a Code Generator into a Research Partner", 20, False, SECONDARY_TEXT),
        ("", 10),
        ("A constraint-based framework for AI-assisted engineering", 16, False, TERTIARY_TEXT),
    ])
    add_shape_with_text(slide, Inches(0.8), Inches(5.2), Inches(3.0), Inches(0.5),
                        "30-Minute Overview  •  Two Parts", ACCENT_BLUE, font_size=12)

    # Industry credibility strip
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
                 "Proven at AVEVA  •  BMW  •  Michelin  •  Hexagon  •  Kubota  •  Nvidia",
                 font_size=14, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "Welcome. Over the next 30 minutes I'll show you why the way most teams use AI "
        "coding assistants is fundamentally broken — and how a constraint-based framework "
        "called RPI fixes it. We'll cover the problem, the framework, a live demo, and "
        "how different roles use it today.\n\n"
        "HVE-Core has been industry-proven at companies like AVEVA, BMW, Michelin, Hexagon, "
        "Kubota, and Nvidia to improve coding accuracy, automate repetitive tasks, and "
        "personalize solutions."
    ))


def slide_02_the_problem(prs):
    """Slide 2: The universally relatable AI failure mode with pain point statistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Problem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_RED)

    add_multi_text(slide, Inches(1.0), Inches(1.5), Inches(7.5), Inches(4.0), [
        ("", 8),
        ('You: "Build me a Terraform module for Azure IoT"', 18, False, DARK_TEXT),
        ("", 12),
        ("AI: *immediately generates 2,000 lines of code*", 18, False, SECONDARY_TEXT),
        ("", 12),
        ("Reality:", 20, True, ACCENT_RED),
        ("  • Missing provider dependencies", 14, False, ACCENT_RED),
        ("  • Wrong variable naming conventions", 14, False, ACCENT_RED),
        ("  • Patterns from 3 years ago", 14, False, ACCENT_RED),
        ("  • Broke two downstream services", 14, False, ACCENT_RED),
        ("  • Used modules that don't exist in your registry", 14, False, ACCENT_RED),
        ("", 12),
        ("It looked right. It compiled. Then you deployed.", 14, True, TERTIARY_TEXT),
    ], default_size=14, default_color=DARK_TEXT)

    # Pain point statistics panel
    add_shape_with_text(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(0.5),
                        "The Cost", ACCENT_RED, font_size=14)
    add_multi_text(slide, Inches(9.0), Inches(2.2), Inches(3.5), Inches(3.0), [
        ("15–20 hrs/week", 20, True, ACCENT_ORANGE),
        ("lost to repetitive tasks and", 12, False, SECONDARY_TEXT),
        ("relearning context", 12, False, SECONDARY_TEXT),
        ("", 10),
        ("40% productivity loss", 20, True, ACCENT_ORANGE),
        ("from unstructured AI interactions", 12, False, SECONDARY_TEXT),
        ("and context-switching", 12, False, SECONDARY_TEXT),
    ])

    add_text_box(slide, Inches(1.0), Inches(6.5), Inches(10), Inches(0.5),
                 "Raise your hand if this has happened to you.", font_size=14, color=TERTIARY_TEXT)
    set_speaker_notes(slide, (
        "Raise your hand if this has happened to you. The AI generates something that looks "
        "right. It compiles. Then you deploy and discover it used patterns from 3 years ago, "
        "missed your naming conventions, and broke two downstream services. "
        "This is the universal failure mode of AI-assisted development.\n\n"
        "The pain is quantifiable: developers lose 15 to 20 hours per week dealing with "
        "repetitive tasks and relearning context. Unstructured AI interactions and constant "
        "context-switching cost about 40% of productivity."
    ))


def slide_03_why_it_happens(prs):
    """Slide 3: Root cause — AI conflates investigation with implementation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Why It Happens", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_shape_with_text(slide, Inches(1.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Investigate?\n\n🔍\n\nSearch • Analyze • Understand",
                        CARD_BG, font_size=16, font_color=ACCENT_BLUE)

    add_shape_with_text(slide, Inches(7.5), Inches(1.8), Inches(4.0), Inches(2.5),
                        "Implement!\n\n⌨️\n\nGenerate • Write • Ship",
                        CARD_BG, font_size=16, font_color=ACCENT_ORANGE)

    add_connector_line(slide, Inches(5.7), Inches(3.0), Inches(7.3), Inches(3.0),
                       color=ACCENT_RED, width=Pt(3))

    add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9), Inches(0.8),
                 "AI Writes First and Thinks Never", font_size=28, bold=True,
                 color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.0), Inches(5.7), Inches(9), Inches(0.5),
                 '"Plausible" and "correct" aren\'t the same thing',
                 font_size=16, color=TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "The root cause: AI treats every request as an implementation request. "
        "It never separates investigation from implementation. It pattern-matches "
        "from training data instead of searching your actual codebase. "
        "The output looks plausible — but plausible and correct aren't the same thing."
    ))


def slide_04_counterintuitive_insight(prs):
    """Slide 4: The key insight — constraints, not intelligence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Counterintuitive Insight", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_BLUE)

    add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(2.5), [
        ('"The solution isn\'t teaching AI to be smarter.', 22, False, DARK_TEXT),
        ("", 8),
        ('It\'s preventing AI from doing certain things', 22, False, DARK_TEXT),
        ('at certain times."', 22, False, DARK_TEXT),
    ])
    add_accent_bar(slide, top=Inches(3.8), color=ACCENT_BLUE)

    add_shape_with_text(slide, Inches(1.5), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for plausible code", ACCENT_RED, font_size=16)
    add_connector_line(slide, Inches(6.2), Inches(5.15), Inches(7.0), Inches(5.15),
                       color=DARK_TEXT, width=Pt(3))
    add_shape_with_text(slide, Inches(7.2), Inches(4.8), Inches(4.5), Inches(0.7),
                        "Optimizing for verified truth", ACCENT_GREEN, font_size=16)

    add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
                 "Constraints are features, not limitations",
                 font_size=16, color=TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "This is the philosophical foundation of HVE-Core. We don't try to make AI smarter. "
        "We constrain what it can do at each stage. A researcher that cannot write code will "
        "search instead of inventing. A planner that cannot implement will sequence instead "
        "of coding. Constraints force better behavior."
    ))


def slide_05_what_is_hve_core(prs):
    """Slide 5: Elevator pitch, component summary, and industry adoption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "What is HVE-Core?", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 "An enterprise-ready prompt engineering framework that transforms "
                 "GitHub Copilot from a code-completion tool into a structured engineering partner.",
                 font_size=15, color=SECONDARY_TEXT)

    components = [
        ("Agents", "22", ACCENT_BLUE),
        ("Prompts", "27", ACCENT_GREEN),
        ("Instructions", "24", ACCENT_ORANGE),
        ("Skills", "1", ACCENT_TEAL),
        ("Collections", "10", ACCENT_PURPLE),
    ]
    card_w = Inches(2.1)
    for i, (name, count, color) in enumerate(components):
        left = Inches(0.8 + i * 2.4)
        add_shape_with_text(slide, left, Inches(2.4), card_w, Inches(1.2),
                            f"{name}\n{count}", color, font_size=18)

    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
                 "Delegation Flow", font_size=18, bold=True, color=DARK_TEXT)

    flow_items = ["User", "Prompt", "Agent", "Instructions"]
    flow_colors = [TERTIARY_TEXT, ACCENT_GREEN, ACCENT_BLUE, ACCENT_ORANGE]
    for i, (item, color) in enumerate(zip(flow_items, flow_colors)):
        left = Inches(1.0 + i * 2.8)
        add_shape_with_text(slide, left, Inches(4.6), Inches(2.0), Inches(0.6),
                            item, CARD_BG, font_size=14, font_color=color)
        if i < len(flow_items) - 1:
            add_connector_line(slide, left + Inches(2.1), Inches(4.9),
                               left + Inches(2.7), Inches(4.9), color=color, width=Pt(2))

    add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
        ("Agents define behavior and constraints. Instructions encode standards.",
         13, False, TERTIARY_TEXT),
        ("Prompts are user entry points. Skills package domain knowledge.",
         13, False, TERTIARY_TEXT),
        ("Collections bundle everything for specific roles.", 13, False, TERTIARY_TEXT),
        ("", 6),
        ("Industry-proven at AVEVA, BMW, Michelin, Hexagon, Kubota, and Nvidia",
         13, True, ACCENT_TEAL),
    ])

    set_speaker_notes(slide, (
        "HVE-Core is an enterprise-ready prompt engineering framework. 22 custom agents, "
        "27 prompts, 24 instruction files, 1 skill, and 10 collections. Agents define what "
        "AI can and cannot do. Instructions encode coding standards. Prompts are how users "
        "interact. Collections bundle artifacts for specific roles.\n\n"
        "This isn't just theory. HVE-Core has been used by industry leaders like AVEVA, "
        "BMW, Michelin, Hexagon, Kubota, and Nvidia."
    ))


def slide_06_rpi_pipeline(prs, rpi_img):
    """Slide 6: The RPI type transformation pipeline diagram with 88% stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "The RPI Pipeline", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(rpi_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "Each phase converts one form of understanding into the next. "
        "Internal trials show up to 88% faster task completion using RPI compared to "
        "unstructured AI prompting."
    ))


def slide_07_phase_deep_dive(prs):
    """Slide 7: Attribute cards for all 4 RPI phases with best practice tip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Deep Dive: Each Phase", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    phases = [
        ("Research", ACCENT_BLUE,
         "Discover what exists",
         "CANNOT implement or modify code",
         "research.md with citations",
         "/task-research"),
        ("Plan", ACCENT_GREEN,
         "Sequence the work",
         "CANNOT implement or modify code",
         "plan.instructions.md + details.md",
         "/task-plan"),
        ("Implement", ACCENT_ORANGE,
         "Execute the plan",
         "MUST follow plan — no creative decisions",
         "Working code matching plan",
         "/task-implement"),
        ("Review", ACCENT_PURPLE,
         "Validate the output",
         "CANNOT modify code — only report findings",
         "8-dimension review report",
         "/task-review"),
    ]

    card_w = Inches(2.9)
    for i, (name, color, purpose, constraint, output, invocation) in enumerate(phases):
        left = Inches(0.5 + i * 3.1)
        top_start = Inches(1.5)

        add_shape_with_text(slide, left, top_start, card_w, Inches(0.5),
                            name, color, font_size=16)

        add_multi_text(slide, left + Inches(0.1), top_start + Inches(0.6), card_w - Inches(0.2), Inches(4.5), [
            ("Purpose", 11, True, color),
            (purpose, 10, False, SECONDARY_TEXT),
            ("", 6),
            ("Core Constraint", 11, True, ACCENT_RED),
            (constraint, 10, True, DARK_TEXT),
            ("", 6),
            ("Key Output", 11, True, color),
            (output, 10, False, SECONDARY_TEXT),
            ("", 6),
            ("Invocation", 11, True, color),
            (invocation, 10, False, TERTIARY_TEXT),
        ])

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 'Best practice: Keep implementation to ~3 files per RPI cycle for manageable scope',
                 font_size=12, bold=True, color=ACCENT_TEAL)

    set_speaker_notes(slide, (
        "The constraints row is the key architectural element. Each constraint forces the AI "
        "into its correct mode of operation. Keep each RPI cycle scoped to about 3 files."
    ))


def slide_08_quality_and_clear(prs, quality_img):
    """Slide 8: Quality comparison + /clear rule + paradigm shift quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Quality Comparison + The /clear Rule", font_size=28, bold=True)
    add_accent_bar(slide, top=Inches(0.85))

    dimensions = [
        ("Pattern matching", "Invents plausible patterns",
         "Finds actual codebase patterns"),
        ("Traceability", "No audit trail", "Research → Plan → Code"),
        ("Knowledge transfer", "Lost between sessions", "Preserved in artifacts"),
        ("Rework", "30-50% discarded", "< 10% rework"),
        ("Validation", "Manual review only", "8-dimension structured review"),
    ]

    add_shape_with_text(slide, Inches(1.5), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Traditional AI", ACCENT_RED, font_size=11)
    add_shape_with_text(slide, Inches(4.7), Inches(1.2), Inches(3.0), Inches(0.35),
                        "Dimension", BORDER_GRAY, font_size=11, font_color=DARK_TEXT)
    add_shape_with_text(slide, Inches(7.9), Inches(1.2), Inches(3.5), Inches(0.35),
                        "RPI Approach", ACCENT_GREEN, font_size=11)

    for i, (dim, trad, rpi) in enumerate(dimensions):
        row_top = Inches(1.6 + i * 0.42)
        add_text_box(slide, Inches(1.5), row_top, Inches(3.0), Inches(0.38),
                     trad, font_size=10, color=ACCENT_RED)
        add_text_box(slide, Inches(4.7), row_top, Inches(3.0), Inches(0.38),
                     dim, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.9), row_top, Inches(3.5), Inches(0.38),
                     rpi, font_size=10, color=ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
                 "The /clear Rule: Context Reset Between Phases", font_size=16, bold=True)

    clear_flow = [
        ("Research", ACCENT_BLUE),
        ("/clear", ACCENT_RED),
        ("Plan", ACCENT_GREEN),
        ("/clear", ACCENT_RED),
        ("Implement", ACCENT_ORANGE),
        ("/clear", ACCENT_RED),
        ("Review", ACCENT_PURPLE),
    ]

    for i, (label, color) in enumerate(clear_flow):
        left = Inches(0.8 + i * 1.7)
        if label == "/clear":
            add_shape_with_text(slide, left, Inches(4.5), Inches(1.0), Inches(0.45),
                                label, ACCENT_RED, font_size=11,
                                shape=MSO_SHAPE.OVAL)
        else:
            add_shape_with_text(slide, left, Inches(4.5), Inches(1.4), Inches(0.45),
                                label, color, font_size=12)
        if i < len(clear_flow) - 1:
            add_connector_line(slide, left + Inches(1.4 if label != "/clear" else 1.0),
                               Inches(4.72),
                               left + Inches(1.6 if label != "/clear" else 1.2),
                               Inches(4.72), color=BORDER_GRAY, width=Pt(2))

    add_multi_text(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(1.2), [
        ("", 6),
        ('Stop asking AI: "Write this code."', 15, True, DARK_TEXT),
        ('Start asking: "Help me research, plan, then implement with evidence."',
         15, True, ACCENT_BLUE),
    ])

    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
                 "Context clearing prevents mode contamination — each phase gets a clean start.",
                 font_size=12, color=TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "The traceability row is the strongest argument for RPI adoption. "
        "The /clear rule sounds counterintuitive — why throw away context? "
        "Because context contamination causes mode confusion.\n\n"
        "The paradigm shift: stop asking AI to write code. Start asking it to help "
        "you research, plan, and then implement with evidence."
    ))


def slide_09_chat_modes(prs):
    """Slide 9: Overview of all 6 custom Copilot chat modes with rpi-agent note."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Six Custom Chat Modes", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1))

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                 "RPI Cycle", font_size=16, bold=True, color=ACCENT_BLUE)

    rpi_modes = [
        ("Task Researcher", "/task-research", "Cannot implement", ACCENT_BLUE),
        ("Task Planner", "/task-plan", "Cannot implement", ACCENT_GREEN),
        ("Task Implementor", "/task-implement",
         "Follows plan only", ACCENT_ORANGE),
        ("Task Reviewer", "/task-review", "Cannot modify code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(rpi_modes):
        top = Inches(2.0 + i * 0.95)
        add_shape_with_text(slide, Inches(0.8), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(3.5), top, Inches(1.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(3.5), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)
        if i < len(rpi_modes) - 1:
            add_connector_line(slide, Inches(2.0), top + Inches(0.75),
                               Inches(2.0), top + Inches(0.9),
                               color=BORDER_GRAY, width=Pt(1))

    add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
                 "Complementary", font_size=16, bold=True, color=ACCENT_TEAL)

    comp_modes = [
        ("Prompt Builder", "/prompt-build", "Orchestrates subagents", ACCENT_TEAL),
        ("PR Review", "Agent picker", "Never modifies code", ACCENT_PURPLE),
    ]

    for i, (name, command, constraint, color) in enumerate(comp_modes):
        top = Inches(2.0 + i * 1.2)
        add_shape_with_text(slide, Inches(7.5), top, Inches(2.5), Inches(0.7),
                            name, color, font_size=12)
        add_text_box(slide, Inches(10.2), top, Inches(2.5), Inches(0.35),
                     command, font_size=11, color=color, font_name="Segoe UI")
        add_text_box(slide, Inches(10.2), top + Inches(0.35), Inches(3.0), Inches(0.35),
                     f"Constraint: {constraint}", font_size=10, color=ACCENT_RED)

    add_multi_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.0), [
        ("Every mode enforces specific constraints — AI cannot bypass its designated role",
         13, False, TERTIARY_TEXT),
        ("", 6),
        ("Advanced: rpi-agent orchestrates all 4 phases in a single session for familiar tasks",
         11, False, ACCENT_TEAL),
    ])

    set_speaker_notes(slide, (
        "Four modes form the RPI cycle. Two additional modes serve specialized functions. "
        "For familiar tasks, the rpi-agent can orchestrate all four phases in one session."
    ))


def slide_10_demo_intro(prs):
    """Slide 10: Setup slide before switching to VS Code for live demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "Let's See It in Action", font_size=40, bold=True,
                 alignment=PP_ALIGN.LEFT)
    add_accent_bar(slide, top=Inches(2.3), color=ACCENT_GREEN)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
                 "Live Demo: RPI Workflow — Researching and Planning a Feature",
                 font_size=20, color=SECONDARY_TEXT)

    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.4),
                 "Watch for:", font_size=16, bold=True, color=DARK_TEXT)

    watch_items = [
        ("1.", "How the researcher finds existing patterns (not invents code)", ACCENT_BLUE),
        ("2.", "How it cites specific files and line numbers", ACCENT_GREEN),
        ("3.", "How the plan references the research artifacts", ACCENT_ORANGE),
        ("4.", "How /clear resets context between phases", ACCENT_RED),
        ("5.", "How the implementor follows the plan step by step", ACCENT_PURPLE),
    ]

    for i, (num, desc, color) in enumerate(watch_items):
        top = Inches(4.1 + i * 0.5)
        add_shape_with_text(slide, Inches(1.0), top, Inches(0.5), Inches(0.4),
                            num, color, font_size=14)
        add_text_box(slide, Inches(1.7), top, Inches(10), Inches(0.4),
                     desc, font_size=14, color=SECONDARY_TEXT)

    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
                 "Demo: ~5 minutes in VS Code  •  Research → /clear → Plan → /clear → Implement",
                 font_size=13, color=TERTIARY_TEXT)

    set_speaker_notes(slide, (
        "Watch for five things: how the researcher searches rather than invents, "
        "how it cites files and line numbers, how the plan references research artifacts, "
        "how /clear resets context, and how the implementor follows the plan."
    ))


# ── Slide Builders: Part 2 — "HVE-Core in Practice" (Slides 11–20) ──

def slide_11_who_uses_hve_core(prs, role_img):
    """Slide 11: 8-role workflow mapping."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Eight Roles, One Framework", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(0.9))
    slide.shapes.add_picture(role_img, Inches(
        0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    set_speaker_notes(slide, (
        "HVE-Core isn't just for developers. Eight distinct roles, each with a dedicated "
        "collection that filters the right tools for their job."
    ))


def slide_12_discovery_design_thinking(prs, drpi_img):
    """Slide 12: D-RPI extended pipeline and Design Thinking alignment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Beyond RPI: Discovery + Design Thinking", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_TEAL)

    slide.shapes.add_picture(drpi_img, Inches(
        0.3), Inches(1.3), Inches(12.7), Inches(3.5))

    add_text_box(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.4),
                 "Design Thinking Alignment", font_size=16, bold=True, color=DARK_TEXT)

    dt_rows = [
        ("Discovery", "Empathize / Define", TERTIARY_TEXT),
        ("Research", "Ideate (evidence-constrained)", ACCENT_BLUE),
        ("Plan", "Prototype (on paper)", ACCENT_GREEN),
    ]
    dt_rows_right = [
        ("Implement", "Build / Test", ACCENT_ORANGE),
        ("Review", "Validate / Feedback", ACCENT_PURPLE),
    ]

    for i, (phase, dt_stage, color) in enumerate(dt_rows):
        row_top = Inches(5.4 + i * 0.35)
        add_shape_with_text(slide, Inches(1.0), row_top, Inches(1.8), Inches(0.3),
                            phase, color, font_size=10)
        add_text_box(slide, Inches(3.0), row_top, Inches(3.0), Inches(0.3),
                     f"→  {dt_stage}", font_size=10, color=SECONDARY_TEXT)

    for i, (phase, dt_stage, color) in enumerate(dt_rows_right):
        row_top = Inches(5.4 + i * 0.35)
        add_shape_with_text(slide, Inches(7.0), row_top, Inches(1.8), Inches(0.3),
                            phase, color, font_size=10)
        add_text_box(slide, Inches(9.0), row_top, Inches(3.0), Inches(0.3),
                     f"→  {dt_stage}", font_size=10, color=SECONDARY_TEXT)

    add_text_box(slide, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
                 "D-RPI and Design Thinking integration are forward-looking extensions. "
                 "Core four-phase RPI is the production workflow today.",
                 font_size=11, color=TERTIARY_TEXT)

    set_speaker_notes(slide, (
        "For complex or ambiguous projects, add a Discovery phase before Research — D-RPI. "
        "Discovery mirrors Empathize and Define stages of Design Thinking. "
        "D-RPI and Design Thinking integration are forward-looking extensions."
    ))


def slide_13_real_results(prs):
    """Slide 13: Concrete case study metrics from field deployments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Real Results with HVE-Core", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    metrics = [
        ("50%", "Faster\nDeployments", "Global telecom:\nAI-driven IaC", ACCENT_GREEN),
        ("90%", "Faster\nDocs", "Architecture docs\nand security plans:\ndays → hours", ACCENT_GREEN),
        ("2 Days", "vs 8 Weeks", "Internal hackathon:\nRPI + Copilot\nprototype", ACCENT_GREEN),
        ("Quality", "Improved", "AI PR reviews\ncatch subtle bugs\npre-merge", ACCENT_BLUE),
        ("Dev\nExperience", "Better", "Less rework,\nless context-switching,\nhigher satisfaction", ACCENT_TEAL),
    ]

    card_w = Inches(2.3)
    for i, (number, subtitle, detail, color) in enumerate(metrics):
        left = Inches(0.5 + i * 2.5)

        add_shape_with_text(slide, left, Inches(1.6), card_w, Inches(1.0),
                            number, color, font_size=36)
        add_shape_with_text(slide, left, Inches(2.6), card_w, Inches(0.8),
                            subtitle, CARD_BG, font_size=14, font_color=SECONDARY_TEXT)
        add_multi_text(slide, left + Inches(0.1), Inches(3.6), card_w - Inches(0.2), Inches(1.5), [
            (detail, 10, False, TERTIARY_TEXT),
        ])

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                 "Real teams.  Real metrics.  Real improvements.",
                 font_size=16, bold=True, color=SECONDARY_TEXT, alignment=PP_ALIGN.CENTER)

    add_shape_with_text(slide, Inches(3.0), Inches(6.2), Inches(7.0), Inches(0.5),
                        "Internal trials: up to 88% faster task completion with RPI",
                        ACCENT_GREEN, font_size=13)

    set_speaker_notes(slide, (
        "These aren't hypothetical — a global telecom cut deployment times by 50%. "
        "Architecture and security docs went from days to hours. A hackathon team delivered "
        "a prototype in 2 days instead of 8 weeks. Developers are happier."
    ))


def slide_14_dogfooding_and_validation(prs):
    """Slide 14: Self-referential credibility + CI/CD pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Built by HVE-Core, Validated by HVE-Core", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_TEAL)

    add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
        ("Self-Referential Credibility", 18, True, ACCENT_TEAL),
        ("", 8),
        ("• This presentation was planned using RPI", 14, False, SECONDARY_TEXT),
        ("• Every agent was built by the Prompt Builder agent", 14, False, SECONDARY_TEXT),
        ("• Instructions files follow their own rules", 14, False, SECONDARY_TEXT),
        ("• Collection manifests are validated by their own CI", 14, False, SECONDARY_TEXT),
        ("• Documentation is generated from the same artifacts", 14, False, SECONDARY_TEXT),
        ("", 8),
        ("If HVE-Core doesn't work on itself,", 14, True, TERTIARY_TEXT),
        ("why would it work on your project?", 14, True, TERTIARY_TEXT),
    ])

    add_multi_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5), [
        ("Validation Pipeline", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("Linting (7 jobs)", 14, True, ACCENT_BLUE),
        ("  Markdown, YAML, frontmatter, links,", 12, False, TERTIARY_TEXT),
        ("  tables, collections metadata, version consistency", 12, False, TERTIARY_TEXT),
        ("", 6),
        ("Analysis (2 jobs)", 14, True, ACCENT_GREEN),
        ("  PowerShell PSScriptAnalyzer, skill validation", 12, False, TERTIARY_TEXT),
        ("", 6),
        ("Security (3 jobs)", 14, True, ACCENT_RED),
        ("  Dependency pinning, SHA staleness, copyright", 12, False, TERTIARY_TEXT),
        ("", 6),
        ("Schema Validation", 14, True, ACCENT_PURPLE),
        ("  Collection YAML → plugin generation → extension packaging", 12, False, TERTIARY_TEXT),
    ])

    set_speaker_notes(slide, (
        "HVE-Core uses its own tools to build itself. This presentation was planned using "
        "the RPI workflow. The validation pipeline runs 12+ automated checks on every PR."
    ))


def slide_15_extension_ecosystem(prs):
    """Slide 15: 8 VS Code extension packages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Extension Ecosystem", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_PURPLE)

    headers = [("Extension", Inches(0.8), Inches(3.0)),
               ("Agents", Inches(4.0), Inches(1.2)),
               ("Prompts", Inches(5.3), Inches(1.2)),
               ("Instructions", Inches(6.6), Inches(1.5)),
               ("Focus", Inches(8.3), Inches(4.0))]

    for name, left, width in headers:
        add_shape_with_text(slide, left, Inches(1.5), width, Inches(0.4),
                            name, ACCENT_BLUE, font_size=11)

    extensions = [
        ("hve-core", "21", "23", "18", "Core RPI + all foundations", True),
        ("hve-ado", "9", "19", "10", "Azure DevOps integration", False),
        ("hve-github", "9", "19", "10", "GitHub backlog management", False),
        ("hve-project-planning", "13", "15", "5",
         "PRD / BRD / ADR workflows", False),
        ("hve-security-planning", "9", "16", "5", "Threat modeling", False),
        ("hve-rpi", "8", "14", "5", "Standalone RPI cycle", False),
        ("hve-prompt-engineering", "8", "14", "5", "Artifact authoring", False),
        ("hve-data-science", "1", "1", "0", "Data science workflows", False),
    ]

    for i, (ext, agents, prompts, instr, focus, highlight) in enumerate(extensions):
        row_top = Inches(2.0 + i * 0.5)
        bg_color = ACCENT_BLUE if highlight else CARD_BG

        if highlight:
            add_shape_with_text(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                                ext, bg_color, font_size=11)
        else:
            add_text_box(slide, Inches(0.8), row_top, Inches(3.0), Inches(0.4),
                         ext, font_size=11, color=SECONDARY_TEXT)

        add_text_box(slide, Inches(4.0), row_top, Inches(1.2), Inches(0.4),
                     agents, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(5.3), row_top, Inches(1.2), Inches(0.4),
                     prompts, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.6), row_top, Inches(1.5), Inches(0.4),
                     instr, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.3), row_top, Inches(4.0), Inches(0.4),
                     focus, font_size=11, color=TERTIARY_TEXT)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Common base: 8 core agents + 5 core instructions shared across all extensions",
                 font_size=13, bold=True, color=ACCENT_TEAL)

    set_speaker_notes(slide, (
        "8 VS Code extensions organized by role. All share a common base of 8 core agents "
        "and 5 core instructions. Pick the extension that matches your role."
    ))


def slide_16_learning_curve(prs):
    """Slide 16: Honest framing + payoff narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Learning Curve & Compounding Value", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_ORANGE)

    add_shape_with_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Honest Truth", ACCENT_ORANGE, font_size=14)
    add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("Your first RPI workflow will feel slower.", 18, True, ACCENT_ORANGE),
        ("", 8),
        ("You'll wonder why you can't just ask AI to code it.", 14, False, SECONDARY_TEXT),
        ("You'll feel like the constraints are slowing you down.", 14, False, SECONDARY_TEXT),
        ("You'll want to skip the research phase.", 14, False, SECONDARY_TEXT),
        ("", 8),
        ("That's normal. Stay with it.", 14, True, TERTIARY_TEXT),
    ])

    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
                        "The Payoff", ACCENT_GREEN, font_size=14)
    add_multi_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.5), [
        ("By your third feature, the workflow", 18, True, ACCENT_GREEN),
        ("feels natural.", 18, True, ACCENT_GREEN),
        ("", 8),
        ("Research artifacts compound across features.", 14, False, SECONDARY_TEXT),
        ("Plans reference previous research findings.", 14, False, SECONDARY_TEXT),
        ("Rework drops below 10%.", 14, False, SECONDARY_TEXT),
        ("", 8),
        ("The framework pays for itself.", 14, True, TERTIARY_TEXT),
    ])

    add_multi_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1.0), [
        ("Research documents accumulate into institutional memory.", 15, True, DARK_TEXT),
        ("New team members can read how past decisions were made. You're not just solving ",
         12, False, SECONDARY_TEXT),
        ("today's problem — you're building the knowledge base that accelerates tomorrow's.",
         12, False, SECONDARY_TEXT),
    ])

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
                 "Evolution Timeline", font_size=14, bold=True, color=DARK_TEXT)

    versions = [
        ("v1.1.0", "Jan 19", "Initial release", ACCENT_BLUE),
        ("v2.0.0", "Jan 28", "Breaking: agent restructure", ACCENT_RED),
        ("v2.3.4", "Feb 13", "Stable + extensions", ACCENT_GREEN),
    ]

    for i, (ver, date, desc, color) in enumerate(versions):
        left = Inches(1.0 + i * 3.8)
        add_shape_with_text(slide, left, Inches(6.6), Inches(1.2), Inches(0.4),
                            ver, color, font_size=11)
        add_text_box(slide, left + Inches(1.4), Inches(6.6), Inches(2.2), Inches(0.4),
                     f"{date} — {desc}", font_size=10, color=TERTIARY_TEXT)
        if i < len(versions) - 1:
            add_connector_line(slide, left + Inches(3.5), Inches(6.8),
                               left + Inches(3.7), Inches(6.8),
                               color=BORDER_GRAY, width=Pt(2))

    set_speaker_notes(slide, (
        "Be honest with your team: the first RPI workflow feels slower. By the third "
        "feature, the workflow feels natural. Research documents accumulate into "
        "institutional memory."
    ))


def slide_17_learning_resources(prs):
    """Slide 17: Three installation paths + learning resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Three Paths to Start + Learn More", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    install_methods = [
        ("VS Code Extension ⭐", "10 seconds", "Individual users, TPMs,\nimmediate access", ACCENT_BLUE),
        ("Peer Clone", "2 minutes", "Developers needing\ncustomization", ACCENT_ORANGE),
        ("Codespaces", "1 click", "Contributors wanting\nzero-config", ACCENT_GREEN),
    ]

    for i, (method, time, desc, color) in enumerate(install_methods):
        left = Inches(0.8 + i * 4.0)
        add_shape_with_text(slide, left, Inches(1.5), Inches(3.5), Inches(0.5),
                            f"{method}  •  {time}", color, font_size=13)
        add_multi_text(slide, left + Inches(0.2), Inches(2.2), Inches(3.2), Inches(1.0), [
            (desc, 12, False, SECONDARY_TEXT),
        ])

    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(11), Inches(0.4),
                 "Learning Resources", font_size=18, bold=True, color=DARK_TEXT)

    resources_left = [
        ("HVE-Learning Repository", "Self-paced modules on prompt engineering, RPI, backlog management", ACCENT_BLUE),
        ("Customer Zero Katas", "Hands-on practice exercises for real-world scenarios", ACCENT_BLUE),
    ]
    for i, (name, desc, color) in enumerate(resources_left):
        row_top = Inches(3.9 + i * 0.6)
        add_shape_with_text(slide, Inches(0.8), row_top, Inches(0.3), Inches(0.3),
                            "•", color, font_size=14, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(1.3), row_top, Inches(5.0), Inches(0.5), [
            (name, 13, True, DARK_TEXT),
            (desc, 11, False, TERTIARY_TEXT),
        ])

    resources_right = [
        ("First Workflow Tutorial", "15-minute guided RPI exercise", ACCENT_GREEN),
        ("Documentation", "aka.ms/hve-core — install, usage, FAQs", ACCENT_GREEN),
    ]
    for i, (name, desc, color) in enumerate(resources_right):
        row_top = Inches(3.9 + i * 0.6)
        add_shape_with_text(slide, Inches(7.0), row_top, Inches(0.3), Inches(0.3),
                            "•", color, font_size=14, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(7.5), row_top, Inches(5.0), Inches(0.5), [
            (name, 13, True, DARK_TEXT),
            (desc, 11, False, TERTIARY_TEXT),
        ])

    add_shape_with_text(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5),
                        "microsoft/hve-learning  •  aka.ms/cz-repo-katas  •  aka.ms/hve-core",
                        ACCENT_BLUE, font_size=13)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Start with one RPI workflow. The research artifacts compound from there.",
                 font_size=13, color=TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "Three paths to get started. The VS Code Extension is the fastest — 10 seconds. "
        "The HVE-Learning repository has self-paced modules. Customer Zero Katas provide "
        "hands-on practice exercises."
    ))


def slide_18_community(prs):
    """Slide 18: Community engagement and contribution pathways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Join the Community", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_PURPLE)

    cards = [
        ("Bi-Weekly Syncs", "HVE Community Syncs — live demos,\nQ&A, roadmap updates,\nand best practice sharing",
         ACCENT_BLUE, Inches(0.8), Inches(1.5)),
        ("Teams Channels", "#Hypervelocity and #SeasonOfHVE\n— real-time discussions,\ntips, troubleshooting",
         ACCENT_TEAL, Inches(7.0), Inches(1.5)),
        ("Good First Issues", 'Look for "good first issue" tags\non HVE-Core and accelerator\nrepos (Edge AI, Robotics)',
         ACCENT_GREEN, Inches(0.8), Inches(3.3)),
        ("Contribution Types", "Code • Prompts • Instructions\n• Skills • Docs • Bug reports\n• Feature requests",
         ACCENT_ORANGE, Inches(7.0), Inches(3.3)),
    ]

    for title, desc, color, left, top in cards:
        add_shape_with_text(slide, left, top, Inches(5.5), Inches(0.5),
                            title, color, font_size=14)
        add_multi_text(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.2), Inches(1.0), [
            (desc, 12, False, SECONDARY_TEXT),
        ])

    flow_items = [
        ("Find Issue", ACCENT_BLUE),
        ("Use RPI to Research", ACCENT_GREEN),
        ("Submit PR", ACCENT_ORANGE),
        ("Get Review", ACCENT_PURPLE),
    ]
    for i, (label, color) in enumerate(flow_items):
        left = Inches(1.0 + i * 2.8)
        add_shape_with_text(slide, left, Inches(5.3), Inches(2.2), Inches(0.5),
                            label, color, font_size=11)
        if i < len(flow_items) - 1:
            add_connector_line(slide, left + Inches(2.3), Inches(5.55),
                               left + Inches(2.7), Inches(5.55),
                               color=BORDER_GRAY, width=Pt(2))

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                 "Use RPI to learn the repo and make your first PR — the agents keep you on track",
                 font_size=13, color=TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide, (
        "Bi-weekly community syncs, Teams channels, and curated good-first-issue tags. "
        "Contributions aren't limited to code — prompts, instructions, documentation, "
        "and bug reports are all valuable."
    ))


def slide_19_getting_started(prs):
    """Slide 19: Streamlined getting started with GitHub URL."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Getting Started", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_GREEN)

    steps = [
        ("1", "Install HVE-Core extension from VS Code Marketplace (10 seconds)", ACCENT_BLUE),
        ("2", "Open Copilot Chat and select a chat mode (e.g., Task Researcher)", ACCENT_GREEN),
        ("3", "Run /task-research on your next feature or bug", ACCENT_ORANGE),
        ("4", "Experience the difference — then try the full RPI cycle", ACCENT_PURPLE),
    ]

    for i, (num, desc, color) in enumerate(steps):
        top = Inches(1.6 + i * 0.8)
        add_shape_with_text(slide, Inches(1.0), top, Inches(0.6), Inches(0.6),
                            num, color, font_size=20, shape=MSO_SHAPE.OVAL)
        add_text_box(slide, Inches(2.0), top + Inches(0.1), Inches(10), Inches(0.5),
                     desc, font_size=16, color=SECONDARY_TEXT)

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
                 "github.com/microsoft/hve-core", font_size=22, bold=True,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    add_shape_with_text(slide, Inches(3.0), Inches(5.8), Inches(7), Inches(0.6),
                        "Try /task-research on your next feature", ACCENT_GREEN, font_size=16)

    set_speaker_notes(slide, (
        "Four steps to get started. Install the extension — 10 seconds. Open Copilot Chat. "
        "Select Task Researcher. Ask it about your next feature."
    ))


def slide_20_key_takeaways(prs):
    """Slide 20: Four key takeaways, closing quote, CTA, and Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Key Takeaways", font_size=32, bold=True)
    add_accent_bar(slide, top=Inches(1.1), color=ACCENT_BLUE)

    takeaways = [
        ("1", "Accelerate delivery and quality", "HVE-Core combines AI tools with structured "
         "practices — tasks up to 88% faster without sacrificing quality or security", ACCENT_BLUE),
        ("2", "RPI is the game-changer", "Research → Plan → Implement → Review turns Copilot "
         "from a nifty helper into a reliable partner. Structured phases kill the AI rework loop", ACCENT_GREEN),
        ("3", "Empower every role", "Not just for coders — PMs, TPMs, security engineers, "
         "data scientists, and OSS contributors each have dedicated workflows", ACCENT_ORANGE),
        ("4", "Start your hypervelocity journey", "It only takes 10 seconds to install. "
         "Pick a small task, run it through RPI, and see the difference", ACCENT_PURPLE),
    ]

    for i, (num, title, desc, color) in enumerate(takeaways):
        top = Inches(1.5 + i * 0.8)
        add_shape_with_text(slide, Inches(0.8), top, Inches(0.5), Inches(0.5),
                            num, color, font_size=16, shape=MSO_SHAPE.OVAL)
        add_multi_text(slide, Inches(1.5), top, Inches(10.5), Inches(0.7), [
            (title, 15, True, DARK_TEXT),
            (desc, 12, False, SECONDARY_TEXT),
        ])

    quote_shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.8), Pt(4), Inches(0.7))
    quote_shp.fill.solid()
    quote_shp.fill.fore_color.rgb = ACCENT_BLUE
    quote_shp.line.fill.background()

    add_text_box(slide, Inches(1.8), Inches(4.8), Inches(10), Inches(0.7),
                 '"The code comes last, after the hard work of understanding is complete."',
                 font_size=18, bold=True, color=DARK_TEXT)

    cta_items = [
        ("Install", ACCENT_GREEN),
        ("Try RPI", ACCENT_BLUE),
        ("Join Community", ACCENT_PURPLE),
    ]
    for i, (label, color) in enumerate(cta_items):
        left = Inches(2.5 + i * 3.0)
        add_shape_with_text(slide, left, Inches(5.8), Inches(2.5), Inches(0.5),
                            label, color, font_size=14)
        if i < len(cta_items) - 1:
            add_connector_line(slide, left + Inches(2.6), Inches(6.05),
                               left + Inches(2.9), Inches(6.05),
                               color=BORDER_GRAY, width=Pt(2))

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.8),
                 "Questions?", font_size=36, bold=True,
                 alignment=PP_ALIGN.CENTER, color=DARK_TEXT)

    set_speaker_notes(slide, (
        "Four takeaways: accelerate delivery and quality, RPI is the game-changer, "
        "empower every role, and start your hypervelocity journey. "
        "The code comes last, after the hard work of understanding is complete. Questions?"
    ))


# ── Main Assembly ──────────────────────────────────────────

def main() -> int:
    """Generate the enhanced HVE-Core and RPI presentation deck (light theme, 20 slides)."""
    print("Generating enhanced HVE-Core and RPI presentation (light theme)...")

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(IMG_DIR, exist_ok=True)

    print("  Creating diagrams (light theme)...")
    rpi_img = create_rpi_pipeline_diagram()
    drpi_img = create_drpi_pipeline_diagram()
    quality_img = create_quality_comparison_diagram()
    role_img = create_role_mapping_diagram()

    print("  Building slides (light theme)...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Part 1: Why RPI? (Slides 1-10)
    slide_01_title(prs)
    slide_02_the_problem(prs)
    slide_03_why_it_happens(prs)
    slide_04_counterintuitive_insight(prs)
    slide_05_what_is_hve_core(prs)
    slide_06_rpi_pipeline(prs, rpi_img)
    slide_07_phase_deep_dive(prs)
    slide_08_quality_and_clear(prs, quality_img)
    slide_09_chat_modes(prs)
    slide_10_demo_intro(prs)

    # Part 2: HVE-Core in Practice (Slides 11-20)
    slide_11_who_uses_hve_core(prs, role_img)
    slide_12_discovery_design_thinking(prs, drpi_img)
    slide_13_real_results(prs)
    slide_14_dogfooding_and_validation(prs)
    slide_15_extension_ecosystem(prs)
    slide_16_learning_curve(prs)
    slide_17_learning_resources(prs)
    slide_18_community(prs)
    slide_19_getting_started(prs)
    slide_20_key_takeaways(prs)

    output_path = os.path.join(OUTPUT_DIR, "hve-core-rpi-presentation-enhanced-light.pptx")
    prs.save(output_path)
    print(f"\n✅ Saved light theme presentation: {output_path}")
    print(f"   {len(prs.slides)} slides generated")
    return EXIT_SUCCESS


if __name__ == "__main__":
    sys.exit(main())
